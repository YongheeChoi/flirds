#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Flirds 세미나 입문 데크 빌더 (재현용).

목적: 연구를 처음 듣는 청중을 위한 "입문 → 연구소개" 흐름.
  FL 동기 → Data Shapley(기초) → In-Run Data Shapley → 우리의 FL 접목(+오차분석)
  → 대표 실험 결과 → 정직한 한계.

결과 데크(../build_deck.py, flirds-results-2026-06-26.pptx)와는 별개 파일.
팔레트·헤더/푸터 스타일은 결과 데크를 그대로 재사용한다(담백한 학술 핸드아웃).

개념·논문 설명 출처(사실 확인):
  research-wiki/wiki/flirds.md
  research-wiki/wiki/concepts/{data-shapley,in-run-data-shapley,shapley-value}.md
  research-wiki/wiki/sources/{ghorbani-zou-data-shapley,in-run-data-shapley}.md
  CLAUDE.md(루트) — 핵심 질문 위계 / 설계 lock
대표 수치 출처:
  ../build_deck.py + flirds-experiment-results-overview-2026-06-25 (개요 문서)

실행: /home/korea_bupj/miniconda3/envs/flirds/bin/python build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (결과 데크와 동일) ----
C_TITLE   = RGBColor(0x1F, 0x2A, 0x37)   # dark slate
C_HEADER  = RGBColor(0x33, 0x47, 0x5B)   # table header fill
C_HEADTXT = RGBColor(0xFF, 0xFF, 0xFF)
C_ROWALT  = RGBColor(0xF2, 0xF4, 0xF7)
C_BODY    = RGBColor(0x22, 0x22, 0x22)
C_MUTED   = RGBColor(0x7A, 0x7A, 0x7A)
C_SETTING = RGBColor(0x3C, 0x46, 0x54)
C_RULE    = RGBColor(0xC9, 0xD2, 0xDC)
C_BG_DARK = RGBColor(0x1F, 0x2A, 0x37)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_ICE     = RGBColor(0xCA, 0xDC, 0xFC)
C_CARD    = RGBColor(0xED, 0xF1, 0xF6)   # 카드 배경 (옅은 청회색)
C_CARDBD  = RGBColor(0xD3, 0xDD, 0xE8)   # 카드 테두리
C_ACCENT  = RGBColor(0x33, 0x47, 0x5B)   # 강조(헤더와 동일 톤)

FONT = "NanumGothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
_pageno = [0]


def _set_run(r, text, size, bold=False, color=C_BODY, italic=False):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_textbox(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, line_space=1.0):
    """lines = list of (text, size, bold, color[, italic])"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_space
        txt, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        ital = ln[4] if len(ln) > 4 else False
        _set_run(p.add_run(), txt, size, bold, color, ital)
    return tb


def _page_number(slide):
    _pageno[0] += 1
    if _pageno[0] == 1:
        return
    tb = slide.shapes.add_textbox(SW - Inches(0.7), SH - Inches(0.38),
                                  Inches(0.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _set_run(p.add_run(), str(_pageno[0]), 9, False, C_MUTED)


def slide_header(slide, title, setting_lines):
    add_textbox(slide, Inches(0.45), Inches(0.30), Inches(12.4), Inches(0.8),
                [(title, 24, True, C_TITLE)])
    y = 1.12
    if setting_lines:
        add_textbox(slide, Inches(0.45), Inches(y), Inches(12.4), Inches(0.7),
                    [(s, 11, False, C_SETTING) for s in setting_lines],
                    line_space=1.05)
        y = 1.12 + 0.22 * len(setting_lines) + 0.06
    return Inches(y)


def slide_footer(slide, source, caption=None):
    add_textbox(slide, Inches(0.45), SH - Inches(0.78), Inches(12.4), Inches(0.7),
                ([(caption, 10, False, C_MUTED, True)] if caption else []) +
                [("출처: " + source, 9, False, C_MUTED)],
                line_space=1.0)


def dark_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG_DARK; bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element); slide.shapes._spTree.insert(2, bg._element)
    return bg


def card(slide, left, top, width, height, fill=C_CARD, line=C_CARDBD):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    sp.adjustments[0] = 0.06
    return sp


def card_text(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP,
              line_space=1.12, pad=0.18):
    return add_textbox(slide, left + Inches(pad), top + Inches(0.12),
                       width - Inches(pad * 2), height - Inches(0.24),
                       lines, anchor=anchor, line_space=line_space)


def chip(slide, left, top, width, height, text, fill=C_ACCENT, txt=C_WHITE, size=12):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    sp.adjustments[0] = 0.5
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), text, size, True, txt)
    return sp


def arrow(slide, left, top, width, height, fill=C_ACCENT):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


# ---- table helpers (결과 데크와 동일) ----
def add_table(slide, top, columns, rows, col_widths_, font_size,
              manual_bold=None, left=Inches(0.45), row_h=Inches(0.3)):
    ncol = len(columns)
    nrow = len(rows) + 1
    total_w = sum(col_widths_)
    height = row_h * nrow
    gfx = slide.shapes.add_table(nrow, ncol, left, top, Emu(int(total_w)), height)
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False
    for c, w in enumerate(col_widths_):
        tbl.columns[c].width = Emu(int(w))
    for c, hdr in enumerate(columns):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = C_HEADER
        cell.margin_left = Pt(3); cell.margin_right = Pt(3)
        cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        _set_run(p.add_run(), hdr, font_size, True, C_HEADTXT)
    bold = manual_bold or set()
    for r, cells in enumerate(rows):
        for c in range(ncol):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_ROWALT if (r % 2 == 1) else C_WHITE
            cell.margin_left = Pt(3); cell.margin_right = Pt(3)
            cell.margin_top = Pt(0); cell.margin_bottom = Pt(0)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            _set_run(p.add_run(), str(cells[c]), font_size, (r, c) in bold, C_BODY)
    for r in range(nrow):
        tbl.rows[r].height = Emu(int(row_h))
    return gfx


def col_widths(weights, total=Inches(12.45)):
    s = sum(weights)
    return [int(total * w / s) for w in weights]


def newslide():
    return prs.slides.add_slide(BLANK)


# ================================================================ Slide 1: 표지
s = newslide()
dark_bg(s)
add_textbox(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(1.3),
            [("Flirds", 46, True, C_WHITE)])
add_textbox(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.0),
            [("연합학습에서 클라이언트별 데이터 기여도 측정", 26, True, C_ICE)])
add_textbox(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(1.0),
            [("In-Run Data Shapley(한 학습 궤적, 1차+2차 Taylor, true Hessian)를 "
              "연합학습(FL)에 접목한 client-level 데이터 가치평가 방법.", 15, False, C_WHITE)],
            line_space=1.2)
add_textbox(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.9),
            [("발표 흐름:  FL 동기  →  Data Shapley(기초)  →  In-Run Data Shapley  "
              "→  FL 접목 + 오차 분석  →  대표 실험 결과  →  novelty  →  정직한 한계", 13, False, C_ICE)],
            line_space=1.25)
add_textbox(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.8),
            [("세미나 입문 발표 · 2026-06-30", 13, True, C_ICE),
             ("연구 처음 듣는 청중 대상. 개념 설명 출처 = research-wiki (flirds.md / concepts / sources); "
              "대표 수치 출처 = 결과 데크 flirds-results-2026-06-26.", 10, False, C_ICE)],
            line_space=1.2)
_page_number(s)

# ================================================================ Slide 2: 목차
s = newslide()
top = slide_header(s, "목차", None)
items = [
    ("1.", "왜 연합학습에서 '기여도 측정'인가", "보상 분배 · free-rider 탐지 · 불량 데이터 식별"),
    ("2.", "Data Shapley — 데이터 가치의 공정한 정의", "협력게임 Shapley value, 원조(Ghorbani & Zou 2019)와 그 한계"),
    ("3.", "In-Run Data Shapley", "재학습 없이 한 학습 궤적에서 기여도를 추적 (Taylor 분해)"),
    ("4.", "우리 연구 — In-Run Data Shapley의 FL 접목 (Flirds)", "client·round 수준 설계 + 근사 오차 분석(왜 2차 항이 필요한가)"),
    ("5.", "대표 실험 결과", "① Fidelity(오라클 대비 정확도; 2차 항 효과 포함) → ② 비용·확장 → ③ 강건성/탐지"),
    ("6.", "우리 연구의 novelty", "기존에 비어 있던 교집합 — 가장 가까운 경쟁자 대비 차별점"),
    ("7.", "정직한 한계", "valuation의 경계 — clean-preserving backdoor"),
]
y = 1.42
for num, title, sub in items:
    card(s, Inches(0.55), Inches(y), Inches(0.68), Inches(0.72), fill=C_ACCENT, line=C_ACCENT)
    add_textbox(s, Inches(0.55), Inches(y), Inches(0.68), Inches(0.72),
                [(num, 19, True, C_WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(1.43), Inches(y + 0.03), Inches(11.4), Inches(0.4),
                [(title, 15.5, True, C_TITLE)])
    add_textbox(s, Inches(1.43), Inches(y + 0.4), Inches(11.4), Inches(0.32),
                [(sub, 11.5, False, C_MUTED)])
    y += 0.84
_page_number(s)

# ================================================================ Slide 3: FL 동기
s = newslide()
top = slide_header(s, "왜 연합학습(FL)에서 '클라이언트별 기여도'를 재는가",
    ["연합학습 = 데이터를 한곳에 모으지 않고, 각 참여자(클라이언트)가 로컬에서 학습한 업데이트만 서버가 모아 모델을 만드는 방식.",
     "이때 자연스러운 질문: \"이 공동 모델에 누가 얼마나 기여했나?\" — 세 가지 실제 필요로 이어진다."])
cards = [
    ("보상 · 인센티브 분배", "참여자에게 정당한 몫을 나누려면 기여도를 수치로 매겨야 한다. "
     "데이터 시장·연합 컨소시엄의 핵심 동기."),
    ("free-rider 탐지", "학습에 실질 기여 없이 모델만 받아가는 무임승차자를 식별. "
     "기여도가 0에 가까운 클라이언트를 가려낸다."),
    ("불량 데이터 클라 식별", "라벨 오류·노이즈·오염 데이터를 가진 클라이언트는 "
     "기여도가 낮거나 음수 → 가중치 하향·배제의 근거."),
]
cw = Inches(3.95); gap = Inches(0.28); x0 = Inches(0.55); cy = Inches(2.35); ch = Inches(2.7)
for i, (h, body) in enumerate(cards):
    x = x0 + i * (cw + gap)
    card(s, x, cy, cw, ch)
    chip(s, x + Inches(0.25), cy + Inches(0.25), Inches(0.5), Inches(0.5), str(i + 1))
    add_textbox(s, x + Inches(0.9), cy + Inches(0.28), cw - Inches(1.1), Inches(0.5),
                [(h, 15, True, C_TITLE)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, x + Inches(0.25), cy + Inches(1.05), cw - Inches(0.5), Inches(1.5),
                [(body, 12.5, False, C_BODY)], line_space=1.22)
add_textbox(s, Inches(0.55), cy + ch + Inches(0.22), Inches(12.2), Inches(0.6),
    [("Flirds의 주 용도는 연구-측 in-run 기여도 측정(valuation)이며, 위 셋(인센티브·탐지)은 그 위에 얹히는 응용이다.",
      11.5, False, C_MUTED, True)])
slide_footer(s, "research-wiki/wiki/flirds.md (Q1 use-case lock) · CLAUDE.md 핵심 질문 위계")
_page_number(s)

# ================================================================ Slide 4: Data Shapley (1) 동기
s = newslide()
top = slide_header(s, "Data Shapley (1) — \"누구의 데이터가 모델에 얼마나 기여했나?\"",
    ["데이터 가치평가(data valuation)의 근본 질문. 2019년 이전엔 leave-one-out(LOO)·영향력 휴리스틱이 주로 쓰였으나 공정성 근거가 약했다."])
# 좌: 질문 / 우: LOO의 한계
lx = Inches(0.55); lw = Inches(6.0); ly = Inches(2.0); lh = Inches(3.7)
card(s, lx, ly, lw, lh)
add_textbox(s, lx + Inches(0.3), ly + Inches(0.25), lw - Inches(0.6), Inches(3.2),
    [("핵심 질문", 15, True, C_ACCENT),
     ("• 한 모델을 여러 데이터(또는 여러 사람의 데이터)로 학습했을 때,", 13, False, C_BODY),
     ("   각 데이터가 최종 성능에 준 '몫'을 어떻게 공정하게 나눌까?", 13, False, C_BODY),
     ("", 6, False, C_BODY),
     ("• 좋은 가치 척도는 이상치·라벨 오류를 드러내고,", 13, False, C_BODY),
     ("   어떤 데이터를 더 사들일지(data acquisition) 결정에 쓰인다.", 13, False, C_BODY),
     ("", 6, False, C_BODY),
     ("• 게임이론의 Shapley value가 \"공정 분배\"의 표준 답을 준다.", 13, True, C_TITLE)],
    line_space=1.18)
rx = Inches(6.85); rw = Inches(5.9)
card(s, rx, ly, rw, lh)
add_textbox(s, rx + Inches(0.3), ly + Inches(0.25), rw - Inches(0.6), Inches(3.2),
    [("LOO(하나 빼보기)로는 왜 부족한가", 15, True, C_ACCENT),
     ("• LOO = '이 데이터 하나를 빼면 성능이 얼마나 떨어지나'.", 13, False, C_BODY),
     ("", 5, False, C_BODY),
     ("• 비슷한 데이터가 둘 있으면, 하나를 빼도 다른 하나가 메워", 13, False, C_BODY),
     ("   둘 다 가치 0으로 평가 — 중복에 취약.", 13, False, C_BODY),
     ("", 5, False, C_BODY),
     ("• 데이터 간 상호작용(누구와 함께 있을 때의 기여)을 무시.", 13, False, C_BODY),
     ("", 5, False, C_BODY),
     ("• Shapley는 '모든 조합에서의 한계기여 평균'으로 이를 해결.", 13, True, C_TITLE)],
    line_space=1.18)
slide_footer(s, "research-wiki/wiki/sources/ghorbani-zou-data-shapley.md · concepts/data-shapley.md")
_page_number(s)

# ================================================================ Slide 5: Data Shapley (2) 정의
s = newslide()
top = slide_header(s, "Data Shapley (2) — 협력게임 Shapley value",
    ["직관: 각 데이터를 '협력게임의 플레이어'로 보고, 가능한 모든 참여 순서(조합)에서 그 데이터를 추가했을 때 늘어난 효용(=한계기여)을 평균낸다."])
# 정의식 카드
fx = Inches(0.55); fw = Inches(12.2); fy = Inches(2.0)
card(s, fx, fy, fw, Inches(1.35), fill=C_BG_DARK, line=C_BG_DARK)
add_textbox(s, fx + Inches(0.4), fy + Inches(0.18), fw - Inches(0.8), Inches(1.0),
    [("φ_i  =  C · Σ_{S ⊆ D∖{i}}  C(n−1, |S|)⁻¹ · ( U(S ∪ {i}) − U(S) )", 20, True, C_WHITE),
     ("데이터 i의 가치 = 모든 부분집합 S에 i를 추가했을 때의 효용 증가분 U(S∪{i})−U(S) 의 가중 평균. "
      "U(S) = S로 학습한 모델의 검증 성능.", 11.5, False, C_ICE)],
    line_space=1.25)
# 4 공리
ay = fy + Inches(1.65)
add_textbox(s, fx, ay, fw, Inches(0.4), [("이 식은 다음 네 공리를 모두 만족하는 유일한 분배다 (Shapley 1953; 데이터로 가져온 게 Ghorbani & Zou 2019):", 12.5, True, C_TITLE)])
axioms = [
    ("Null(영 기여)", "어디에 넣어도 한계기여가 0인 데이터는 가치 0."),
    ("Symmetry(대칭)", "교환 가능한(동일 역할) 데이터는 같은 가치."),
    ("Additivity(가법성)", "효용에 대해 선형 — 두 과제 가치의 합."),
    ("Efficiency(효율)", "모든 가치의 합 = 전체 효용 (남김없이 분배)."),
]
cw = Inches(2.95); gap = Inches(0.13); cy = ay + Inches(0.5); ch = Inches(1.85)
for i, (h, body) in enumerate(axioms):
    x = fx + i * (cw + gap)
    card(s, x, cy, cw, ch)
    add_textbox(s, x + Inches(0.2), cy + Inches(0.18), cw - Inches(0.4), Inches(0.5),
                [(h, 13, True, C_ACCENT)])
    add_textbox(s, x + Inches(0.2), cy + Inches(0.72), cw - Inches(0.4), Inches(1.0),
                [(body, 12, False, C_BODY)], line_space=1.18)
slide_footer(s, "research-wiki/wiki/concepts/shapley-value.md · sources/ghorbani-zou-data-shapley.md (ICML 2019, arXiv:1904.02868)")
_page_number(s)

# ================================================================ Slide 6: Data Shapley (3) 한계
s = newslide()
top = slide_header(s, "Data Shapley (3) — 원조의 한계: 2ᴺ 재학습",
    ["정의대로 계산하려면 데이터의 모든 부분집합 S마다 모델을 다시 학습해 U(S)를 재야 한다 → 데이터 n개면 2ⁿ 번 재학습.",
     "근사 추정기(TMC-Shapley 등)도 결국 '많은 재학습'을 요구 → foundation model 규모에선 비현실적. 이것이 다음 절(In-Run)의 다리."])
cols = ["접근", "비용", "적용 범위"]
rows = [
    ["Exact Shapley", "O(2ⁿ) 재학습", "장난감 규모만"],
    ["Monte-Carlo 순열 샘플링", "다수 재학습", "소규모"],
    ["Truncated MC (Ghorbani & Zou)", "재학습 수 절감", "소~중규모"],
    ["KNN-Shapley 닫힌형", "O(n log n), KNN 한정", "제한적"],
    ["In-Run Data Shapley", "≈ 학습 1회", "foundation model ← 다음 절"],
]
mb = {(4, 0), (4, 1), (4, 2)}
add_table(s, Inches(2.2), cols, rows, col_widths([3.2, 2.4, 3.0]), 13, manual_bold=mb, row_h=Inches(0.5))
add_textbox(s, Inches(0.45), Inches(5.3), Inches(12.4), Inches(1.1),
    [("두 가지 문제 (In-Run 논문이 지적):", 13, True, C_TITLE),
     ("① 계산 — 수많은 부분집합 재학습은 대형 모델에서 불가능.", 12.5, False, C_BODY),
     ("② 개념 — 원조는 '학습 알고리즘'에 대한 Shapley라 초기화·배치 순서의 무작위성을 평균낸다. "
      "그러나 실무자는 '내가 실제로 학습한 그 모델'에 대한 기여를 알고 싶다 → model-specific 정의가 필요.",
      12.5, False, C_BODY)],
    line_space=1.2)
slide_footer(s, "research-wiki/wiki/concepts/data-shapley.md (computational ladder) · sources/in-run-data-shapley.md")
_page_number(s)


# ================================================================ Slide 7: In-Run (1) 아이디어
s = newslide()
top = slide_header(s, "In-Run Data Shapley (1) — 한 번의 학습 안에서 추적",
    ["핵심 전환: 부분집합마다 재학습하는 대신, '실제로 진행된 한 학습 과정' 안에서 각 데이터가 검증 손실 변화에 준 기여를 스텝마다 누적한다.",
     "Wang·Mittal·Song·Jia, \"Data Shapley in One Training Run\", ICML 2024 (arXiv:2406.11011). 최적화 구현은 표준 학습과 거의 같은 비용."])
# 좌: 원조 vs 우: in-run 대비
lx = Inches(0.55); lw = Inches(5.95); ly = Inches(2.05); lh = Inches(2.7)
card(s, lx, ly, lw, lh)
add_textbox(s, lx + Inches(0.28), ly + Inches(0.2), lw - Inches(0.56), Inches(2.3),
    [("원조 (retraining-based)", 14, True, C_ACCENT),
     ("• 학습 '알고리즘'에 대한 Shapley — 무작위성 평균.", 12.5, False, C_BODY),
     ("• 부분집합마다 모델을 처음부터 재학습.", 12.5, False, C_BODY),
     ("• 비용 O(2ⁿ) 재학습 → 대형 모델 불가.", 12.5, False, C_BODY)],
    line_space=1.25)
rx = Inches(6.85); rw = Inches(5.9)
card(s, rx, ly, rw, lh, fill=RGBColor(0xE6,0xEF,0xE6), line=RGBColor(0xC4,0xD8,0xC4))
add_textbox(s, rx + Inches(0.28), ly + Inches(0.2), rw - Inches(0.56), Inches(2.3),
    [("In-Run (model-specific)", 14, True, RGBColor(0x2C,0x5F,0x2D)),
     ("• '내가 실제로 학습한 그 모델'에 대한 기여.", 12.5, False, C_BODY),
     ("• 재학습 없음 — 학습하며 스텝별 기여를 누적.", 12.5, False, C_BODY),
     ("• 비용 ≈ 학습 1회 → GPT-2 / Pythia-410M 규모 입증.", 12.5, False, C_BODY)],
    line_space=1.25)
# 누적 식
fy = ly + lh + Inches(0.25)
card(s, Inches(0.55), fy, Inches(12.2), Inches(0.95), fill=C_BG_DARK, line=C_BG_DARK)
add_textbox(s, Inches(0.95), fy + Inches(0.13), Inches(11.4), Inches(0.7),
    [("φᵢ(in-run)  =  Σₜ φᵢ⁽ᵗ⁾        (스텝 t의 기여 φᵢ⁽ᵗ⁾ 를 학습 전 구간에 걸쳐 합산)", 18, True, C_WHITE)],
    anchor=MSO_ANCHOR.MIDDLE)
slide_footer(s, "research-wiki/wiki/sources/in-run-data-shapley.md · concepts/in-run-data-shapley.md")
_page_number(s)

# ================================================================ Slide 8: In-Run (2) Taylor 분해
s = newslide()
top = slide_header(s, "In-Run Data Shapley (2) — 스텝 효용 변화의 Taylor 분해",
    ["한 SGD 스텝이 검증 효용을 U(θₜ₊₁)−U(θₜ) 만큼 바꾼다. 이 변화를 배치 안 데이터에 대해 Taylor 전개하면 스텝별 Shapley가 닫힌형으로 나온다."])
# 1차 / 2차 항 카드
lx = Inches(0.55); lw = Inches(6.0); ly = Inches(2.05); lh = Inches(2.55)
card(s, lx, ly, lw, lh)
add_textbox(s, lx + Inches(0.3), ly + Inches(0.2), lw - Inches(0.6), Inches(2.2),
    [("1차 항 (gradient · gradient)", 15, True, C_ACCENT),
     ("φᵢ⁽ᵗ⁾  ∝  ⟨ g_val , gᵢ ⟩", 17, True, C_TITLE),
     ("검증 gradient와 데이터 i의 gradient의 내적.", 12.5, False, C_BODY),
     ("두 방향이 일치할수록(검증 손실을 낮추는 쪽) 기여가 크다.", 12.5, False, C_BODY)],
    line_space=1.3)
rx = Inches(6.85); rw = Inches(5.9)
card(s, rx, ly, rw, lh)
add_textbox(s, rx + Inches(0.3), ly + Inches(0.2), rw - Inches(0.6), Inches(2.2),
    [("2차 항 (gradient · Hessian · gradient)", 15, True, C_ACCENT),
     ("φᵢ⁽ᵗ⁾  에  gᵀ H g  꼴이 추가", 17, True, C_TITLE),
     ("곡률(Hessian)을 반영 — 데이터 간 상호작용/중복을 포착.", 12.5, False, C_BODY),
     ("변위가 클수록 1차만으론 부정확 → 2차가 보정.", 12.5, False, C_BODY)],
    line_space=1.3)
add_textbox(s, Inches(0.55), ly + lh + Inches(0.25), Inches(12.2), Inches(1.2),
    [("원 논문(IRDS)의 관찰: 중앙집중 학습의 per-step 정밀(작은 η) 영역에선 2차 항이 정확도를 크게 높이지 않는다 (Appx E.2.2).",
      12.5, False, C_BODY),
     ("→ 우리 연구의 출발점: FL은 라운드당 변위가 커서 무대가 다르다. (다음 절의 오차 분석에서 정확히 다룬다.)",
      12.5, True, C_TITLE),
     ("※ 본 발표는 ghost dot-product 등 per-sample-gradient 효율화 트릭은 다루지 않는다(본 연구와 직접 무관).",
      11, False, C_MUTED, True)],
    line_space=1.25)
slide_footer(s, "research-wiki/wiki/sources/in-run-data-shapley.md (Method: per-step Shapley via Taylor expansion)")
_page_number(s)

# ================================================================ Slide 9: FL (1) 세팅
s = newslide()
top = slide_header(s, "우리 연구 (1) — 연합학습(FL) 세팅 한 장",
    ["FedAvg: 매 라운드 r 마다 ① 서버가 현재 모델 wʳ 을 클라이언트에 보냄 → ② 각 클라이언트 k가 자기 데이터로 여러 로컬 스텝 학습해 업데이트 Δwₖ 산출 "
     "→ ③ 서버가 이를 모아(가중 평균) 다음 모델 wʳ⁺¹ 을 만든다."])
# 흐름 다이어그램: 서버 -> 클라이언트들 -> 집계
boxes = [
    (Inches(0.7),  "서버 모델  wʳ", "전체에 배포"),
    (Inches(4.5),  "클라이언트 k", "로컬 E스텝 학습 → Δwₖ"),
    (Inches(8.3),  "서버 집계", "wʳ⁺¹ = wʳ + Σ pₖ·Δwₖ"),
]
by = Inches(2.3); bw = Inches(3.4); bh = Inches(1.5)
for i, (bx, h, sub) in enumerate(boxes):
    card(s, bx, by, bw, bh, fill=C_CARD if i != 1 else RGBColor(0xE6,0xEF,0xE6),
         line=C_CARDBD if i != 1 else RGBColor(0xC4,0xD8,0xC4))
    add_textbox(s, bx + Inches(0.2), by + Inches(0.28), bw - Inches(0.4), Inches(0.5),
                [(h, 16, True, C_TITLE)], align=PP_ALIGN.CENTER)
    add_textbox(s, bx + Inches(0.2), by + Inches(0.85), bw - Inches(0.4), Inches(0.5),
                [(sub, 12.5, False, C_BODY)], align=PP_ALIGN.CENTER)
    if i < 2:
        arrow(s, bx + bw + Inches(0.0), by + Inches(0.55), Inches(0.4), Inches(0.4))
add_textbox(s, Inches(0.55), Inches(4.35), Inches(12.2), Inches(1.7),
    [("서버가 이미 받는 것 = 각 클라이언트의 업데이트 Δwₖ 뿐.", 14, True, C_TITLE),
     ("• Flirds는 바로 이 Δwₖ 만으로 client-level 기여도를 계산한다 → 추가 통신·추가 통계량 0 (vanilla FedAvg 위에 그대로 얹힘).", 12.5, False, C_BODY),
     ("• 단위: 데이터포인트가 아니라 '클라이언트'. 분해: 스텝이 아니라 '라운드'. (다음 장에서 이유 설명)", 12.5, False, C_BODY),
     ("• pₖ = 집계 가중치(보통 데이터 크기 비례). E = 라운드당 로컬 스텝 수.", 11.5, False, C_MUTED)],
    line_space=1.25)
slide_footer(s, "research-wiki/wiki/flirds.md (One-line method statement · Locked design decisions)")
_page_number(s)

# ================================================================ Slide 10: FL (2) 설계 3선택
s = newslide()
top = slide_header(s, "우리 연구 (2) — Flirds의 핵심 설계 선택 3가지",
    ["In-Run Data Shapley를 FL로 가져올 때의 세 가지 결정. 그 결과가 아래 라운드별 닫힌형 공식."])
choices = [
    ("(a) 단위: 데이터포인트 → 클라이언트", "FL에서 가장 자연스러운 귀속 단위. 서버는 클라이언트 업데이트만 보므로 client-level이 직접 대상."),
    ("(b) 분해: 스텝 → 라운드", "라운드마다 각 클라이언트 업데이트가 검증 손실에 준 델타를 1차+2차 Taylor로 추정. 라운드당 HVP(Hessian-vector product) 단 1회."),
    ("(c) 효용=검증 손실, 곡률=true Hessian, SGD momentum=0", "게임의 효용은 검증 손실. 곡률은 진짜 Hessian(GGN/Fisher 아님 — 테스트해보니 더 나빴음). 2차 Taylor가 성립하도록 plain SGD 가정."),
]
cy = Inches(1.9); ch = Inches(1.04); cgap = Inches(0.11)
for i, (h, body) in enumerate(choices):
    y = cy + i * (ch + cgap)
    card(s, Inches(0.55), y, Inches(12.2), ch)
    add_textbox(s, Inches(0.8), y + Inches(0.12), Inches(11.7), Inches(0.42),
                [(h, 14, True, C_ACCENT)])
    add_textbox(s, Inches(0.8), y + Inches(0.55), Inches(11.7), Inches(0.45),
                [(body, 12, False, C_BODY)], line_space=1.12)
# core formula
fy = cy + 3 * (ch + cgap) + Inches(0.06)
card(s, Inches(0.55), fy, Inches(12.2), Inches(0.92), fill=C_BG_DARK, line=C_BG_DARK)
add_textbox(s, Inches(0.9), fy + Inches(0.12), Inches(11.5), Inches(0.7),
    [("φₖ⁽ʳ⁾ ≈ −∇ℓ(wʳ, z_val)·Δwₖ  +  ½ Δwₖᵀ H_val(wʳ)·ΔW⁽ʳ⁾        (라운드 합:  φₖ = Σᵣ φₖ⁽ʳ⁾)", 16, True, C_WHITE),
     ("1차항(검증 gradient·업데이트 내적) + 2차항(클라이언트 상호작용; ΔW⁽ʳ⁾=참여 클라 업데이트 합). 비용 = HVP 1회 + N번 내적.", 10.5, False, C_ICE)],
    anchor=MSO_ANCHOR.MIDDLE, line_space=1.2)
slide_footer(s, "research-wiki/wiki/flirds.md (Core formula · Locked design decisions) · CLAUDE.md §code conventions")
_page_number(s)

# ================================================================ Slide 11: FL (3) 왜 in-run
s = newslide()
top = slide_header(s, "우리 연구 (3) — 왜 in-run이어야 하나 (FL에서의 비용)",
    ["FL에서 '정확한' 클라이언트 Shapley를 직접 재려면, 매 라운드마다 N개 클라이언트의 모든 부분집합(2ᴺ)에 대해 모델을 다시 만들어 검증 손실을 재야 한다."])
lx = Inches(0.55); lw = Inches(6.0); ly = Inches(2.1); lh = Inches(2.5)
card(s, lx, ly, lw, lh)
add_textbox(s, lx + Inches(0.3), ly + Inches(0.2), lw - Inches(0.6), Inches(2.1),
    [("정확 오라클 (왜 비싼가)", 15, True, C_ACCENT),
     ("• 라운드마다 2ᴺ 부분집합 평가 → R 라운드면 2ᴺ·R.", 12.5, False, C_BODY),
     ("• N=5 → 32×, N=10 → 1024× (라운드마다).", 12.5, False, C_BODY),
     ("• 실험 검증용으론 쓰되(아래 (b) 오라클), 실제 운용엔 불가.", 12.5, False, C_BODY)],
    line_space=1.28)
rx = Inches(6.85); rw = Inches(5.9)
card(s, rx, ly, rw, lh, fill=RGBColor(0xE6,0xEF,0xE6), line=RGBColor(0xC4,0xD8,0xC4))
add_textbox(s, rx + Inches(0.3), ly + Inches(0.2), rw - Inches(0.6), Inches(2.1),
    [("Flirds (in-run 추정)", 15, True, RGBColor(0x2C,0x5F,0x2D)),
     ("• 라운드당 HVP 1회 + N번 내적 (닫힌형).", 12.5, False, C_BODY),
     ("• 서버가 이미 받은 Δwₖ 만 사용 → 추가 통신 0.", 12.5, False, C_BODY),
     ("• LoRA로 파라미터 차원이 작아 HVP가 저렴.", 12.5, False, C_BODY)],
    line_space=1.28)
add_textbox(s, Inches(0.55), ly + lh + Inches(0.25), Inches(12.2), Inches(1.4),
    [("실험에서 쓰는 두 가지 '정답(oracle)' — 추정이 맞는지 검증하기 위해서만 (비싸므로 작은 N에서):", 13, True, C_TITLE),
     ("• (b) in-run 정확 오라클: 한 학습 궤적에서 2ᴺ 분해로 정확 Shapley를 계산 — Flirds가 근사하려는 바로 그 대상.", 12.5, False, C_BODY),
     ("• (a) retrain 정확 오라클: 부분집합마다 처음부터 재학습 — 데이터-가치 커뮤니티의 고전적 정답(별개 게임).", 12.5, False, C_BODY)],
    line_space=1.25)
slide_footer(s, "research-wiki/wiki/flirds.md (Resolved Q5: dual oracle) · CLAUDE.md (LLM backend: (b) oracle cost = 2ᴺ·R·val·seq)")
_page_number(s)

# ================================================================ Slide 12: FL (4) 오차 분석 [필수]
s = newslide()
top = slide_header(s, "우리 연구 (4) — 근사 오차 분석: 왜 FL에선 2차 항이 필요한가",
    ["Taylor 전개의 잔차(버리는 고차항)는 '변위 크기'에 좌우된다. 중앙집중 IRDS와 FL은 변위 규모가 달라, 같은 1차 근사라도 오차 규모가 다르다."])
# 두 레짐 비교
lx = Inches(0.55); lw = Inches(6.0); ly = Inches(2.0); lh = Inches(2.95)
card(s, lx, ly, lw, lh)
add_textbox(s, lx + Inches(0.28), ly + Inches(0.18), lw - Inches(0.56), Inches(2.6),
    [("원조 IRDS — 중앙집중 per-step", 15, True, C_ACCENT),
     ("• 검증 사이 변위 = SGD '한 스텝' (작은 η) → ‖Δθ‖ 작음.", 12.5, False, C_BODY),
     ("• 1차 Taylor 잔차 = O(‖Δθ‖²) → 이미 매우 작음.", 12.5, False, C_BODY),
     ("• 따라서 2차 항의 역할이 미미 (논문 Appx E.2.2와 일치).", 12.5, False, C_BODY)],
    line_space=1.3)
rx = Inches(6.85); rw = Inches(5.9)
card(s, rx, ly, rw, lh, fill=RGBColor(0xF7,0xEC,0xE8), line=RGBColor(0xE3,0xC9,0xC0))
add_textbox(s, rx + Inches(0.28), ly + Inches(0.18), rw - Inches(0.56), Inches(2.6),
    [("FL — 라운드당 multi-step", 15, True, RGBColor(0xB8,0x50,0x42)),
     ("• 클라가 여러 로컬 스텝 후 집계 → 라운드 변위 ‖Δw_round‖ ≫ per-step.", 12.5, False, C_BODY),
     ("• 1차 잔차 O(‖Δw_round‖²) 가 커진다 → 1차만으론 부정확.", 12.5, False, C_BODY),
     ("• 2차(Hessian) 항을 넣으면 잔차가 O(‖Δw‖³) 로 낮아져 fidelity 유지.", 12.5, True, C_TITLE)],
    line_space=1.3)
# 하단: 스케일 사다리 + Prop 2
fy = ly + lh + Inches(0.2)
add_textbox(s, Inches(0.55), fy, Inches(12.2), Inches(1.0),
    [("스케일 논증(증명 아님):  변위가 ε배 커지면 1차 잔차는 ε² 배, 2차 잔차는 ε³ 배로 줄어든다. "
      "per-step에선 ε≈작아 둘 다 무시 가능하지만, 라운드 변위에선 ε²가 더 이상 작지 않다 → 2차 항이 비로소 load-bearing.", 12.5, False, C_BODY),
     ("Proposition 2 (informal): 라운드 변위 반경 Rᵣ에 대해 client drift residual은 Rᵣ의 3차(cubic) 함수로 bound된다. → 이것이 우리가 2차를 LLM 스케일에서 검증한 이유.",
      12.5, True, C_TITLE)],
    line_space=1.3)
slide_footer(s, "research-wiki/wiki/flirds.md (FL multi-step deviation: drift residual; Prop 1·2) · sources/in-run-data-shapley.md (Notes 2026-06-03)")
_page_number(s)


# ================================================================ Slide 13: 2차가 1차를 이긴 실측
s = newslide()
top = slide_header(s, "우리 연구 (4-보강) — \"2차 항이 1차를 이긴\" 실측 근거",
    ["오차 분석의 예측(FL에선 2차가 의미 있다)을 실제 데이터가 뒷받침한다. 두 개의 구체적 데이터포인트."])
lx = Inches(0.55); lw = Inches(6.0); ly = Inches(2.05); lh = Inches(3.0)
card(s, lx, ly, lw, lh)
add_textbox(s, lx + Inches(0.3), ly + Inches(0.2), lw - Inches(0.6), Inches(2.6),
    [("CNN, plain SGD — 추정 vs (b)오라클 Spearman", 14, True, C_ACCENT),
     ("• 2차 포함:  0.96", 14, True, C_TITLE),
     ("• 1차만:      0.92", 14, False, C_BODY),
     ("", 5, False, C_BODY),
     ("• momentum=0.9 에선 역전 (2차 0.73 < 1차 0.81):", 12.5, False, C_BODY),
     ("   모멘텀의 속도 꼬리가 'Taylor가 전개하는 변위'와 어긋남 → plain SGD가 convention인 이유.", 12, False, C_MUTED)],
    line_space=1.3)
rx = Inches(6.85); rw = Inches(5.9)
card(s, rx, ly, rw, lh, fill=RGBColor(0xF7,0xEC,0xE8), line=RGBColor(0xE3,0xC9,0xC0))
add_textbox(s, rx + Inches(0.3), ly + Inches(0.2), rw - Inches(0.6), Inches(2.6),
    [("LLM 1B, cross-silo N=5 — poison 탐지 AUROC", 14, True, RGBColor(0xB8,0x50,0x42)),
     ("• Flirds (2차):    0.917", 14, True, C_TITLE),
     ("• Flirds-1st:       0.000  (완전 회피)", 14, False, C_BODY),
     ("", 5, False, C_BODY),
     ("• clean-보존 backdoor 공격자를 1차 항은 전혀 못 잡고,", 12.5, False, C_BODY),
     ("   2차 Hessian 항이 부호를 일부 복원해 잡아낸다 (LLM 스케일 최초 데이터포인트).", 12, False, C_MUTED)],
    line_space=1.3)
add_textbox(s, Inches(0.55), ly + lh + Inches(0.22), Inches(12.2), Inches(0.7),
    [("정리: 2차 항의 이득은 '변위가 큰' FL 무대에서만 드러난다 — IRDS가 보고한 '중앙 per-step에선 이득 미미'와 모순이 아니라 무대가 다른 것.",
      12, False, C_MUTED, True)])
slide_footer(s, "research-wiki/wiki/flirds.md (Phase 0.5 findings) · ../build_deck.py Slide14 (silo5 poison) · 개요 §3.4.1")
_page_number(s)

# ================================================================ Slide 14: 결과 (1) Fidelity
s = newslide()
top = slide_header(s, "대표 결과 (1·핵심) — Fidelity: 오라클 대비 기여도 정확도",
    ["핵심 질문 위계의 1차. LLM 표준 실험(LoRA; OpenFedLLM 표준 레시피) · 정답 = in-run 정확 오라클 (b) · 3-seed mean. 값 클수록 충실(↑), Spearman = 순위 상관.",
     "발췌(전체 표·열은 결과 데크 참조). 표준 스테이지 = 참여자 20·라운드당 2; 정밀 스테이지 = 참여자 5·전원."])
cols = ["method", "1B 표준 Sp ↑", "3B 표준 Sp ↑", "7B 표준 Sp ↑", "1B 정밀 Sp ↑", "7B 정밀 Sp ↑"]
rows = [
    ["Flirds",      "1.000", "1.000", "0.999", "1.000", "1.000"],
    ["Flirds-1st",  "0.999", "1.000", "0.998", "1.000", "1.000"],
    ["loss-heur",   "1.000", "0.999", "0.999", "1.000", "1.000"],
    ["GTG-Shapley", "0.975", "0.988", "0.977", "1.000", "1.000"],
    ["FedSV",       "0.910", "0.952", "0.968", "0.700", "0.933"],
    ["ShapleyFL",   "0.194", "0.227", "0.406", "0.700", "0.833"],
    ["FedIF",       "0.157", "0.211", "0.480", "0.067", "0.200"],
]
mb = {(0,1),(0,2),(0,4),(0,5),(2,1),(2,3)}
add_table(s, Inches(2.55), cols, rows, col_widths([1.8,1.5,1.5,1.5,1.5,1.5]), 11, manual_bold=mb, row_h=Inches(0.34))
add_textbox(s, Inches(0.45), Inches(5.55), Inches(12.4), Inches(1.0),
    [("• Flirds·Flirds-1st·loss-heur가 천장(≈1.000) — exact in-run Shapley의 순위를 그대로 재현. GTG/FedSV가 그 다음.", 12.5, False, C_BODY),
     ("• FedIF/ShapleyFL/ComFedSV는 설계상 낮음(영향도·surrogate·low-rank 가정). N=5(정밀)는 전수 2⁵라 다수 방법이 1.000 동률.", 12.5, False, C_BODY)],
    line_space=1.25)
slide_footer(s, "../build_deck.py Slide 3·4 (개요 §3.1.1, runs/track_d/fidelity.csv)",
             "Sp=Spearman. ±std·Kendall·Pearson·CNN 표는 결과 데크 참조. 정밀 스테이지엔 (a) retrain 오라클 대비 0.933 일치도도 측정됨.")
_page_number(s)

# ================================================================ Slide 14b: Fidelity — 2차 항 효과 (CNN)
s = newslide()
top = slide_header(s, "대표 결과 (1·보강) — Fidelity에서 2차 항의 효과 (CNN 실험)",
    ["LLM(IID·near-additive)에선 1차·2차 모두 ≈1.000으로 포화해 차이가 안 보인다 → 2차 항의 fidelity 기여는 변위가 큰 CNN 무대에서 드러난다.",
     "CNN 실험 cross-silo N=10 전원 · 정답 = in-run 정확 오라클 (b) · 값 클수록 충실(↑). Flirds(1차+2차) vs Flirds-1st 비교."])
cols = ["데이터셋 / 시나리오", "Flirds (1차+2차) ↑", "Flirds-1st ↑", "차이"]
rows = [
    ["cifar10 / feature_noise", "1.00", "0.89", "+0.11"],
    ["cifar10 / iid",          "0.95", "0.54", "+0.41"],
    ["cifar10 / label_skew",   "0.98", "0.92", "+0.06"],
    ["mnist / feature_noise",  "0.79", "0.70", "+0.09"],
    ["mnist / label_skew",     "0.71", "0.61", "+0.10"],
    ["pool 평균 (10 시나리오 × 3 seed)", "0.919", "0.832", "+0.087"],
]
mb = {(0,1),(1,1),(2,1),(3,1),(4,1),(5,1),(5,0),(5,3)}
add_table(s, Inches(2.55), cols, rows, col_widths([3.4, 2.2, 1.9, 1.4]), 11.5, manual_bold=mb, row_h=Inches(0.36))
add_textbox(s, Inches(0.45), Inches(5.45), Inches(12.4), Inches(1.05),
    [("• CNN(전체 모델 학습 → 라운드 변위 큼)에선 2차 Hessian 항이 모든 시나리오에서 fidelity를 끌어올린다 (pool +0.087).", 12.5, True, C_TITLE),
     ("• 이것이 §오차 분석의 직접 증거 — '변위가 큰 FL 무대에서 2차가 load-bearing'. LLM에서 1·2차가 같아 보이는 건 near-additive 포화 때문.", 12.5, False, C_BODY)],
    line_space=1.25)
slide_footer(s, "../build_deck.py Slide 6·7 (개요 §3.1.2, runs/track_c/fidelity.csv · c1/*/metrics.json)",
             "pool은 신호 없는 iid 칸 포함이라 보수적. ±std·전체 10 시나리오·(a) retrain 오라클 대비는 결과 데크 참조.")
_page_number(s)

# ================================================================ Slide 15: 결과 (2) 비용·확장
s = newslide()
top = slide_header(s, "대표 결과 (2) — 비용: 정확 오라클 대비 5~15× 저렴 + 대규모 확장",
    ["같은 순위 정확도(Fidelity 1.000)를 훨씬 싸게. LLM 1B cross-silo N=5, 라운드당 1 HVP vs 2⁵ 부분집합 평가 · 런타임은 3-seed 대표값."])
cols = ["방법", "런타임(N=5, 1B)", "성격"]
rows = [
    ["Flirds-1st",            "~35 s",   "1차만 — 가장 쌈"],
    ["Flirds (1차+2차)",      "~107 s",  "추정기 본체"],
    ["loss-heur",             "~170 s",  "휴리스틱 floor"],
    ["GTG / FedSV / ShapleyFL / Banzhaf", "~530 s", "부분집합 스윕 baseline"],
    ["(b) in-run 정확 오라클", "~530 s",  "정답(검증용)"],
    ["Ripple",                "~4515 s", "최약·최고가 → 지배됨"],
]
mb = {(0,1),(1,1)}
add_table(s, Inches(2.1), cols, rows, col_widths([4.0,2.3,3.0]), 11.5, manual_bold=mb, row_h=Inches(0.42))
# 확장 카드
cy = Inches(5.15)
card(s, Inches(0.55), cy, Inches(12.2), Inches(1.25), fill=RGBColor(0xE6,0xEF,0xE6), line=RGBColor(0xC4,0xD8,0xC4))
add_textbox(s, Inches(0.8), cy + Inches(0.16), Inches(11.7), Inches(1.0),
    [("대규모 단말 연합(cross-device, N=100)까지 확장", 14, True, RGBColor(0x2C,0x5F,0x2D)),
     ("• N=100, 라운드당 10명 참여(Dirichlet α-sweep)에서도 Flirds vs (b) per-round 오라클 Spearman +1.000 — near-additivity가 스케일에서도 유지.", 12.5, False, C_BODY),
     ("• 오라클은 forward 771ms/회 → R200이면 ~11h/4-GPU(α=0.5 anchor 1점만). Flirds는 이를 검증된 proxy로 대체해 전 α축을 저렴하게 커버.", 12, False, C_BODY)],
    line_space=1.2)
slide_footer(s, "../build_deck.py (개요 §3.1.1 runtime) · CLAUDE.md baseline (task1 runtime · task7c cross-device +1.000)")
_page_number(s)

# ================================================================ Slide 16: 결과 (3) 강건성/탐지
s = newslide()
top = slide_header(s, "대표 결과 (3·선택) — 강건성: 오염 클라 탐지 (cross-silo N=5)",
    ["핵심 질문 위계의 2차③(마지막). 위협별 1명 오염 · AUROC = corrupt를 high-φ로 잡는 탐지력(↑) · Llama-3.2-1B, (b)=exact 2⁵ · 3-seed.",
     "발췌. noisy=answer-swap(데이터 품질), free-rider=zero/random 업데이트, poison=clean-보존 backdoor(ASR≈1.00)."])
cols = ["method", "noisy ↑", "free-rider(random) ↑", "free-rider(zero) ↑", "poison ↑"]
rows = [
    ["Flirds (1차+2차)", "1.000", "1.000", "1.000", "0.917"],
    ["Flirds-1st",       "1.000", "1.000", "1.000", "0.000"],
    ["loss-heur",        "1.000", "1.000", "1.000", "1.000"],
    ["GTG-Shapley",      "1.000", "1.000", "1.000", "1.000"],
    ["(b) oracle",       "1.000", "1.000", "1.000", "1.000"],
    ["FLDetector(탐지기)", "0.750", "1.000", "0.750", "1.000"],
    ["STD-DAGMM(탐지기)",  "0.417", "1.000", "0.250", "0.750"],
]
mb = {(0,4)}
add_table(s, Inches(2.45), cols, rows, col_widths([2.4,1.6,2.3,2.1,1.6]), 11, manual_bold=mb, row_h=Inches(0.34))
add_textbox(s, Inches(0.45), Inches(5.5), Inches(12.4), Inches(1.05),
    [("• noisy·free-rider는 거의 모두 AUROC 1.0 (near-additive 영역) — 전용 탐지기(FLDetector/STD-DAGMM)보다 valuation이 오히려 강함.", 12.5, False, C_BODY),
     ("• poison(clean-보존 backdoor)이 분리점: Flirds-1st는 완전 회피(0.000), 2차 Flirds가 0.917로 버팀 → 다음 장(정직한 한계).", 12.5, True, C_TITLE)],
    line_space=1.25)
slide_footer(s, "../build_deck.py Slide 14 (개요 §3.4.1, runs/phase2_matrix/analysis · RESULTS.md)",
             "탐지기 행은 Spearman 없음(φ 미산출). 전체 8지표·런타임·cross-device α-sweep은 결과 데크 참조.")
_page_number(s)

# ================================================================ Slide 16b: novelty
s = newslide()
top = slide_header(s, "우리 연구의 novelty — 기존에 비어 있던 교집합",
    ["Flirds의 새로움은 '최초의 연합 in-run'이 아니다(그건 Ripple Shapley가 이미). 아래 6가지 속성을 모두 동시에 만족하는 방법이 선행 연구에 없다는 '교집합'이 novelty다.",
     "최신 문헌 스캔(내부 4 검색 + 외부 서베이 교차검증)으로 확인 — 단일 논문이 전체 교집합을 차지하지 않음."])
# 6 속성 칩 그리드 (2행 x 3열)
props = [
    "client-level (데이터포인트 아닌 클라이언트)",
    "in-run (한 학습 궤적, 재학습 없음)",
    "닫힌형 1차+2차 Taylor",
    "2차 HVP 클라이언트 상호작용 항",
    "추가 통신 0 (vanilla FedAvg 위)",
    "LoRA · LLM 스케일",
]
gx0 = Inches(0.55); gw = Inches(3.95); ggap = Inches(0.13); gy = Inches(2.15); gh = Inches(0.72)
for i, p in enumerate(props):
    r, c = divmod(i, 3)
    x = gx0 + c * (gw + ggap)
    y = gy + r * (gh + Inches(0.13))
    card(s, x, y, gw, gh, fill=RGBColor(0xE6,0xEF,0xE6), line=RGBColor(0xC4,0xD8,0xC4))
    add_textbox(s, x + Inches(0.18), y, gw - Inches(0.36), gh,
                [(p, 12, True, RGBColor(0x2C,0x5F,0x2D))], anchor=MSO_ANCHOR.MIDDLE)
# 차별점: 가장 가까운 경쟁자 대비
dy = gy + 2 * (gh + Inches(0.13)) + Inches(0.18)
add_textbox(s, Inches(0.55), dy, Inches(12.2), Inches(0.36),
    [("가장 가까운 경쟁자 대비 차별점:", 13, True, C_TITLE)])
cols = ["방법", "한계 (Flirds가 채우는 지점)"]
rows = [
    ["Ripple Shapley", "sample-level · 재귀 Jacobian 저랭크 근사 · 2차항 없음 → Flirds는 client-level·닫힌형 1차+2차"],
    ["FedIF (2025)", "1차 TracIn만 · CNN 전용 · 집계 방식을 바꿈 → Flirds는 2차 곡률 + vanilla FedAvg 위 사후 valuation"],
    ["FedSV / GTG / ShapleyFL", "라운드마다 부분집합·재학습 대리 효용 평가 → Flirds는 라운드당 HVP 1회 닫힌형"],
    ["IRDS (중앙집중)", "FL 아님 · per-step 영역이라 2차 이득 미미 → Flirds는 FL 라운드 변위에서 2차가 load-bearing"],
]
add_table(s, dy + Inches(0.42), cols, rows, col_widths([2.6, 9.85]), 10.5, left=Inches(0.55), row_h=Inches(0.42))
slide_footer(s, "research-wiki/wiki/flirds.md (Differentiators · Recent prior-work scan 2026-06-03) · ../build_deck.py Slide 30 (기여)")
_page_number(s)

# ================================================================ Slide 17: 정직한 한계
s = newslide()
top = slide_header(s, "정직한 한계 — valuation의 경계",
    ["방법을 '탐지기'가 아니라 '기여도 측정기'로 정직하게 보면, 다음은 실패가 아니라 정의상 옳은 답이다."])
card(s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(2.5), fill=C_CARD)
add_textbox(s, Inches(0.85), Inches(2.25), Inches(11.6), Inches(2.0),
    [("clean-preserving backdoor (clean 검증 손실을 낮추는 공격자)", 16, True, C_ACCENT),
     ("• 이 공격자는 깨끗한 검증 손실을 실제로 낮춘다 → valuation은 정직하게 '기여 높음(φ 큼)'으로 평가한다.", 13, False, C_BODY),
     ("• 즉 1차 Taylor가 부호를 놓치는 것은 '탐지 실패'라기보다, 검증 손실을 게임 효용으로 쓴 정직한 결과에 가깝다.", 13, False, C_BODY),
     ("• 같은 검증-손실 게임을 쓰는 (b) 정확 오라클·loss-heur도 이를 high-φ로 평가 → 방법론의 일관된 경계이지 버그가 아니다.", 13, False, C_BODY)],
    line_space=1.28)
add_textbox(s, Inches(0.55), Inches(4.75), Inches(12.2), Inches(1.6),
    [("핵심 질문 위계와 일치하는 메시지:", 13, True, C_TITLE),
     ("• 기여도(valuation)와 오염 탐지(detection)는 완전 직결이 아니다 — 그래서 탐지는 위계의 '마지막'에 둔다.", 12.5, False, C_BODY),
     ("• 이 경계는 보편이 아니라 조건부(큰 scaled-update × 큰 모델; cross-device 희석 시 회피 안 됨) → 상보적 탐지기가 필요한 지점을 정직하게 표시.", 12.5, False, C_BODY),
     ("• 2차 Hessian 항이 일부(1B 0.917) 복원한다는 점은 '2차가 FL에서 load-bearing'이라는 본 연구의 주장과도 맞물린다.", 12.5, False, C_BODY)],
    line_space=1.25)
slide_footer(s, "research-wiki/wiki/flirds.md (핵심 질문 위계 §poison) · ../build_deck.py Slide 20·25 (poison 회피 발견)")
_page_number(s)

# ================================================================ Slide 18: 마무리
s = newslide()
dark_bg(s)
add_textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.9),
            [("정리 — 한 문장", 30, True, C_WHITE)])
add_textbox(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.8),
    [("Flirds는 연합학습의 한 학습 궤적에서, 서버가 이미 받는 클라이언트 업데이트 Δwₖ 만으로 "
      "client-level 데이터 기여도를 1차+2차 Taylor(true Hessian)로 추정한다 — 추가 통신 0, "
      "정확 오라클의 순위를 충실히(Spearman ≈ 1.000) 5~15× 저렴하게 재현하며, "
      "FL 특유의 큰 라운드 변위 때문에 2차 항이 비로소 의미를 갖는다.", 17, False, C_ICE)],
    line_space=1.35)
add_textbox(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.6),
    [("발표 흐름 요약", 14, True, C_WHITE),
     ("FL 동기(보상·탐지·불량식별) → Data Shapley(공정 분배, 2ᴺ 한계) → In-Run(한 궤적·Taylor 분해) "
      "→ Flirds(client·round 설계 + 왜 2차가 필요한가) → 결과(Fidelity 1.000 · 5~15× 저렴 · 강건성) → novelty(교집합) → 정직한 한계.",
      12.5, False, C_ICE),
     ("개념 출처 = research-wiki / 대표 수치 출처 = 결과 데크 flirds-results-2026-06-26. 슬라이드별 각주에 경로 표기.",
      10.5, False, C_ICE)],
    line_space=1.3)
_page_number(s)

# ---------------------------------------------------------------- save
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "flirds-seminar-intro.pptx")
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
