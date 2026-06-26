#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Flirds 실험 결과 전체 발표 데크 빌더 (재현용).

출처(수치 단일 소스): research-wiki/survey/flirds-experiment-results-overview-2026-06-25.md
결론 섹션 서술 출처:  research-wiki/survey/flirds-results-analysis-2026-06-26.md
모든 수치는 위 개요 문서에서 그대로 전사. 재집계/재실행 없음.

볼드 규칙(자동): 각 비교 열(또는 행)에서 방향(↑=최대, ↓=최소)의 최고값을 볼드.
  - 오라클/기준((a)/(b)/Flirds-proxy)·설정/산출값 행은 비교에서 제외(eligible=False).
  - 동률이 4개 이상이면 볼드 생략(포화 열) — '동률 다수' 캡션으로 대체.
  - arm/수렴 표는 (scale·stage) 블록별 수동 볼드.

실행: /home/korea_bupj/miniconda3/envs/flirds/bin/python build_deck.py
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- palette (담백한 학술 핸드아웃) ----
C_TITLE   = RGBColor(0x1F, 0x2A, 0x37)   # dark slate
C_HEADER  = RGBColor(0x33, 0x47, 0x5B)   # table header fill
C_HEADTXT = RGBColor(0xFF, 0xFF, 0xFF)
C_ROWALT  = RGBColor(0xF2, 0xF4, 0xF7)
C_BODY    = RGBColor(0x22, 0x22, 0x22)
C_MUTED   = RGBColor(0x7A, 0x7A, 0x7A)
C_SETTING = RGBColor(0x3C, 0x46, 0x54)
C_RULE    = RGBColor(0xC9, 0xD2, 0xDC)
C_BG_DARK = RGBColor(0x1F, 0x2A, 0x37)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_ICE     = RGBColor(0xCA, 0xDC, 0xFC)

FONT = "NanumGothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
_pageno = [0]


def _set_run(r, text, size, bold=False, color=C_BODY, italic=False):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    # ensure east-asian font too
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_textbox(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, line_space=1.0):
    """lines = list of (text, size, bold, color[, italic])"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_space
        txt, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        ital = ln[4] if len(ln) > 4 else False
        _set_run(p.add_run(), txt, size, bold, color, ital)
    return tb


def _page_number(slide):
    _pageno[0] += 1
    if _pageno[0] == 1:
        return
    tb = slide.shapes.add_textbox(SW - Inches(0.7), SH - Inches(0.38),
                                  Inches(0.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _set_run(p.add_run(), str(_pageno[0]), 9, False, C_MUTED)


def slide_header(slide, title, setting_lines):
    add_textbox(slide, Inches(0.45), Inches(0.30), Inches(12.4), Inches(0.8),
                [(title, 24, True, C_TITLE)])
    y = 1.12
    if setting_lines:
        add_textbox(slide, Inches(0.45), Inches(y), Inches(12.4), Inches(0.7),
                    [(s, 11, False, C_SETTING) for s in setting_lines],
                    line_space=1.05)
        y = 1.12 + 0.22 * len(setting_lines) + 0.06
    return Inches(y)


def slide_footer(slide, source, caption=None):
    parts = []
    if caption:
        parts.append((caption, 10, False, C_MUTED, True))
    add_textbox(slide, Inches(0.45), SH - Inches(0.78), Inches(12.4), Inches(0.7),
                ([(caption, 10, False, C_MUTED, True)] if caption else []) +
                [("출처: " + source, 9, False, C_MUTED)],
                line_space=1.0)


def pnum(s):
    if s is None:
        return None
    s = str(s).strip()
    if s in ('–', '-', '⬚', '', '(truth)', '(proxy truth)', '(proxy 기준)', '(anchor만)'):
        return None
    s = s.replace('−', '-')
    s = s.split('±')[0].replace('~', '').replace('s', '')
    s = s.split('–')[0]
    m = re.match(r'-?\d*\.?\d+', s.strip())
    return float(m.group()) if m else None


def compute_bold(columns, rows, bold_axis):
    """columns: list of (header, dir). rows: list of (cells, eligible).
       returns set of (r,c) to bold."""
    bold = set()
    if bold_axis == 'col':
        for c, (_, d) in enumerate(columns):
            if d not in ('up', 'down'):
                continue
            vals = []
            for r, (cells, elig) in enumerate(rows):
                if not elig:
                    continue
                v = pnum(cells[c])
                if v is not None:
                    vals.append((r, v))
            if not vals:
                continue
            best = max(v for _, v in vals) if d == 'up' else min(v for _, v in vals)
            tie = [r for r, v in vals if abs(v - best) < 1e-9]
            if 1 <= len(tie) <= 3:
                for r in tie:
                    bold.add((r, c))
    elif bold_axis == 'row':
        for r, (cells, elig) in enumerate(rows):
            if not elig:
                continue
            d = None
            vals = []
            for c, (_, dd) in enumerate(columns):
                if dd not in ('up', 'down'):
                    continue
                d = dd
                v = pnum(cells[c])
                if v is not None:
                    vals.append((c, v))
            if not vals:
                continue
            best = max(v for _, v in vals) if d == 'up' else min(v for _, v in vals)
            tie = [c for c, v in vals if abs(v - best) < 1e-9]
            if 1 <= len(tie) <= 3:
                for c in tie:
                    bold.add((r, c))
    return bold


def add_table(slide, top, columns, rows, col_widths, font_size,
              bold_axis='col', manual_bold=None, left=Inches(0.45), row_h=Inches(0.3)):
    ncol = len(columns)
    nrow = len(rows) + 1
    total_w = sum(col_widths)
    height = row_h * nrow
    gfx = slide.shapes.add_table(nrow, ncol, left, top, Emu(int(total_w)), height)
    tbl = gfx.table
    # disable banded style
    tbl.first_row = False
    tbl.horz_banding = False
    for c, w in enumerate(col_widths):
        tbl.columns[c].width = Emu(int(w))
    # header
    for c, (hdr, _) in enumerate(columns):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = C_HEADER
        cell.margin_left = Pt(3); cell.margin_right = Pt(3)
        cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        _set_run(p.add_run(), hdr, font_size, True, C_HEADTXT)
    # bold map
    if manual_bold is not None:
        bold = manual_bold
    else:
        bold = compute_bold(columns, rows, bold_axis)
    for r, (cells, _) in enumerate(rows):
        for c in range(ncol):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_ROWALT if (r % 2 == 1) else C_WHITE
            cell.margin_left = Pt(3); cell.margin_right = Pt(3)
            cell.margin_top = Pt(0); cell.margin_bottom = Pt(0)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            is_b = (r, c) in bold or (c == 0 and False)
            _set_run(p.add_run(), str(cells[c]), font_size, is_b, C_BODY)
    for r in range(nrow):
        tbl.rows[r].height = Emu(int(row_h))
    return gfx


def col_widths(weights, total=Inches(12.45)):
    s = sum(weights)
    return [int(total * w / s) for w in weights]


SLIDES = []  # registry of builder funcs in order

# ---------------------------------------------------------------- Slide 1: 표지
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = C_BG_DARK; bg.line.fill.background()
bg.shadow.inherit = False
s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
add_textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2),
            [("Flirds — 연합학습 클라이언트 기여도 측정", 38, True, C_WHITE)])
add_textbox(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(1.0),
            [("연합학습(FL) 한 학습 궤적에서 In-Run Data Shapley(1차+2차 Taylor, true Hessian)로 "
              "client-level 데이터 기여도를 추정한다.", 16, False, C_ICE)], line_space=1.15)
add_textbox(s, Inches(0.9), Inches(4.25), Inches(11.5), Inches(1.4),
            [("스코프: 네 실험 트랙 — LLM 표준(Llama 1B/3B/7B, LoRA) · CNN · "
              "강건성(오염·공격 주입) · 기초 검증", 14, False, C_WHITE),
             ("핵심 질문 위계: 1차 Fidelity → 2차 ①성능 ②수렴 ③탐지 → 비용", 14, False, C_WHITE)],
            line_space=1.3)
add_textbox(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.8),
            [("2026-06-26 · 교수님 미팅 발표", 13, True, C_ICE),
             ("모든 수치 출처: flirds-experiment-results-overview-2026-06-25 (git e89af94, RESULTS↔CSV 교차검증본)",
              10, False, C_ICE)], line_space=1.2)
_page_number(s)

# ---------------------------------------------------------------- Slide 2: 실험 지도
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "실험 지도 — 무엇을 돌렸나 (마스터 표)",
                   ["네 트랙에서 실제로 돌아 디스크에 남은 실험(실행 15개). 계획·미실행 6행(P1–P6)은 생략."])
cols = [("# 실험", None), ("모델 · N · 참여", None), ("검증 목적", None),
        ("seeds", None), ("status", None)]
rows = [
 (["1  LLM 표준 1B 표준스테이지", "Llama-3.2-1B · N=20 · 2/20", "Fidelity·성능·수렴·비용", "3", "● 실측"], False),
 (["2  LLM 표준 1B 정밀스테이지", "Llama-3.2-1B · N=5 · full", "Fidelity·성능·수렴·비용", "3", "● 실측"], False),
 (["3  LLM 표준 3B 표준스테이지", "Llama-3.2-3B · N=20 · 2/20", "Fidelity·성능·수렴·비용", "3", "● 실측"], False),
 (["4  LLM 표준 3B 정밀스테이지", "Llama-3.2-3B · N=5 · full", "Fidelity·성능·수렴·비용", "3", "● 실측"], False),
 (["5  LLM 표준 7B 표준스테이지", "Llama-2-7B · N=20 · 2/20", "Fidelity·성능·수렴·비용", "3", "● 실측"], False),
 (["6  LLM 표준 7B 정밀스테이지", "Llama-2-7B · N=5 · full", "Fidelity·성능·수렴·비용", "3", "● 실측"], False),
 (["7  CNN cross-silo N=10", "소형 CNN · N=10 · full", "Fidelity·안정성·비용", "3", "● 실측"], False),
 (["8  CNN cross-device N=100", "소형 CNN · N=100 · 10/100", "성능·탐지·수렴", "3", "● 실측"], False),
 (["9  강건성 1B cross-silo N=5", "Llama-3.2-1B · N=5 · full", "Fidelity·탐지·비용", "3", "● 실측"], False),
 (["10 강건성 1B cross-device α-sweep", "Llama-3.2-1B · N=100 · 10/100", "탐지·Fidelity", "3", "● 실측"], False),
 (["11 강건성 1B cross-device α=0.5 anchor", "Llama-3.2-1B · N=100 · 10/100", "Fidelity·탐지·비용", "3", "● 실측"], False),
 (["12 강건성 1B cross-device poison", "Llama-3.2-1B · N=100 · 10/100", "탐지·Fidelity", "3", "● 실측"], False),
 (["13 강건성 3B cross-silo N=5", "Llama-3.2-3B · N=5 · full", "Fidelity·탐지", "1", "◐ 부분"], False),
 (["14 기초 1B 첫 clean run", "Llama-3.2-1B · N=5 · full", "탐지·성능(selection)", "3×2lr", "● 부록"], False),
 (["15 기초 1B LR sweep", "Llama-3.2-1B · N=5 · full", "탐지·성능(selection)", "1×4lr", "● 부록"], False),
]
add_table(s, top, cols, rows, col_widths([3.5, 2.7, 2.4, 0.9, 1.0]), 9.5)
slide_footer(s, "flirds-experiment-results-overview-2026-06-25 §2",
             "마커: ● 실측 · ◐ 부분(1-seed 등) · ⬚ 미실행.  오라클: (a)=조합마다 재학습한 retrain 정확 오라클, "
             "(b)=한 궤적에서 2^N 분해한 in-run 정확 오라클.  참여 k/N=라운드당 k명, full=전원.")
_page_number(s)

# ---------------------------------------------------- Slide 3: LLM std20 순위·값 상관
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "Fidelity (1차) — LLM 표준, 표준 스테이지: 순위·값 상관",
    ["LLM 표준 실험(LoRA; OpenFedLLM 표준 레시피) · 표준 스테이지(N=20, 라운드당 2). alpaca-gpt4 20k IID, R=200, lr=1e-3, SGD mom=0, fp32 · 3-seed mean±std",
     "정답 = in-run 정확 오라클 (b). 값 클수록 충실(↑)."])
cols = [("method", None), ("1B Sp ↑", 'up'), ("1B Ke ↑", 'up'), ("1B Pe ↑", 'up'),
        ("3B Sp ↑", 'up'), ("3B Ke ↑", 'up'), ("3B Pe ↑", 'up'),
        ("7B Sp ↑", 'up'), ("7B Ke ↑", 'up'), ("7B Pe ↑", 'up')]
rows = [
 (["Flirds","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","0.999±.001","0.996±.005","1.000±.000"], True),
 (["Flirds-1st","0.999±.001","0.996±.005","1.000±.000","1.000±.000","1.000±.000","1.000±.000","0.998±.001","0.986±.010","1.000±.000"], True),
 (["loss-heur","1.000±.000","1.000±.000","1.000±.000","0.999±.001","0.996±.005","1.000±.000","0.999±.001","0.996±.005","1.000±.000"], True),
 (["GTG","0.975±.018","0.916±.043","0.995±.001","0.988±.005","0.944±.022","0.996±.000","0.977±.017","0.916±.045","0.989±.006"], True),
 (["FedSV","0.910±.073","0.786±.117","0.959±.013","0.952±.018","0.853±.026","0.972±.004","0.968±.010","0.881±.030","0.976±.006"], True),
 (["FedIF","0.157±.303","0.111±.199","0.229±.222","0.211±.184","0.139±.115","0.262±.137","0.480±.101","0.323±.061","0.508±.054"], True),
 (["ShapleyFL","0.194±.351","0.133±.244","0.245±.283","0.227±.143","0.161±.092","0.246±.143","0.406±.081","0.274±.054","0.431±.026"], True),
 (["ComFedSV","0.093±.146","0.060±.108","0.095±.193","-0.129±.066","-0.105±.034","-0.093±.038","0.039±.171","0.039±.110","0.048±.115"], True),
]
add_table(s, top, cols, rows, col_widths([1.5]+[1.0]*9), 8.5)
slide_footer(s, "개요 §3.1.1 (runs/track_d/fidelity.csv)",
             "Sp=Spearman, Ke=Kendall, Pe=Pearson. Flirds·Flirds-1st·loss-heur가 천장(≈1.000); FedIF/ShapleyFL/ComFedSV는 설계상 낮음(영향도·surrogate·low-rank 가정).")
_page_number(s)

# ---------------------------------------------------- Slide 4: LLM anchor5 순위·값 상관
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "Fidelity (1차) — LLM 표준, 정밀 스테이지: 순위·값 상관",
    ["LLM 표준 실험 · 정밀 스테이지(N=5, 전원 참여). R=30, 그 외 동일 · 3-seed mean±std · 정답 = in-run 정확 오라클 (b)",
     "맨 아래 (a)oracle 행 = retrain 정확 오라클 (a)와 (b)의 일치도(1B만 실행); 비교 대상 아님."])
cols = [("method", None), ("1B Sp ↑", 'up'), ("1B Ke ↑", 'up'), ("1B Pe ↑", 'up'),
        ("3B Sp ↑", 'up'), ("3B Ke ↑", 'up'), ("3B Pe ↑", 'up'),
        ("7B Sp ↑", 'up'), ("7B Ke ↑", 'up'), ("7B Pe ↑", 'up')]
rows = [
 (["Flirds","1.000±.000","1.000±.000","1.000±.000","0.967±.047","0.933±.094","1.000±.000","1.000±.000","1.000±.000","1.000±.000"], True),
 (["Flirds-1st","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000"], True),
 (["loss-heur","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000"], True),
 (["Banzhaf","1.000±.000","1.000±.000","1.000±.000","0.967±.047","0.933±.094","1.000±.000","1.000±.000","1.000±.000","1.000±.000"], True),
 (["GTG","1.000±.000","1.000±.000","1.000±.000","1.000±.000","1.000±.000","0.999±.000","1.000±.000","1.000±.000","1.000±.000"], True),
 (["FedSV","0.700±.163","0.600±.163","0.824±.116","0.700±.000","0.600±.000","0.882±.056","0.933±.047","0.867±.094","0.961±.024"], True),
 (["ShapleyFL","0.700±.283","0.600±.283","0.764±.258","0.100±.000","0.000±.000","0.353±.415","0.833±.125","0.733±.189","0.903±.067"], True),
 (["ComFedSV","0.500±.432","0.467±.340","0.563±.356","0.600±.294","0.533±.249","0.475±.272","0.600±.216","0.467±.189","0.588±.256"], True),
 (["FedIF","0.067±.531","0.067±.411","-0.068±.626","0.067±.492","0.000±.432","0.335±.475","0.200±.616","0.200±.490","0.368±.509"], True),
 (["(a)oracle","0.933±.047","0.867±.094","0.933±.054","⬚","⬚","⬚","⬚","⬚","⬚"], False),
]
add_table(s, top, cols, rows, col_widths([1.5]+[1.0]*9), 8.5)
slide_footer(s, "개요 §3.1.1 (runs/track_d/fidelity.csv)",
             "N=5 전수(2^5). 다수 방법이 1.000 동률(볼드 생략 열 있음). (a)oracle vs (b) = 0.933 → 두 정답 정의가 거의 일치(3B/7B (a)는 ⬚ 미실행).")
_page_number(s)

# ---------------------------------------------------- Slide 5: anchor5 vs (a) retrain
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "Fidelity (1차) — LLM 정밀 스테이지: retrain 오라클(a) 대비",
    ["LLM 표준 1B 정밀 스테이지(N=5) · 모든 방법을 retrain 정확 오라클 (a)와 직접 비교 · 3-seed mean±std",
     "(a) = 조합마다 처음부터 재학습한 2^5 Shapley(문헌 공백). 맨 오른쪽 열은 참고용 vs (b)."])
cols = [("method", None), ("Sp vs (a) ↑", 'up'), ("Ke vs (a) ↑", 'up'),
        ("Pe vs (a) ↑", 'up'), ("max_diff vs (a) ↓", 'down'), ("(참고) Sp vs (b) ↑", None)]
rows = [
 (["Flirds","0.933±.047","0.867±.094","0.933±.055",".001","1.000"], True),
 (["Flirds-1st","0.933±.047","0.867±.094","0.929±.060",".001","1.000"], True),
 (["loss-heur","0.933±.047","0.867±.094","0.931±.057",".001","1.000"], True),
 (["Banzhaf","0.933±.047","0.867±.094","0.933±.054",".001","1.000"], True),
 (["GTG","0.933±.047","0.867±.094","0.937±.052",".002","1.000"], True),
 (["FedSV","0.733±.170","0.600±.163","0.685±.249",".003","0.700"], True),
 (["ShapleyFL","0.767±.330","0.733±.377","0.916±.084",".983","0.700"], True),
 (["ComFedSV","0.467±.450","0.467±.411","0.598±.280",".014","0.500"], True),
 (["FedIF","0.167±.613","0.200±.490","0.048±.585",".984","0.067"], True),
]
add_table(s, top, cols, rows, col_widths([1.5, 1.6, 1.6, 1.6, 1.9, 1.9]), 11)
slide_footer(s, "개요 §3.1.1 (runs/track_d/rundirs/1B_anchor5_seed*/phi.parquet)",
             "(b)와 거의 일치하는 방법들은 vs (a)도 0.933 동률(천장 효과 = (b)-vs-(a) 일치도 0.933이 상한).")
_page_number(s)

# ---------------------------------------------------- Slide 6: CNN pool 듀얼 오라클
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "Fidelity (1차) — CNN 실험: 듀얼 오라클 pool 평균",
    ["CNN 실험 cross-silo N=10 전원 · mnist=LeNet5 / cifar10=FedSVCNN(전체 모델 학습) · 10 시나리오 × 3 seed pool",
     "정답 둘 다 실측: (b) in-run 2^10 · (a) retrain 2^10. 값 클수록 충실(↑)."])
cols = [("method", None), ("vs (b) Sp ↑", 'up'), ("vs (b) Pe ↑", 'up'),
        ("vs (a) Sp ↑", 'up'), ("vs (a) Pe ↑", 'up')]
rows = [
 (["Flirds","0.919±.134","0.934±.128","0.352±.462","0.354±.461"], True),
 (["Flirds-1st","0.832±.194","0.853±.159","0.408±.435","0.412±.421"], True),
 (["loss-heur","0.860±.154","0.885±.134","0.425±.429","0.423±.408"], True),
 (["Banzhaf","0.989±.019","0.998±.004","0.355±.441","0.357±.459"], True),
 (["GTG","0.569±.343","0.612±.317","0.374±.412","0.332±.456"], True),
 (["FedSV","0.401±.410","0.410±.406","0.284±.479","0.215±.466"], True),
 (["ComFedSV","0.348±.377","0.328±.396","0.338±.398","0.309±.431"], True),
 (["ShapleyFL","0.391±.385","0.392±.425","0.453±.380","0.443±.410"], True),
 (["FedIF","0.491±.391","0.506±.427","0.380±.393","0.368±.431"], True),
 (["Ripple","0.373±.444","0.404±.437","0.213±.462","0.158±.470"], True),
]
add_table(s, top, cols, rows, col_widths([2.0, 2.6, 2.6, 2.6, 2.6]), 11)
slide_footer(s, "개요 §3.1.2 (runs/track_c/fidelity.csv)",
             "pool은 iid 셀 포함 → 깎임(iid 제외 시 Flirds vs (b) 0.928). LLM과 달리 CNN은 clean에서도 큰 spread. (a)와 (b)는 갈림(어떤 방법도 vs (a) 0.45 미만).")
_page_number(s)

# ---------------------------------------------------- Slide 7: CNN 시나리오별 vs (b)
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "Fidelity (1차) — CNN 실험: 시나리오별 vs (b) Spearman",
    ["CNN 실험 cross-silo N=10 · 10 시나리오(데이터셋 × skew/flip/noise) · 3-seed 평균 · 값 클수록 충실(↑)",
     "행별 최고 Spearman을 볼드(셀 우열)."])
cols = [("데이터셋 / 시나리오", None), ("Flirds", 'up'), ("Flirds-1st", 'up'), ("GTG", 'up'),
        ("FedSV", 'up'), ("ComFedSV", 'up'), ("Banzhaf", 'up'), ("ShapleyFL", 'up'),
        ("FedIF", 'up'), ("loss-heur", 'up'), ("Ripple", 'up')]
rows = [
 (["cifar10 / feature_noise","1.00","0.89","0.59","0.40","0.19","1.00","0.20","0.62","0.90","-0.01"], True),
 (["cifar10 / iid","0.95","0.54","0.21","0.22","0.12","0.98","0.18","0.45","0.69","0.31"], True),
 (["cifar10 / label_flip","1.00","0.95","0.64","0.54","0.31","1.00","0.37","0.74","0.95","0.32"], True),
 (["cifar10 / label_skew","0.98","0.92","0.49","0.53","0.31","1.00","0.29","0.68","0.88","0.26"], True),
 (["cifar10 / quantity_skew","0.99","0.96","0.78","0.56","0.67","1.00","0.44","-0.20","0.98","0.68"], True),
 (["mnist / feature_noise","0.79","0.70","0.41","0.13","0.21","0.95","0.48","0.57","0.78","0.00"], True),
 (["mnist / iid","0.81","0.78","0.47","0.04","0.10","0.98","0.47","0.73","0.84","0.18"], True),
 (["mnist / label_flip","1.00","0.99","0.99","0.97","0.95","1.00","0.98","0.98","0.99","0.97"], True),
 (["mnist / label_skew","0.71","0.61","0.33","-0.01","0.14","0.98","-0.02","0.41","0.63","0.06"], True),
 (["mnist / quantity_skew","0.96","0.98","0.78","0.63","0.49","1.00","0.52","-0.07","0.96","0.96"], True),
]
add_table(s, top, cols, rows, col_widths([2.3]+[1.02]*10), 8.5, bold_axis='row')
slide_footer(s, "개요 §3.1.2 (runs/track_c/c1/*/metrics.json)",
             "Banzhaf(exact semivalue)·Flirds가 대부분 시나리오 최강. 신호 없는 iid 칸은 fidelity가 의미상 낮음.")
_page_number(s)

# ---------------------------------------------------- Slide 8: CNN 데이터셋별 평균
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "Fidelity (1차) — CNN 실험: 데이터셋별 평균",
    ["CNN 실험 cross-silo N=10 · 시나리오 표를 데이터셋으로 묶음(각 5 시나리오 × 3 seed) · 값 클수록 충실(↑)",
     "행별 최고를 볼드. vs (b)=in-run 오라클, vs (a)=retrain 오라클."])
cols = [("데이터셋", None), ("기준", None), ("Flirds", 'up'), ("Flirds-1st", 'up'), ("GTG", 'up'),
        ("FedSV", 'up'), ("ComFedSV", 'up'), ("Banzhaf", 'up'), ("ShapleyFL", 'up'),
        ("FedIF", 'up'), ("loss-heur", 'up'), ("Ripple", 'up')]
rows = [
 (["cifar10","vs (b) Sp","0.98","0.85","0.54","0.45","0.32","1.00","0.30","0.46","0.88","0.31"], True),
 (["mnist","vs (b) Sp","0.85","0.81","0.60","0.35","0.38","0.98","0.49","0.52","0.84","0.44"], True),
 (["cifar10","vs (a) Sp","0.26","0.29","0.37","0.30","0.40","0.27","0.30","0.20","0.34","-0.01"], True),
 (["mnist","vs (a) Sp","0.44","0.53","0.38","0.27","0.27","0.44","0.60","0.56","0.51","0.44"], True),
]
add_table(s, top, cols, rows, col_widths([1.3, 1.2]+[0.99]*10), 9, bold_axis='row')
slide_footer(s, "개요 §3.1.2 (runs/track_c/fidelity.csv)",
             "vs (b)는 Banzhaf/Flirds 우위. vs (a)는 전 방법이 낮고(≤0.60) 순위가 뒤바뀜 — retrain 오라클이 in-run과 다른 게임임을 시사.")
_page_number(s)

# ---------------------------------------------------- Slide 9: 개입 arm 정의
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "개입 arm 정의 — 측정한 기여도로 어떻게 개입하나",
    ["온라인 점수기가 라운드별 기여도를 EMA 누적(s ← β·s + (1−β)·raw) → 다음 라운드 FedAvg 가중을 바꿈.",
     "각 baseline은 자기 논문 방식(가중식 + 점수원 + β)을 그대로 사용 → 공정 비교."])
cols = [("arm", None), ("개입 메커니즘", None), ("관행 / 점수원", None)]
rows = [
 (["base", "학습 전 베이스 모델 (개입 없음)", "기준점"], False),
 (["vanilla", "표준 FedAvg (데이터 크기 n 가중)", "기준"], False),
 (["flirds_w", "곱셈 가중  w ∝ n · s   (EMA β=0.5)", "Flirds 기본 (Yonghee 규칙)"], False),
 (["flirds_sel", "softmax(s/T)로 k명 선택 (cohort 진부분집합=표준 스테이지만)", "S-FedAvg 관행"], False),
 (["shapleyfl_w", "교체 가중  w ∝ s   (β=0.3)", "ShapleyFL 관행, 점수=per-round exact Shapley"], False),
 (["fedif_w", "교체 가중  w ∝ s   (β=0.7=1−γ)", "FedIF 관행, 점수=per-round 1차 influence"], False),
]
add_table(s, top, cols, rows, col_widths([1.5, 6.2, 4.75]), 12)
add_textbox(s, Inches(0.45), top + Inches(2.7), Inches(12.4), Inches(1.6),
    [("가중 규칙 4종(점수원과 무관한 결합식):", 12, True, C_TITLE),
     ("• multiplicative  w ∝ n·s — FedAvg 크기 가중에 기여도를 곱함 (Flirds 기본)", 11.5, False, C_BODY),
     ("• replacement  w ∝ s — n 가중을 기여도로 대체 (FedIF·ShapleyFL 관행)", 11.5, False, C_BODY),
     ("• additive  w = λ·s/Σs + (1−λ)·n/Σn, λ=0.5 — 혼합 (Ripple 관행, CNN flirds_add)", 11.5, False, C_BODY),
     ("• selection  softmax(s/T)로 k명 비복원 샘플 (S-FedAvg 관행)", 11.5, False, C_BODY),
     ("주의: 모든 n_i가 같으면 multiplicative==replacement → IID 무대에선 flirds_w·shapleyfl_w 가중식이 같고 점수원·β만 다름.", 10.5, False, C_MUTED, True)],
    line_space=1.18)
slide_footer(s, "개요 §3.2.1 (codes/flirds/fl/intervene.py)")
_page_number(s)

# ---------------------------------------------------- Slide 10: LLM arms MMLU/ROUGE
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "성능 / 집계 (2차①) — LLM 개입 arm: MMLU · ROUGE-L",
    ["LLM 표준 실험, §Fidelity 동일 궤적 · MMLU full-test(14,042) 0-shot + Alpaca-test(1k) ROUGE-L · 3-seed mean±std · 좌=표준/우=정밀 스테이지 · arm 중 최고 블록별 볼드(base 제외)."])
cols = [("스케일·스테이지", None), ("arm", None), ("MMLU ↑", None), ("ROUGE-L ↑", None),
        ("스케일·스테이지", None), ("arm", None), ("MMLU ↑", None), ("ROUGE-L ↑", None)]
rows = [
 (["1B 표준","base","0.4822±.0000","0.2168±.0019","1B 정밀","base","0.4822±.0000","0.2168±.0019"], False),
 (["","vanilla","0.4742±.0001","0.2841±.0051","","vanilla","0.4801±.0003","0.2725±.0032"], True),
 (["","flirds_w","0.4745±.0003","0.2848±.0050","","flirds_w","0.4802±.0007","0.2741±.0025"], True),
 (["","flirds_sel","0.4739±.0005","0.2838±.0041","","shapleyfl_w","0.4802±.0007","0.2741±.0026"], True),
 (["","shapleyfl_w","0.4742±.0005","0.2845±.0050","","fedif_w","0.4797±.0008","0.2713±.0037"], True),
 (["","fedif_w","0.4741±.0003","0.2847±.0046","","","",""], True),
 (["3B 표준","base","0.6230±.0000","0.2219±.0015","3B 정밀","base","0.6230±.0000","0.2219±.0015"], False),
 (["","vanilla","0.6147±.0006","0.3017±.0024","","vanilla","0.6215±.0001","0.2749±.0035"], True),
 (["","flirds_w","0.6137±.0006","0.3015±.0018","","flirds_w","0.6214±.0002","0.2755±.0042"], True),
 (["","flirds_sel","0.6139±.0014","0.3029±.0039","","shapleyfl_w","0.6214±.0002","0.2755±.0042"], True),
 (["","shapleyfl_w","0.6136±.0005","0.3016±.0018","","fedif_w","0.6213±.0002","0.2730±.0035"], True),
 (["","fedif_w","0.6139±.0007","0.3022±.0024","","","",""], True),
 (["7B 표준","base","0.4175±.0000","0.1496±.0024","7B 정밀","base","0.4175±.0000","0.1496±.0024"], False),
 (["","vanilla","0.4038±.0024","0.2778±.0026","","vanilla","0.4206±.0012","0.1651±.0016"], True),
 (["","flirds_w","0.4026±.0028","0.2780±.0027","","flirds_w","0.4210±.0014","0.1680±.0023"], True),
 (["","flirds_sel","0.4025±.0022","0.2790±.0044","","shapleyfl_w","0.4210±.0014","0.1680±.0023"], True),
 (["","shapleyfl_w","0.4027±.0027","0.2787±.0028","","fedif_w","0.4204±.0008","0.1656±.0011"], True),
 (["","fedif_w","0.4030±.0023","0.2763±.0033","","","",""], True),
]
mb = {(2,2),(2,3),(2,6),(3,6),(2,7),(3,7),
      (7,2),(9,3),(7,6),(8,7),(9,7),
      (13,2),(15,3),(14,6),(15,6),(14,7),(15,7)}
add_table(s, top, cols, rows, col_widths([1.5,1.3,1.55,1.55,1.5,1.3,1.55,1.55]), 7, manual_bold=mb, row_h=Inches(0.21))
slide_footer(s, "개요 §3.2.1 (runs/track_d/rundirs/*/metrics.json)",
             "모든 개입 arm이 vanilla와 ±0.001~0.003 이내 = clean-IID parity(향상 아님). 학습 자체는 base 대비 ROUGE 크게↑, MMLU는 SFT로 분포-밖 소폭↓. 정밀 스테이지는 전원 참여라 flirds_sel 없음.")
_page_number(s)

# ---------------------------------------------------- Slide 11: CNN C2 정확도
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "성능 / 집계 (2차①) — CNN cross-device: 최종 정확도",
    ["CNN 실험 cross-device N=100, 라운드당 10% 참여, R=120 · cifar10/fmnist × iid/dir1/shard × 위협(clean/label_flip/free_rider/grad_noise) · 30 셀 · 3 seed",
     "값=최종 test 정확도(↑). 열=위협 그룹(partition·강도·dataset pool, 셀 수 표기). 열별 최고 arm 볼드."])
cols = [("arm", None), ("clean (6셀) ↑", 'up'), ("free_rider (6셀) ↑", 'up'),
        ("grad_noise (8셀) ↑", 'up'), ("label_flip (10셀) ↑", 'up')]
rows = [
 (["vanilla","0.686±.127","0.646±.146","0.499±.241","0.583±.184"], True),
 (["flirds_mult","0.698±.122","0.662±.144","0.609±.187","0.626±.170"], True),
 (["flirds_repl (dir1)","0.734±.096","0.704±.111","0.621±.185","0.652±.149"], True),
 (["flirds_add (dir1)","0.733±.094","0.702±.110","0.604±.195","0.635±.161"], True),
 (["flirds_select","0.679±.150","0.656±.148","0.548±.231","0.618±.172"], True),
 (["shapleyfl","0.702±.126","0.645±.136","0.645±.183","0.622±.168"], True),
 (["fedif","0.685±.127","0.654±.154","0.624±.178","0.623±.169"], True),
 (["sfedavg","0.695±.128","0.655±.139","0.510±.252","0.598±.186"], True),
]
add_table(s, top, cols, rows, col_widths([2.2, 2.55, 2.55, 2.6, 2.55]), 11)
slide_footer(s, "개요 §3.2.2 (runs/track_c/c2/*/metrics.json)",
             "그룹 평균이라 std 큼(셀별 30칸은 RESULTS.txt). 오염(grad_noise/label_flip)에서 기여도-가중이 정확도 회복; clean은 parity~소폭↑. flirds_repl/add는 size-skew(dir1) 전용.")
_page_number(s)

# ---------------------------------------------------- Slide 12: LLM 수렴
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "수렴 (2차②) — LLM 표준: final val-loss · rounds-to-target",
    ["LLM 표준 실험, §성능과 동일 로그 · rounds-to-target = vanilla 최종 val-loss에 도달한 첫 라운드 · 3-seed mean±std",
     "좌=표준 스테이지(R=200), 우=정밀 스테이지(R=30). 둘 다 작을수록 좋음(↓). arm 중 최저를 블록별 볼드."])
cols = [("스케일·스테이지", None), ("arm", None), ("val-loss ↓", None), ("rounds ↓", None),
        ("스케일·스테이지", None), ("arm", None), ("val-loss ↓", None), ("rounds ↓", None)]
rows = [
 (["1B 표준","vanilla","1.2653±.0216","200.0±0.0","1B 정밀","vanilla","1.2977±.0209","30.0±0.0"], True),
 (["","flirds_w","1.2653±.0215","199.0±1.0","","flirds_w","1.2964±.0209","29.7±0.5"], True),
 (["","flirds_sel","1.2654±.0217","199.0±0.0","","shapleyfl_w","1.2964±.0209","29.7±0.5"], True),
 (["","shapleyfl_w","1.2653±.0215","199.0±1.0","","fedif_w","1.2976±.0204","30.0±0.0"], True),
 (["","fedif_w","1.2652±.0215","198.0±0.0","","","",""], True),
 (["3B 표준","vanilla","1.1483±.0270","198.3±2.4","3B 정밀","vanilla","1.1970±.0272","30.0±0.0"], True),
 (["","flirds_w","1.1479±.0271","192.3±3.3","","flirds_w","1.1961±.0272","30.0±0.0"], True),
 (["","flirds_sel","1.1481±.0273","193.5±0.5","","shapleyfl_w","1.1960±.0272","30.0±0.0"], True),
 (["","shapleyfl_w","1.1479±.0271","192.3±3.3","","fedif_w","1.1970±.0269","30.0±0.0"], True),
 (["","fedif_w","1.1478±.0271","191.0±3.6","","","",""], True),
 (["7B 표준","vanilla","1.0357±.0244","184.7±18.3","7B 정밀","vanilla","1.0941±.0221","30.0±0.0"], True),
 (["","flirds_w","1.0348±.0245","153.0±18.8","","flirds_w","1.0904±.0224","28.3±0.5"], True),
 (["","flirds_sel","1.0351±.0243","153.0±23.6","","shapleyfl_w","1.0904±.0224","28.3±0.5"], True),
 (["","shapleyfl_w","1.0347±.0245","151.3±19.9","","fedif_w","1.0932±.0207","29.5±0.5"], True),
 (["","fedif_w","1.0348±.0246","158.0±12.7","","","",""], True),
]
mb = {(4,2),(4,3),(1,6),(2,6),(1,7),(2,7),
      (9,2),(9,3),(7,6),
      (13,2),(13,3),(11,6),(12,6),(11,7),(12,7)}
add_table(s, top, cols, rows, col_widths([1.5,1.3,1.55,1.55,1.5,1.3,1.55,1.55]), 9, manual_bold=mb)
slide_footer(s, "개요 §3.3.1 (runs/track_d/rundirs/*/metrics.json)",
             "clean-IID라 arm 간 거의 동률. 두드러진 칸은 7B 표준: 개입 arm이 vanilla 184.7→~151~158 라운드(~14~18% 빠름, 단 std 겹침). 정밀 3B는 전부 30.0 동률.")
_page_number(s)

# ---------------------------------------------------- Slide 13: CNN C2 rounds
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "수렴 (2차②) — CNN cross-device: rounds-to-target",
    ["CNN 실험 cross-device N=100 · 값=목표 정확도 도달 라운드(↓) · 열=위협 그룹(셀 pool) · 열별 최저 arm 볼드"])
cols = [("arm", None), ("clean ↓", 'down'), ("free_rider ↓", 'down'),
        ("grad_noise ↓", 'down'), ("label_flip ↓", 'down')]
rows = [
 (["vanilla","35.7±27.3","41.2±39.1","27.2±35.0","19.3±12.1"], True),
 (["flirds_mult","34.1±27.7","38.8±34.8","13.7±13.6","31.3±30.9"], True),
 (["flirds_repl (dir1)","41.5±34.9","7.7±0.5","6.3±0.5","10.9±4.0"], True),
 (["flirds_add (dir1)","39.5±32.8","34.2±44.9","6.8±0.7","11.8±4.0"], True),
 (["flirds_select","35.3±26.9","50.4±43.2","21.7±24.4","16.1±7.7"], True),
 (["shapleyfl","35.6±28.1","17.4±11.7","18.7±26.1","33.1±35.0"], True),
 (["fedif","36.0±28.7","44.7±35.8","23.9±24.5","30.6±30.4"], True),
 (["sfedavg","35.3±28.3","42.4±38.0","21.9±24.5","19.3±10.6"], True),
]
add_table(s, top, cols, rows, col_widths([2.2, 2.55, 2.55, 2.6, 2.55]), 11)
slide_footer(s, "개요 §3.3.2 (runs/track_c/c2/*/metrics.json)",
             "그룹 평균(target 미달 셀에서 분산 큼). grad_noise에서 기여도-가중이 도달을 크게 앞당김(vanilla 27.2 → flirds_repl 6.3).")
_page_number(s)

# ---------------------------------------------------- Slide 14: silo5 N=5 1B
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "탐지 (2차③) — 강건성 cross-silo N=5 (1B)",
    ["소수기관 연합 cross-silo, N=5 전원 · Llama-3.2-1B, R=10 · 위협별 1명 오염(noisy/free-rider random·zero/poison) · (b)=exact 2^5 · 3-seed",
     "AUROC = corrupt를 high-φ로 잡는 탐지력(↑). Sp = Spearman vs (b)(↑). runtime ↓. poison = clean-보존 backdoor(ASR≈1.00)."])
cols = [("method", None), ("noisy AUROC↑", 'up'), ("noisy Sp↑", 'up'),
        ("frrand AUROC↑", 'up'), ("frrand Sp↑", 'up'), ("frzero AUROC↑", 'up'),
        ("frzero Sp↑", 'up'), ("poison AUROC↑", 'up'), ("poison Sp↑", 'up'), ("runtime↓", 'down')]
rows = [
 (["Flirds","1.000","1.000","1.000","1.000","1.000","1.000","0.917±.118","0.967±.047","~107s"], True),
 (["Flirds-1st","1.000","1.000","1.000","1.000","1.000","1.000","0.000","0.000","~35s"], True),
 (["loss-heur","1.000","1.000","1.000","1.000","1.000","1.000","1.000","1.000","~170s"], True),
 (["FedIF","1.000","0.933±.05","1.000","0.900±.08","1.000","0.933±.05","1.000","0.967±.05","~37s"], True),
 (["GTG","1.000","1.000","1.000","1.000","1.000","1.000","1.000","0.867±.12","~540s"], True),
 (["FedSV","1.000","1.000","1.000","0.933±.05","1.000","1.000","1.000","0.367±.26","~535s"], True),
 (["ShapleyFL","1.000","1.000","1.000","1.000","1.000","1.000","1.000","1.000","~530s"], True),
 (["Banzhaf","1.000","1.000","1.000","1.000","1.000","1.000","1.000","1.000","~535s"], True),
 (["(b)oracle","1.000","(truth)","1.000","(truth)","1.000","(truth)","1.000","(truth)","~530s"], False),
 (["FLDetector","0.750","–","1.000","–","0.750","–","1.000","–","~30s"], True),
 (["STD-DAGMM","0.417±.31","–","1.000","–","0.250±.20","–","0.750±.20","–","~120s"], True),
 (["FLTrust","1.000","–","1.000","–","1.000","–","1.000","–","~37s"], True),
 (["FedDQC","0.917±.12","–","0.750","–","0.750","–","1.000","–","~22s"], True),
]
add_table(s, top, cols, rows, col_widths([1.5,1.28,1.05,1.28,1.05,1.28,1.05,1.28,1.05,1.0]), 8)
slide_footer(s, "개요 §3.4.1 (runs/phase2_matrix/analysis · RESULTS.md)",
             "noisy·free-rider는 거의 전부 AUROC 1.0(near-additive). poison(clean-보존 backdoor)이 분리점: Flirds-1st AUROC/Sp=0.000 완전 회피, 2차 Flirds 0.917로 일부 버팀. FedSV Sp가 poison서 0.367로 추락.")
_page_number(s)

# ---------------------------------------------------- Slide 15: device100 noisy α
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "탐지 (2차③) — 강건성 cross-device N=100: noisy 탐지 AUROC (α-sweep)",
    ["대규모 단말 연합 cross-device, N=100, 라운드당 10명 · Llama-3.2-1B, R=30 · Dirichlet α∈{0,0.01,0.1,0.5=anchor,5.0} · 오염 5명 · 3-seed",
     "값=noisy 탐지 AUROC(↑). α=0.5 anchor에서만 (b) per-round 오라클 + coalition baseline(GTG/FedSV/ShapleyFL) 켬. 열별 최고 볼드."])
cols = [("method", None), ("α=0.0", 'up'), ("α=0.01", 'up'), ("α=0.1", 'up'),
        ("α=0.5 (anchor)", 'up'), ("α=5.0", 'up')]
rows = [
 (["(b)oracle","–","–","–","0.604±.041","–"], False),
 (["Flirds","0.774±.058","0.575±.055","0.605±.056","0.604±.041","0.596±.039"], True),
 (["Flirds-1st","0.772±.058","0.575±.057","0.606±.055","0.605±.042","0.597±.038"], True),
 (["loss-heur","0.772±.058","0.574±.056","0.607±.056","0.605±.042","0.597±.038"], True),
 (["FedIF","0.973±.017","0.568±.106","0.693±.126","0.830±.085","0.973±.022"], True),
 (["GTG","–","–","–","0.734±.112","–"], True),
 (["FedSV","–","–","–","0.708±.142","–"], True),
 (["ShapleyFL","–","–","–","0.762±.095","–"], True),
 (["ComFedSV","0.442±.115","0.419±.054","0.432±.032","0.371±.028","0.396±.002"], True),
 (["FLDetector","0.535±.048","0.482±.085","0.525±.070","0.539±.055","0.532±.058"], True),
 (["STD-DAGMM","0.856±.037","0.652±.190","0.659±.147","0.671±.142","0.760±.040"], True),
 (["FLTrust","1.000","0.602±.096","0.720±.136","0.854±.090","0.994±.008"], True),
 (["FedDQC","0.960±.057","1.000","1.000","1.000","1.000"], True),
]
add_table(s, top, cols, rows, col_widths([2.0, 2.05, 2.05, 2.05, 2.25, 2.05]), 9.5)
slide_footer(s, "개요 §3.4.2 (b1) (runs/phase2_matrix · master_metrics.csv)",
             "데이터-품질 전용 FedDQC가 1.0으로 최강. valuation φ는 0.57~0.77로 침식(비-IID서 정상 소수 클라가 상위). 결정적: exact (b) 오라클 자체도 0.604 — 침식은 근사 결함이 아닌 valuation 본질.")
_page_number(s)

# ---------------------------------------------------- Slide 16: device100 free-rider
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "탐지 (2차③) — 강건성 cross-device N=100: free-rider 탐지 AUROC",
    ["cross-device N=100, Llama-3.2-1B · free-rider 두 모드(random=benign-std / zero) · 대표 α(0.0/0.5) · 3-seed",
     "값=AUROC(↑). 열별 최고 볼드."])
cols = [("method", None), ("frrand α=0.0", 'up'), ("frrand α=0.5", 'up'),
        ("frzero α=0.0", 'up'), ("frzero α=0.5", 'up')]
rows = [
 (["Flirds / Flirds-1st / loss-heur","1.000","1.000","1.000","1.000"], True),
 (["FedIF","0.983±.009","0.981±.003","0.989±.004","0.987"], True),
 (["ComFedSV","0.449±.130","0.383±.045","0.441±.122","0.367±.028"], True),
 (["FLDetector","0.606±.030","0.617±.039","0.529±.024","0.540±.061"], True),
 (["STD-DAGMM","0.960±.029","0.588±.205","0.870±.094","0.963±.036"], True),
 (["FLTrust","1.000","1.000","1.000","1.000"], True),
 (["FedDQC","0.140±.014","0.573±.113","0.140±.014","0.573±.113"], True),
]
add_table(s, top, cols, rows, col_widths([3.5, 2.25, 2.25, 2.25, 2.2]), 11)
slide_footer(s, "개요 §3.4.2 (b2) (runs/phase2_matrix · master_metrics.csv)",
             "gradient 쓰는 방법(Flirds/Flirds-1st/loss-heur/FLTrust)이 1.0으로 깔끔. model-free STD-DAGMM은 가변, FedDQC는 off-threat(free-rider는 데이터-품질 아님)이라 낮음.")
_page_number(s)

# ---------------------------------------------------- Slide 17: device100 Spearman
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "Fidelity — 강건성 cross-device N=100: Spearman vs truth",
    ["cross-device N=100, Llama-3.2-1B · 3-seed. truth가 셀마다 다름:",
     "좌 3열(α=0.5 anchor) truth=(b) per-round 정확 오라클(genuine). 우측 열 truth=Flirds 기준 대용(proxy) → '1.000'은 'Flirds와 동일 순위'의 뜻(vs 오라클 아님)."])
cols = [("method", None), ("noisy α=0.5\n(truth=(b))", 'up'), ("frrand α=0.5\n(truth=(b))", 'up'),
        ("frzero α=0.5\n(truth=(b))", 'up'), ("off-anchor α\n(truth=Flirds proxy)", None)]
rows = [
 (["Flirds","1.000","1.000","1.000","(proxy 기준)"], True),
 (["Flirds-1st","1.000","1.000","1.000","모든 α 1.000"], True),
 (["loss-heur","1.000","1.000","1.000","모든 α 1.000"], True),
 (["FedIF","0.721±.027","0.827±.022","0.824±.017","α별 0.62~0.83"], True),
 (["GTG","0.784±.021","0.817±.022","0.843±.026","(anchor만)"], True),
 (["FedSV","0.752±.020","0.795±.020","0.814±.018","(anchor만)"], True),
 (["ShapleyFL","0.582±.075","0.685±.054","0.681±.049","(anchor만)"], True),
 (["ComFedSV","-0.023±.127","-0.051±.153","-0.051±.142","모든 α ≈0"], True),
]
add_table(s, top, cols, rows, col_widths([1.8, 2.6, 2.6, 2.6, 2.85]), 10)
slide_footer(s, "개요 §3.4.2 (b3) (runs/phase2_matrix · RESULTS.md)",
             "진짜 오라클 검증은 α=0.5 anchor 한 칸뿐(거기서 Flirds vs (b) per-round = 1.000, genuine). off-anchor 1.0은 Flirds 자기참조이므로 fidelity 증거로 쓰면 순환 — 주의.")
_page_number(s)

# ---------------------------------------------------- Slide 18: device100 poison
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "탐지 (2차③) — 강건성 cross-device N=100: poison",
    ["cross-device N=100, Llama-3.2-1B · poison(model-replacement backdoor, R=60, frac=0.8) · α=0.0(ASR≈1.00) / α=0.5(ASR≈0.50) · truth=Flirds proxy · 3-seed",
     "AUROC=탐지력(↑), Sp=Spearman vs Flirds proxy(↑). 열별 최고 볼드(Flirds는 proxy 기준이라 Sp 제외)."])
cols = [("method", None), ("α=0.0 AUROC↑", 'up'), ("α=0.0 Sp↑", 'up'),
        ("α=0.5 AUROC↑", 'up'), ("α=0.5 Sp↑", 'up')]
rows = [
 (["Flirds","1.000","(proxy truth)","1.000","(proxy truth)"], True),
 (["Flirds-1st","1.000","0.997±.002","0.670±.467","0.980±.028"], True),
 (["loss-heur","1.000","0.997±.002","1.000","0.999"], True),
 (["FedIF","0.542±.258","0.620±.204","0.458±.284","0.439±.071"], True),
 (["ComFedSV","0.778±.314","0.104±.054","0.727±.386","0.025±.098"], True),
 (["FLDetector","0.987±.019","–","0.983±.024","–"], True),
 (["STD-DAGMM","1.000","–","0.983±.024","–"], True),
 (["FLTrust","0.650±.180","–","0.498±.281","–"], True),
 (["FedDQC","1.000","–","1.000","–"], True),
]
add_table(s, top, cols, rows, col_widths([2.0, 2.65, 2.6, 2.65, 2.55]), 11)
slide_footer(s, "개요 §3.4.3 (master_metrics.csv 03_device100_poison)",
             "device100 poison은 cross-device 희석으로 설치가 약함(α0.5 ASR 0.50). 여기선 Flirds(2차) AUROC 1.0으로 회피 안 됨 = 설정 의존(cross-silo와 대비). caveat: tiny val=10.")
_page_number(s)

# ---------------------------------------------------- Slide 19: 3B silo5 1-seed
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "탐지 (2차③) — 강건성 cross-silo N=5 (3B, 1 seed ◐)",
    ["소수기관 연합 cross-silo, N=5 전원 · Llama-3.2-3B, R=10 · seeds=[0]만(3-seed 미완 ◐) · (b)=exact 2^5",
     "값=AUROC(↑)/Sp(↑)/runtime(↓), 1 seed라 ± 없음. 열별 최고 볼드((b)oracle 제외)."])
cols = [("method", None), ("noisy AUROC↑", 'up'), ("noisy Sp↑", 'up'),
        ("frrand AUROC↑", 'up'), ("frzero AUROC↑", 'up'),
        ("poison AUROC↑", 'up'), ("poison Sp↑", 'up'), ("runtime(noisy)↓", 'down')]
rows = [
 (["Flirds","1.000","1.000","1.000","1.000","0.000","0.000","~251s"], True),
 (["Flirds-1st","1.000","1.000","1.000","1.000","0.000","0.000","~82s"], True),
 (["loss-heur","1.000","1.000","1.000","1.000","1.000","1.000","~384s"], True),
 (["FedIF","1.000","0.600","1.000","1.000","1.000","0.600","~82s"], True),
 (["(b)oracle","1.000","(truth)","1.000","1.000","1.000","(truth)","~1244s"], False),
 (["FLDetector","1.000","–","1.000","1.000","1.000","–","~146–382s"], True),
 (["STD-DAGMM","0.250","–","1.000","0.000","0.750","–","~206–745s"], True),
 (["FLTrust","1.000","–","1.000","1.000","1.000","–","~83–91s"], True),
 (["FedDQC","1.000","–","0.750","0.750","1.000","–","~46–50s"], True),
]
add_table(s, top, cols, rows, col_widths([1.6,1.55,1.2,1.55,1.55,1.55,1.2,1.65]), 9)
slide_footer(s, "개요 §3.4.4 (master_metrics.csv 05_scale_3b)",
             "3B poison: Flirds·Flirds-1st 둘 다 AUROC/Sp 0.000 = clean-보존 backdoor에 완전 회피(1B는 2차가 0.917로 일부 버텼으나 3B는 둘 다 0). loss-heur·(b)·FedIF·탐지기는 1.0. 1 seed 한계.")
_page_number(s)

# ---------------------------------------------------- Slide 20: poison 회피 발견
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "탐지 — poison 회피 발견: clean-보존 backdoor가 Flirds-1st를 회피",
    ["cross-silo poison(clean-보존 backdoor) 칸만 모음 · 1B(3-seed) / 3B(1 seed ◐) · 값=poison AUROC(↑)·Spearman(↑)",
     "공격자가 clean val-loss를 낮춰 φ가 '기여 높음'으로 나옴 → 1차 Taylor가 부호를 놓침."])
cols = [("method", None), ("1B AUROC↑", None), ("1B Sp↑", None),
        ("3B AUROC↑", None), ("3B Sp↑", None)]
rows = [
 (["Flirds (1차+2차)","0.917±.118","0.967±.047","0.000","0.000"], True),
 (["Flirds-1st (1차만)","0.000","0.000","0.000","0.000"], True),
 (["loss-heur","1.000","1.000","1.000","1.000"], True),
 (["FedIF","1.000","0.967±.05","1.000","0.600"], True),
 (["(b) in-run 오라클","1.000","(truth)","1.000","(truth)"], False),
 (["FLDetector","1.000","–","1.000","–"], True),
 (["FLTrust","1.000","–","1.000","–"], True),
]
add_table(s, top, cols, rows, col_widths([3.2, 2.3, 2.3, 2.3, 2.35]), 11)
add_textbox(s, Inches(0.45), top + Inches(2.55), Inches(12.4), Inches(1.2),
    [("• Flirds-1st AUROC 0.000 = 완전 회피. 1B는 2차 Hessian 항이 0.917로 부호 일부 복원, 3B는 둘 다 0.000(1-seed).", 12, False, C_BODY),
     ("• 같은 val-loss 게임을 쓰는 (b) 정확 오라클·loss-heur는 AUROC 1.0으로 잡음 → '정직한 valuation 답'보다 1차 Taylor 부호 실패에 가까움.", 12, False, C_BODY),
     ("• cross-device에선 회피 안 됨(설치 약함) → 경계는 보편이 아니라 '큰 scaled-update × 큰 모델' 조건부. 상보적 탐지기 필요.", 12, False, C_BODY)],
    line_space=1.25)
slide_footer(s, "개요 §3.4.1 · §3.4.4")
_page_number(s)

# ---------------------------------------------------- Slide 21: LLM runtime
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "비용 (2차 이후) — LLM 표준 runtime (방법 × 스케일·스테이지)",
    ["LLM 표준 실험 · wall-clock 초(↓) · 3-seed mean±std · 표준=N=20 라운드당 2 / 정밀=N=5 전원",
     "(b)/(a) 오라클은 기준점(볼드 제외). 방법 중 최저를 열별 볼드."])
cols = [("method", None), ("1B 표준", 'down'), ("1B 정밀", 'down'), ("3B 표준", 'down'),
        ("3B 정밀", 'down'), ("7B 표준", 'down'), ("7B 정밀", 'down')]
rows = [
 (["Flirds-1st","1531±37","231±5","3630±76","547±10","6485±87","975±13"], True),
 (["Flirds","4697±112","707±16","11147±228","1674±30","20180±250","3027±36"], True),
 (["FedIF","1534±37","232±5","3638±75","549±10","6495±88","978±13"], True),
 (["loss-heur","2913±72","1093±26","6909±150","2585±48","12299±167","4613±62"], True),
 (["ComFedSV","2330±22","2557±215","5526±34","6043±525","9839±110","10792±1044"], True),
 (["GTG","3647±90","3552±82","8647±189","8393±153","15393±206","14972±193"], True),
 (["FedSV","3646±90","3536±86","8647±189","8356±161","15393±208","14907±270"], True),
 (["ShapleyFL","2917±72","3513±83","6916±150","8303±155","12312±167","14812±199"], True),
 (["Banzhaf","– (N=20)","3527±83","–","8329±155","–","14844±198"], True),
 (["(b)oracle","2917±72","3528±83","6916±151","8329±156","12310±165","14839±196"], False),
 (["(a)oracle","–","30817±244","–","⬚","–","⬚"], False),
]
add_table(s, top, cols, rows, col_widths([1.7, 1.79, 1.79, 1.79, 1.79, 1.79, 1.8]), 9.5)
slide_footer(s, "개요 §3.5.1 (runs/track_d/rundirs/*/metrics.json)",
             "Flirds-1st 항상 최저. 2차 Flirds 비용은 라운드당 cohort에 무관(1 HVP); (b) 오라클은 cohort에 지수적(2^k). (a) retrain = (b)의 ~9배. 다음 슬라이드에 비용 모델 정리.")
_page_number(s)

# ---------------------------------------------------- Slide 22: 비용 모델 요점
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "비용 — 비용 모델 요점: cohort 크기(2^k)가 우열을 가른다", [])
add_textbox(s, Inches(0.45), top + Inches(0.05), Inches(12.4), Inches(3.2),
    [("• Flirds-1st = 항상 최저가 (1 val-gradient/round, Hessian 없음).", 14, False, C_BODY),
     ("• 2차 Flirds 비용 = 라운드당 cohort 크기와 무관 (1 HVP/round 고정).", 14, True, C_TITLE),
     ("• (b) in-run 오라클 비용 = 라운드당 cohort에 지수적 (2^k coalition-eval/round).", 14, True, C_TITLE),
     ("", 6, False, C_BODY),
     ("함의 — 우열이 무대마다 갈린다:", 13.5, True, C_TITLE),
     ("   cohort 크면 Flirds 압승:  device100 anchor(K=10) 157s vs (b) 25,000s ≈ 160× · 정밀 스테이지(N=5 full) 707s vs 3528s ≈ 5×",
      13, False, C_BODY),
     ("   cohort 작으면 (b)가 더 쌈:  표준 스테이지는 라운드당 2명(2²=4 eval) → 1B 표준 (b) 2917s < Flirds(2차) 4697s. 이 레짐은 Flirds-1st만 우위.",
      13, False, C_BODY),
     ("   (a) retrain 오라클 = (b)의 ~9배 (1B 정밀 30,817s vs 3528s) — fidelity 비교군이 아닌 별도 정답.", 13, False, C_BODY),
     ("   CNN: Flirds·Flirds-1st가 FL 학습 자체(traj_time ~80~94s)보다 2~3 자릿수 싸다. Ripple은 압도적 dominated(학습보다 10~130×).",
      13, False, C_BODY)],
    line_space=1.28)
# 작은 Robustness runtime 요약 표
cols = [("무대", None), ("Flirds-1st", None), ("Flirds", None), ("(b)·coalition", None), ("탐지기", None)]
rows = [
 (["cross-silo N=5", "~35s", "~107s", "(b)·coalition ~530s", "22~136s"], False),
 (["cross-device N=100 anchor", "~53s", "~157s", "(b)perround ~25,000s / GTG ~16–18k / ShapleyFL ~24.9k / FedSV ~4970s", "–"], False),
]
add_table(s, top + Inches(3.55), cols, rows, col_widths([3.0, 1.5, 1.5, 5.0, 1.45]), 10.5)
slide_footer(s, "개요 §3.5.1 · §3.5.2")
_page_number(s)

# ---------------------------------------------------- Slide 23: CNN runtime MNIST
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "비용 — CNN cross-silo runtime: MNIST",
    ["CNN 실험 cross-silo N=10 · wall-clock 초(↓) · 3-seed mean±std · 시나리오별",
     "(b)oracle·traj_time(FL 학습 자체)은 기준점(볼드 제외). 방법 중 최저를 열별 볼드."])
cols = [("method", None), ("iid", 'down'), ("feature-noise", 'down'), ("label-flip", 'down'),
        ("label-skew", 'down'), ("quantity-skew", 'down')]
rows = [
 (["Flirds-1st","0.08±0.00","0.08±0.00","0.08±0.00","0.08±0.00","0.08±0.00"], True),
 (["Flirds","0.73±0.03","0.58±0.05","0.59±0.03","0.61±0.06","0.55±0.03"], True),
 (["FedIF","0.18±0.00","0.16±0.01","0.17±0.01","0.17±0.01","0.16±0.01"], True),
 (["loss-heur","0.24±0.00","0.24±0.00","0.24±0.00","0.24±0.00","0.24±0.00"], True),
 (["ComFedSV","5.13±0.18","5.17±0.04","5.13±0.18","5.12±0.14","5.10±0.15"], True),
 (["GTG","18.11±0.11","19.77±0.48","16.29±1.51","25.01±0.90","17.26±0.95"], True),
 (["FedSV","6.23±0.07","6.29±0.11","6.29±0.14","6.22±0.11","6.21±0.04"], True),
 (["ShapleyFL","30.65±0.26","30.80±0.43","30.58±0.64","30.47±0.56","30.31±0.31"], True),
 (["Banzhaf","31.16±0.29","31.34±0.49","31.21±0.71","30.85±0.53","31.04±0.45"], True),
 (["(b)oracle","31.09±0.36","31.33±0.47","31.04±0.62","30.89±0.55","30.91±0.48"], False),
 (["Ripple","2254.83±758.57","1147.51±70.79","1160.83±36.82","1121.67±22.25","1075.64±52.27"], True),
 (["traj_time (FL 학습)","93.88±1.32","92.85±1.02","91.81±1.16","84.57±1.56","85.36±1.38"], False),
]
add_table(s, top, cols, rows, col_widths([2.4, 2.0, 2.05, 2.0, 2.0, 2.0]), 9)
slide_footer(s, "개요 §3.5.2 (runs/track_c/c1/*/metrics.json)",
             "Flirds-1st 최저(0.08s). Ripple은 학습 자체보다 10~130× 느려 압도적 dominated.")
_page_number(s)

# ---------------------------------------------------- Slide 24: CNN runtime CIFAR
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "비용 — CNN cross-silo runtime: CIFAR-10",
    ["CNN 실험 cross-silo N=10 · wall-clock 초(↓) · 3-seed mean±std · 시나리오별",
     "(b)oracle·traj_time은 기준점(볼드 제외). 방법 중 최저를 열별 볼드."])
cols = [("method", None), ("iid", 'down'), ("feature-noise", 'down'), ("label-flip", 'down'),
        ("label-skew", 'down'), ("quantity-skew", 'down')]
rows = [
 (["Flirds-1st","0.34±0.01","0.37±0.04","0.38±0.05","0.34±0.01","0.34±0.01"], True),
 (["Flirds","1.20±0.08","1.62±0.49","1.44±0.34","1.27±0.17","5.78±6.24"], True),
 (["FedIF","0.44±0.01","0.48±0.02","0.50±0.07","0.54±0.15","0.49±0.06"], True),
 (["loss-heur","1.27±0.01","1.30±0.03","1.27±0.00","1.28±0.01","1.30±0.02"], True),
 (["ComFedSV","18.45±0.36","19.18±1.10","19.07±0.82","18.83±0.19","19.04±0.51"], True),
 (["GTG","88.98±2.02","87.84±5.34","88.25±3.10","102.46±0.58","75.50±0.96"], True),
 (["FedSV","22.29±0.27","23.10±0.73","23.19±0.91","23.87±1.76","22.81±0.25"], True),
 (["ShapleyFL","110.83±1.28","115.58±4.74","112.25±1.82","113.03±1.14","115.42±3.39"], True),
 (["Banzhaf","110.93±1.26","115.84±5.30","115.94±6.90","112.69±1.10","112.70±0.45"], True),
 (["(b)oracle","110.48±0.93","119.41±4.78","114.54±3.85","112.87±1.16","168.69±71.77"], False),
 (["Ripple","7453.15±710.09","9592.28±2713.45","11105.31±2385.57","10702.14±898.62","10072.95±2706.30"], True),
 (["traj_time (FL 학습)","80.21±0.45","87.28±5.38","86.70±7.24","83.12±2.05","91.23±3.36"], False),
]
add_table(s, top, cols, rows, col_widths([2.4, 2.0, 2.05, 2.05, 2.0, 2.05]), 8.5)
slide_footer(s, "개요 §3.5.2 (runs/track_c/c1/*/metrics.json)",
             "순서는 MNIST와 동일(Flirds-1st 최저 < Flirds ≪ exact 2^10급). Ripple ~7.5~11k s = 압도적 dominated.")
_page_number(s)

# ---------------------------------------------------- Slide 25: 한계·주의
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "한계 · 주의 (caveats)", [])
add_textbox(s, Inches(0.45), top + Inches(0.05), Inches(12.4), Inches(5.0),
    [("• 3B 강건성 = 1 seed (◐).  3B 강건성 전 수치가 단일 seed → 통계 없음(특히 poison 둘 다 0.000).", 13.5, False, C_BODY),
     ("• retrain 오라클 (a) = 1B 정밀 스테이지만.  3B/7B는 fidelity·runtime만 있고 (a) 미실행(⬚).", 13.5, False, C_BODY),
     ("• cross-device 비-anchor truth = Flirds 기준 대용(proxy).  정확 (b)가 칸당 ~25,000s라 α=0.5만 실측 → 그 칸의 Spearman은 vs Flirds(자기참조).", 13.5, False, C_BODY),
     ("• CNN fidelity pool 평균은 iid 포함 → 깎임.  iid 셀은 오염·skew 신호가 없어 fidelity가 의미상 낮다(별도 iid-제외 값 병기).", 13.5, False, C_BODY),
     ("• CNN cross-device 그룹 표는 partition·강도·dataset을 위협 내에서 pool → std 큼.  셀별 30칸은 RESULTS.txt.", 13.5, False, C_BODY),
     ("• tiny val.  cross-silo val=20 / cross-device val=10 → AUROC가 coarse(특히 noisy φ-as-detector).", 13.5, False, C_BODY),
     ("• poison ASR은 deployed-model 기준(cross-silo≈1.00, cross-device α0≈1.00/α0.5≈0.50, 3B≈1.00).", 13.5, False, C_BODY),
     ("• 7B 정밀 스테이지 arm(MMLU/ROUGE/val-loss)은 2026-06-26 추가 완료 → LLM 표준 6 셀 전부 arm 포함.", 13.5, False, C_BODY)],
    line_space=1.4)
slide_footer(s, "개요 §4.2")
_page_number(s)

# ---------------------------------------------------- Slide 26: 부록 phase1 AUROC
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "부록 — 기초 검증: 첫 clean run (AUROC + selection)",
    ["기초 검증(첫 clean run + 학습률 sweep) · Llama-3.2-1B, N=5 전원, K=3 · noisy=client0 answer-swap / free-rider=client1 zero-update",
     "full: R=50, lr∈{1e-3,3e-3}×3 seed · sweep: R=20, lr 4종×1 seed. flirds_keep/random_keep은 설정·산출값(볼드 없음)."])
cols = [("group", None), ("noisy AUROC ↑", None), ("free-rider AUROC ↑", None),
        ("flirds_keep (seed별)", None), ("random_keep", None)]
rows = [
 (["full lr1e-3 (3 seed)","0.750±.000","1.000±.000","[3,2,4] 매 seed (=clean)","seed별 가변"], False),
 (["full lr3e-3 (3 seed)","1.000±.000","0.750±.000","[2,3,4] 매 seed (=clean)","seed별 가변"], False),
 (["sweep lr1e-4 (1 seed)","0.750","1.000","[3,2,4]","[2,3,4]"], False),
 (["sweep lr3e-4 (1 seed)","0.750","1.000","[3,2,4]","[2,3,4]"], False),
 (["sweep lr1e-3 (1 seed)","0.750","1.000","[3,2,4]","[2,3,4]"], False),
 (["sweep lr3e-3 (1 seed)","0.750","1.000","[3,2,4]","[2,3,4]"], False),
]
add_table(s, top, cols, rows, col_widths([3.0, 2.2, 2.6, 2.85, 1.8]), 11)
slide_footer(s, "개요 §4.1 (runs/phase1/rundirs/*/metrics.json)",
             "lr 의존 반전: full lr1e-3은 noisy 0.75/FR 1.0, lr3e-3은 noisy 1.0/FR 0.75. selection: flirds_keep이 매 seed 정확히 clean 클라 3개(noisy·free-rider 항상 드롭).")
_page_number(s)

# ---------------------------------------------------- Slide 27: 부록 phase1 arms
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "부록 — 기초 검증: selection arms (final val-loss)",
    ["기초 검증 full · final val-loss(↓) · 3-seed mean±std · 행별 최저 arm 볼드(full=전원, flirds_topk=Flirds 상위 선택, random_k=무작위 k)"])
cols = [("group", None), ("full (전원) ↓", 'down'), ("flirds_topk ↓", 'down'), ("random_k ↓", 'down')]
rows = [
 (["full lr1e-3","2.4064±.0234","2.3978±.0226","2.4111±.0133"], True),
 (["full lr3e-3","2.3931±.0223","2.3926±.0219","2.4055±.0100"], True),
]
add_table(s, top, cols, rows, col_widths([3.0, 3.15, 3.15, 3.15]), 12, bold_axis='row')
slide_footer(s, "개요 §4.1 (runs/phase1/rundirs/*/metrics.json)",
             "flirds_topk val-loss ≤ random_k (양 lr) 그리고 ≤ full(오염 드롭이 도움) → 'random은 hard bar'를 넘김.")
_page_number(s)

# ---------------------------------------------------- Slide 28: 결론 한 문단
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = C_BG_DARK; bg.line.fill.background()
bg.shadow.inherit = False
s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
add_textbox(s, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.8),
            [("결론 — 결과 한 문단 요약", 26, True, C_WHITE)])
add_textbox(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.6),
    [("Flirds는 연합학습 한 학습 궤적에서 client-level 데이터 기여도를 1차+2차 Taylor(true Hessian)로 추정한다. "
      "네 트랙에서, clean·near-additive 레짐의 fidelity는 사실상 포화된다 — LLM에서 Flirds·Flirds-1st·loss-heur·"
      "GTG·Banzhaf가 거의 모두 exact in-run 오라클 대비 Spearman 1.000으로 동률이라, 이 레짐에서 Flirds의 차별점은 "
      "정확도가 아니라 비용이다.", 15, False, C_ICE),
     ("2차 Flirds의 비용은 라운드당 cohort 크기와 무관(1 HVP/round)인 반면 exact 오라클은 cohort에 지수적(2^k/round)이라, "
      "참여가 많은 무대에서 Flirds가 오라클을 160×까지(cross-device anchor 157s vs 25,000s) 앞선다. 정확도 차별은 "
      "near-additivity가 깨지는 곳에서만 드러난다: (i) poison(clean-보존 backdoor)에서 FedSV(0.367)·Flirds-1st(0.000)가 "
      "무너질 때 2차 Flirds가 0.967로 버티고, (ii) cross-device non-IID에서 Flirds 1.000 vs GTG 0.78/FedSV 0.75.", 15, False, C_ICE),
     ("단 CNN은 거동이 다르다 — clean에서도 fidelity가 큰 spread를 보이고(Flirds 0.919 vs GTG 0.57), 2차항이 benign에서도 "
      "돕고, retrain 오라클 (a)와 in-run 오라클 (b)가 갈린다(Flirds vs (b)=0.919 / vs (a)=0.35) — LLM(1B)에서 둘이 0.933으로 "
      "거의 일치하는 것과 대조된다. 2차 검증은 위계대로다: clean-IID 성능·수렴은 do-no-harm parity, 오염이 있어야 회복하며, "
      "탐지는 마지막 — valuation φ는 전용 탐지기가 아니다(non-IID φ 0.57~0.77 vs FedDQC 1.0; exact (b)도 0.660).", 15, False, C_ICE),
     ("종합하면 가장 단단한 주장은 'Flirds는 exact in-run Shapley를 충실·저렴하게 근사하며 그 비용 우위는 cohort에 무관하다'이고, "
      "정확도 우위·retrain 충실도·poison 강건성은 무대 의존적 조건부 주장이다.", 15, True, C_WHITE)],
    line_space=1.22)
slide_footer(s, "flirds-results-analysis-2026-06-26 §0")
_page_number(s)

# ---------------------------------------------------- Slide 29: 위계별 핵심 발견
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "결론 — 위계별 핵심 발견", [])
add_textbox(s, Inches(0.45), top + Inches(0.0), Inches(12.5), Inches(5.4),
    [("1차 Fidelity:", 13.5, True, C_TITLE),
     ("   clean·near-additive에서 fidelity 포화 — LLM 다수 방법 Spearman 1.000, CNN은 spread(Flirds 0.919 / GTG 0.57). "
      "차별점은 비용. 정확도 차별은 비가법(poison/non-IID)에서만.", 12.5, False, C_BODY),
     ("2차① 성능 / 집계:", 13.5, True, C_TITLE),
     ("   clean-IID는 do-no-harm parity(±0.001~0.003). 오염이 있으면 기여도-가중이 정확도 회복(CNN grad_noise 0.499→0.609). "
      "단 Flirds가 arm 중 최강은 아님(shapleyfl 0.645).", 12.5, False, C_BODY),
     ("2차② 수렴:", 13.5, True, C_TITLE),
     ("   clean-IID 거의 동률. 7B 표준만 ~14~18% 빠름(분산 큼, under-powered). CNN은 오염서 크게 앞당김(27.2→6.3 라운드).", 12.5, False, C_BODY),
     ("2차③ 탐지 (위계상 마지막):", 13.5, True, C_TITLE),
     ("   noisy·free-rider는 거의 AUROC 1.0. valuation ≠ 탐지기 — non-IID φ 0.57~0.77 < FedDQC 1.0, exact (b)도 0.660(설계상 예측). "
      "poison clean-보존 backdoor가 Flirds-1st 회피(0.000), 2차가 일부 복원.", 12.5, False, C_BODY),
     ("비용:", 13.5, True, C_TITLE),
     ("   Flirds-1st 항상 최저. 2차 Flirds 비용 cohort-독립 → 참여 많으면 오라클 160× 압승, 작으면 (b)가 쌈(Flirds-1st만 우위).", 12.5, False, C_BODY),
     ("CNN ↔ LLM:", 13.5, True, C_TITLE),
     ("   near-additivity 강도 차이가 천장/spread·2차항 역할·(a)/(b) 오라클 괴리를 한 번에 설명(가설, 단정 X).", 12.5, False, C_BODY)],
    line_space=1.18)
slide_footer(s, "flirds-results-analysis-2026-06-26 §1 · §2")
_page_number(s)

# ---------------------------------------------------- Slide 30: novelty + 열린 질문
s = prs.slides.add_slide(BLANK)
top = slide_header(s, "결론 — novelty · 기여 + 남은 열린 질문", [])
add_textbox(s, Inches(0.45), top + Inches(0.0), Inches(12.5), Inches(2.7),
    [("기여 (verification·taxonomy와 엮은 새로운 것):", 14, True, C_TITLE),
     ("1. LLM-scale client-level in-run Shapley의 fidelity 검증 — 선행 공백(관찰된 빈칸). 인접 칸 점유자는 valuation이 아닌 "
      "다른 문제(품질·selection·market). 1B→7B 3 스케일로 채움.", 12.5, False, C_BODY),
     ("2. cohort-독립 비용 모델 — Shapley-충실하면서 라운드당 참여 수와 무관한 저비용(cross-device 160× 정량화).", 12.5, False, C_BODY),
     ("3. poison-회피 경계 + 2차항의 FL-특이적 효용 — IRDS '중앙 per-step 2차 이득 미미'와 무대가 달라 모순 아님(직접 증거).", 12.5, False, C_BODY),
     ("4. (부수) dual-oracle 괴리 발견 — in-run (b)과 retrain (a) Shapley가 CNN full-model에서 다른 게임(vs (a)≤0.45).", 12.5, False, C_BODY)],
    line_space=1.2)
add_textbox(s, Inches(0.45), top + Inches(3.0), Inches(12.5), Inches(2.6),
    [("남은 열린 질문 (단정하지 않음):", 14, True, C_TITLE),
     ("• poison-회피는 1차 Taylor 한계인가 매트릭스 경계인가 — 실제-config + (b) 오라클 재확인, 3B 3-seed 필요.", 12.5, False, C_BODY),
     ("• (a)≠(b) 괴리가 CNN-고유인가 LLM서도 스케일로 나타나나 — 3B/7B (a) 미실행이라 확인 불가.", 12.5, False, C_BODY),
     ("• near-additivity가 clean fidelity 1.000을 trivial하게 만드나 — 헤드라인은 clean-1.0이 아니라 비가법 차별 + 비용에.", 12.5, False, C_BODY),
     ("• proxy-truth 순환성 — off-anchor 1.0은 Flirds 자기참조; 진짜 오라클 검증은 anchor 한 칸뿐.", 12.5, False, C_BODY)],
    line_space=1.2)
slide_footer(s, "flirds-results-analysis-2026-06-26 §4 · §5")
_page_number(s)

# ---------------------------------------------------------------- save
import os
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "flirds-results-2026-06-26.pptx")
prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))
