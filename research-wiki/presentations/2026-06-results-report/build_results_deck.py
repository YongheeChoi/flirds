#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Flirds 실험 결과 보고 데크 (독립 상세판) — 표 중심, 담백한 학술 핸드아웃 톤.

빌드:  python3 build_results_deck.py
출력:  flirds-results-2026-06.pptx  (16:9, 편집 가능 네이티브 표)

구성: 모델(CNN→LLM)이 아니라 **검증 질문별**로 묶는다 — 1차 Fidelity(Shapley 추정 정확도)
      → 2차-① Performance → ② Convergence → ③ Detection. (핵심 질문 위계)

데이터 출처:
  - Track C C1 (b)fidelity / 거리 : runs/track_c/c1/<cell>/metrics.json  직접 재집계 (10종 비교 + 오라클)
  - Track C C1 (a)retrain / C2    : runs/track_c/RESULTS.txt 파싱 (3-seed 집계분)
  - Track D 1B/3B                 : runs/track_d/rundirs/<cell>/metrics.json 직접 재집계
  - device100 α=0.5 anchor        : runs/phase2_matrix/rundirs/1B_device100-a0.5_*_anchor/ 재집계 (1 seed)
  - Phase 2 silo5 / device100 / 3B: runs/phase2_matrix/RESULTS.md 의 검증 수치 (상수)
톤: presentation-style-plain-facts — 흰 배경, 강조색 1개, 표/수식 중심, 과장수식어 금지,
    모든 수치에 설정 병기. 표 수치 부호 + 생략(음수만 −). IID/Non-IID 분할 성격 병기.
    어색한 한글 기술용어(intervention/under threat/weighting arm 등)는 영어로 둠.
"""
import glob, json, re, statistics
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

ROOT = Path("/NHNHOME/26msit001_A/edge_ai_lab/yonghee/flirds")
HERE = Path(__file__).resolve().parent
OUT  = HERE / "flirds-results-2026-06.pptx"

# ── 색/폰트 (plain-facts) ───────────────────────────────────────────
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SUB   = RGBColor(0x55, 0x55, 0x55)
FAINT = RGBColor(0x88, 0x88, 0x88)
ACCENT= RGBColor(0x17, 0x50, 0x8C)   # 강조색 1개
LINE  = RGBColor(0xC8, 0xC8, 0xC8)
HDRBG = RGBColor(0xEF, 0xEF, 0xEF)
ZEBRA = RGBColor(0xF8, 0xF8, 0xF8)
EMPH  = RGBColor(0xEA, 0xF0, 0xF7)   # 그룹 평균 행 강조 (연한 강조색)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Malgun Gothic"
NO_STYLE_GRID = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # No Style, Table Grid

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PW = 13.333

# ── helpers ─────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    return s

def _box(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def _set(par, text, size, color=INK, bold=False, font=FONT, align=PP_ALIGN.LEFT):
    par.alignment = align
    r = par.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return r

def header(s, kicker, title, sub=None):
    if kicker:
        tb = _box(s, 0.55, 0.34, 12.2, 0.34); tf = tb.text_frame; tf.word_wrap = True
        tf.margin_top = 0; tf.margin_bottom = 0
        _set(tf.paragraphs[0], kicker.upper(), 12, ACCENT, bold=True)
    tb = _box(s, 0.55, 0.62, 12.2, 0.7); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_top = 0; tf.margin_bottom = 0
    _set(tf.paragraphs[0], title, 24, INK, bold=True)
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.30), Inches(12.23), Pt(1.6))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    y = 1.44
    if sub:
        tb = _box(s, 0.55, 1.40, 12.2, 0.5); tf = tb.text_frame; tf.word_wrap = True
        tf.margin_top = 0
        _set(tf.paragraphs[0], sub, 13.5, SUB)
        y = 1.86
    return y

def caption(s, left, top, w, text, color=ACCENT, size=11.5):
    tb = _box(s, left, top, w, 0.3); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_top = 0; tf.margin_bottom = 0; tf.margin_left = 0
    _set(tf.paragraphs[0], text, size, color, bold=True)

def footnote(s, text):
    tb = _box(s, 0.55, 7.08, 12.2, 0.34); tf = tb.text_frame; tf.word_wrap = True
    _set(tf.paragraphs[0], text, 10, FAINT)

def notes(s, l, t, w, lines, size=12.5, gap=6):
    tb = _box(s, l, t, w, 0.4); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_top = 0; tf.margin_left = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        if isinstance(ln, tuple):
            txt, bold = ln[0], ln[1]; col = ln[2] if len(ln) > 2 else INK
            _set(p, txt, size, col, bold=bold)
        else:
            _set(p, ln, size, INK)
    return tb

def sb(s, left, top, w, spec, metrics, read=None, size=9.4, gap=2.5):
    """결과표 옆 사이드바: 명세 / 지표 / 읽기 블록."""
    lines = [("명세", True, ACCENT)] + list(spec) + ["", ("지표", True, ACCENT)] + list(metrics)
    if read:
        lines += ["", ("읽기", True, ACCENT)] + list(read)
    return notes(s, left, top, w, lines, size=size, gap=gap)

def table(s, left, top, total_w, rows, ratios, aligns=None, size=10.5,
          header_row=True, row_h=0.0, zebra=False, accent_first=False, emph_rows=None):
    nr, nc = len(rows), len(rows[0])
    aligns = aligns or (["l"] + ["c"] * (nc - 1))
    emph_rows = emph_rows or set()
    h = Inches(row_h * nr) if row_h else Inches(0.3 * nr)
    gtbl = s.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(total_w), h)
    tb = gtbl.table
    tblPr = tb._tbl.find(qn('a:tblPr'))
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is None:
        sid = tblPr.makeelement(qn('a:tableStyleId'), {}); tblPr.append(sid)
    sid.text = NO_STYLE_GRID
    tot = sum(ratios)
    for j, rt in enumerate(ratios):
        tb.columns[j].width = Emu(int(Inches(total_w) * rt / tot))
    amap = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
    for i in range(nr):
        if row_h: tb.rows[i].height = Inches(row_h)
        for j in range(nc):
            cell = tb.cell(i, j)
            cell.margin_left = Pt(4); cell.margin_right = Pt(4)
            cell.margin_top = Pt(1.5); cell.margin_bottom = Pt(1.5)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_hdr = header_row and i == 0
            is_emph = i in emph_rows
            if is_hdr:
                cell.fill.solid(); cell.fill.fore_color.rgb = HDRBG
            elif is_emph:
                cell.fill.solid(); cell.fill.fore_color.rgb = EMPH
            elif zebra and i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ZEBRA
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            txt = str(rows[i][j])
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            bold = is_hdr or is_emph or (accent_first and j == 0 and i > 0)
            col = ACCENT if is_hdr else INK
            _set(p, txt, size, col, bold=bold, align=amap[aligns[j]])
    return tb

def fmt(v, nd=3):
    """표 수치 포맷 — 부호 '+'는 표기하지 않음(음수만 −). None→'–', 문자열은 그대로."""
    if v is None: return "–"
    if isinstance(v, str): return v
    return f"{v:.{nd}f}"

def _avg(vals):
    xs = [v for v in vals if isinstance(v, (int, float))]
    return sum(xs) / len(xs) if xs else None

# ════════════════════════════════════════════════════════════════════
# 데이터 — Track C C1 (b)/거리 + Track D : rundir 직접 재집계
# ════════════════════════════════════════════════════════════════════
SCEN_ORDER = ["iid", "label_skew", "quantity_skew", "label_flip", "feature_noise"]
def _scen_key(ds, sc): return f"{ds}_{sc}"

def load_c1():
    sp = defaultdict(lambda: defaultdict(list)); eu = defaultdict(lambda: defaultdict(list))
    for c in glob.glob(str(ROOT / "runs/track_c/c1/*/metrics.json")):
        m = json.load(open(c)); k = _scen_key(m["dataset"], m["scenario"])
        for meth, e in m["methods"].items():
            if isinstance(e, dict):
                if "spearman_b" in e: sp[k][meth].append(e["spearman_b"])
                if "euc_b" in e: eu[k][meth].append(e["euc_b"])
    mean = lambda d: {k: {me: sum(v)/len(v) for me, v in mv.items()} for k, mv in d.items()}
    return mean(sp), mean(eu)
C1_SP, C1_EU = load_c1()

def load_track_d(model, regime):
    fid = defaultdict(lambda: defaultdict(list)); arms = defaultdict(lambda: defaultdict(list)); n = 0
    for r in sorted(glob.glob(str(ROOT / f"runs/track_d/rundirs/{model}_{regime}_seed*/metrics.json"))):
        d = json.load(open(r)); m = d[list(d.keys())[0]]; n += 1
        for mt in ("spearman", "kendall", "cosine_d", "euclid_d", "max_diff", "runtime"):
            for me, v in m.get(mt, {}).items(): fid[mt][me].append(v)
        for arm, av in m.get("arms", {}).items():
            for f, v in av.items(): arms[arm][f].append(v)
    fm = {mt: {me: sum(v)/len(v) for me, v in mv.items()} for mt, mv in fid.items()}
    am = {}
    for arm, av in arms.items():
        am[arm] = {f: (_avg(vs)) for f, vs in av.items()}
    return fm, am, n
D_STD_F, D_STD_A, D_STD_N = load_track_d("1B", "std20")
D_ANC_F, D_ANC_A, D_ANC_N = load_track_d("1B", "anchor5")
D3_STD_F, _, D3_STD_N = load_track_d("3B", "std20")
D3_ANC_F, D3_ANC_A, D3_ANC_N = load_track_d("3B", "anchor5")

# ── Track C C1 (a)retrain + C2 : RESULTS.txt 파싱 (3-seed 집계분) ──────
def _c2key(cell):
    for ds in ("cifar10", "fmnist"):
        if cell.startswith(ds + "_"):
            rest = cell[len(ds) + 1:]
            for st in ("strmain", "str0.05", "str0.6", "str0.8"):
                if rest.endswith("_" + st):
                    toks = rest[:-len(st) - 1].split("_")
                    return (ds, toks[0], "_".join(toks[1:]), st)
    return (cell,)

def parse_track_c_results():
    txt = (ROOT / "runs/track_c/RESULTS.txt").read_text()
    c1a = {}
    for cell, body in re.findall(r"=== \(a\)-oracle fidelity (\S+) \(3 seeds\) ===\n(.*?)(?=\n===|\n====|\Z)", txt, re.S):
        d = {}
        for ln in body.strip().splitlines():
            mt = re.match(r"\s*(\S+)\s+(-?\d+\.\d+)\s+(-?\d+\.\d+)", ln)
            if mt: d[mt.group(1)] = float(mt.group(2))
        c1a[cell] = d
    c2 = {}
    for cell, body in re.findall(r"=== C2 outcome stability (\S+) \(3 seed\(s\)\) ===\n(.*?)(?=\n===|\n====|\Z)", txt, re.S):
        arms = {}
        for ln in body.strip().splitlines():
            mt = re.match(r"\s*(\S+)\s+(\d+\.\d+)\s+(\d+\.\d+)\s+(\S+)\s+(\S+)", ln)
            if mt:
                au = mt.group(4)
                arms[mt.group(1)] = (float(mt.group(2)), None if au == "--" else float(au))
        c2[_c2key(cell)] = arms
    return c1a, c2
C1_A, C2F = parse_track_c_results()

# ── device100 α=0.5 anchor (1 seed) : rundir 재집계 ───────────────────
def load_dev_anchor():
    out = {}
    for threat, suf in (("noisy", "noisy_anchor"), ("frrand", "frrand_anchor"),
                        ("frzero", "frzero_anchor"), ("poison", "poison")):
        au = defaultdict(list); sp = defaultdict(list)
        for f in glob.glob(str(ROOT / f"runs/phase2_matrix/rundirs/1B_device100-a0.5_{suf}*/metrics.json")):
            m = json.load(open(f))
            for v in m.values():
                for me, x in v.get("auroc", {}).items(): au[me].append(x)
                for me, x in v.get("spearman", {}).items(): sp[me].append(x)
        out[threat] = ({me: _avg(v) for me, v in au.items()}, {me: _avg(v) for me, v in sp.items()})
    return out
DEV_ANCHOR = load_dev_anchor()

# ════════════════════════════════════════════════════════════════════
# 데이터 — 상수 (Phase 2 RESULTS.md 검증분)
# ════════════════════════════════════════════════════════════════════
SILO_M = ["in-run oracle","Flirds","Flirds1st","FedIF","GTG","FedSV","ShapleyFL","Banzhaf","loss-heur","FLDetector","STD-DAGMM","FLTrust","FedDQC"]
SILO_AUROC = {
 "in-run oracle":(1.000,1.000,1.000,1.000),"Flirds":(1.000,1.000,1.000,0.917),"Flirds1st":(1.000,1.000,1.000,0.000),
 "FedIF":(1.000,1.000,1.000,1.000),"GTG":(1.000,1.000,1.000,1.000),"FedSV":(1.000,1.000,1.000,1.000),
 "ShapleyFL":(1.000,1.000,1.000,1.000),"Banzhaf":(1.000,1.000,1.000,1.000),"loss-heur":(1.000,1.000,1.000,1.000),
 "FLDetector":(0.750,1.000,0.750,1.000),"STD-DAGMM":(0.417,1.000,0.250,0.750),"FLTrust":(1.000,1.000,1.000,1.000),
 "FedDQC":(0.917,0.750,0.750,1.000),
}
SILO_SP = {
 "Flirds":(1.000,1.000,1.000,0.967),"Flirds1st":(1.000,1.000,1.000,0.000),"FedIF":(0.933,0.900,0.933,0.967),
 "GTG":(1.000,1.000,1.000,0.867),"FedSV":(1.000,0.933,1.000,0.367),"ShapleyFL":(1.000,1.000,1.000,1.000),
 "Banzhaf":(1.000,1.000,1.000,1.000),"loss-heur":(1.000,1.000,1.000,1.000),
}
SILO_RT = {"in-run oracle":532,"Flirds":106,"Flirds1st":35,"FedIF":35,"GTG":538,"FedSV":533,"ShapleyFL":530,
 "Banzhaf":533,"loss-heur":164,"FLDetector":30,"STD-DAGMM":136,"FLTrust":36,"FedDQC":22}

ALPHAS = ["0.0","0.01","0.1","0.5","5.0"]
DEV_NOISY = {
 "Flirds":(0.774,0.575,0.605,0.604,0.596),"FedIF":(0.973,0.568,0.693,0.830,0.973),
 "STD-DAGMM":(0.856,0.652,0.659,0.671,0.760),"FLTrust":(1.000,0.602,0.720,0.854,0.994),
 "FedDQC":(0.960,1.000,1.000,1.000,1.000),"FLDetector":(0.535,0.482,0.525,0.539,0.532),
 "ComFedSV":(0.442,0.419,0.432,0.371,0.396),
}
DEV_POISON = {
 "Flirds":(1.000,1.000),"Flirds1st":(1.000,0.670),"loss-heur":(1.000,1.000),"FLDetector":(0.987,0.983),
 "STD-DAGMM":(1.000,0.983),"FedDQC":(1.000,1.000),"FLTrust":(0.650,0.498),"FedIF":(0.542,0.458),"ComFedSV":(0.778,0.727),
}

M3B = ["in-run oracle","Flirds","Flirds1st","FedIF","loss-heur","FLDetector","STD-DAGMM","FLTrust","FedDQC"]
M3B_AUROC = {"in-run oracle":(1.000,1.000,1.000,1.000),"Flirds":(1.000,1.000,1.000,0.000),"Flirds1st":(1.000,1.000,1.000,0.000),
 "FedIF":(1.000,1.000,1.000,1.000),"loss-heur":(1.000,1.000,1.000,1.000),"FLDetector":(1.000,1.000,1.000,1.000),
 "STD-DAGMM":(0.250,1.000,0.000,0.750),"FLTrust":(1.000,1.000,1.000,1.000),"FedDQC":(1.000,0.750,0.750,1.000)}
M3B_SP = {"Flirds":(1.000,1.000,1.000,0.000),"Flirds1st":(1.000,1.000,1.000,0.000),"FedIF":(0.600,0.600,0.600,0.600),
 "loss-heur":(1.000,1.000,1.000,1.000)}
M3B_RT = {"in-run oracle":1240,"Flirds":251,"Flirds1st":82,"FedIF":82,"loss-heur":384,"FLDetector":250,"STD-DAGMM":530,"FLTrust":85,"FedDQC":49}

# C1 비교 방법 10종 (순위 표시 순; in-run oracle 은 기준=1.0 이라 열에서 제외)
C1_METHS = ["Flirds","Flirds1st","Banzhaf","GTG","FedSV","ComFedSV","ShapleyFL","FedIF","loss-heur","Ripple"]
# LLM scale fidelity ρ 표시 순
SCALE_METHS = ["Flirds","Flirds1st","loss-heur","Banzhaf","GTG","FedSV","ShapleyFL","ComFedSV","FedIF"]
# C2 분할/arm
PART_LABEL = {"iid":"iid (IID)", "dir1":"dir1 (Non-IID)", "shard":"shard (Non-IID 강)"}
DS_LABEL   = {"cifar10":"CIFAR-10", "fmnist":"FMNIST"}
C2_COMMON  = ["vanilla","flirds_mult","flirds_select","shapleyfl","fedif","sfedavg"]
C2_DIR1ARM = ["vanilla","flirds_mult","flirds_repl","flirds_add","flirds_select","shapleyfl","fedif","sfedavg"]
C2_THREATS = ["clean","grad_noise","label_flip","free_rider"]

# ════════════════════════════════════════════════════════════════════
# 슬라이드  (검증 질문별 구성)
# ════════════════════════════════════════════════════════════════════

# ── S1 표지 ─────────────────────────────────────────────────────────
s = slide()
tb = _box(s, 0.9, 2.55, 11.5, 1.1); tf = tb.text_frame; tf.word_wrap = True
_set(tf.paragraphs[0], "Flirds — 연합학습 클라이언트 기여도 측정", 36, INK, bold=True)
tb = _box(s, 0.9, 3.55, 11.5, 0.7)
_set(tb.text_frame.paragraphs[0], "실험 결과 보고 — 검증 질문별 구성", 21, ACCENT, bold=True)
tb = _box(s, 0.9, 4.35, 11.5, 0.9); tf = tb.text_frame; tf.word_wrap = True
_set(tf.paragraphs[0], "방법: 검증손실의 1차+2차 테일러 전개로 클라이언트 Shapley 추정 (라운드당 HVP 1회)", 14.5, SUB)
p = tf.add_paragraph(); _set(p, "구성: 1차 Fidelity → 2차 ① Performance ② Convergence ③ Detection (모델 무관, 질문별 묶음)", 13, FAINT)
ln = s.shapes.add_shape(1, Inches(0.92), Inches(2.42), Inches(4.2), Pt(2.2))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
tb = _box(s, 0.9, 6.5, 11.5, 0.5)
_set(tb.text_frame.paragraphs[0], "2026-06-19 · 모든 수치 ⓑ실측(세팅 병기), 미실행분 ⓒ · 표 수치 부호 +생략(음수만 −) · 분할 IID/Non-IID 병기", 12.5, FAINT)

# ── S2 실험 개요 (검증 질문별) ──────────────────────────────────────
s = slide()
header(s, "Overview", "실험 개요 — 무엇을 검증했는가 (질문별 분류)")
rows = [["검증 질문","무엇을 확인하나","실험 (무대)"],
 ["1차  Fidelity","기여도(Shapley)를 오라클 대비 얼마나 정확히 추정하나","CNN C1 (듀얼 오라클, N=10) · LLM standard (N=20/5) · LLM scale (silo·device·3B)"],
 ["2차-①  Performance","측정한 기여도로 가중/선택 시 일반 성능이 오르나·무해한가","CNN intervention (N=100, IID·Non-IID) · LLM standard (N=20)"],
 ["2차-②  Convergence","수렴 속도 — 목표 도달 라운드·최종 손실","LLM standard (N=20)"],
 ["2차-③  Detection","오염 클라이언트를 분리하나 (AUROC)","CNN intervention · LLM cross-silo·cross-device·3B"]]
table(s, 0.55, 1.55, 12.23, rows, [1.5,4.0,4.6], ["l","l","l"], size=11.5, row_h=0.62, accent_first=True)
notes(s, 0.55, 5.05, 12.2, [
 ("핵심 질문 위계 (발표 순서 = 이 순서):", True),
 "  1차가 가장 기본 — 기여도를 정확히 재는가. 2차는 그 기여도의 실효성(성능 → 수렴 → 탐지 순).",
 ("듀얼 오라클: in-run oracle(학습궤적 exact 2ᴺ 분해) · retrain oracle(부분집합 재학습→검증손실). 비교 방법 11종 + 탐지기 4종.", False, SUB),
 ("같은 robustness run이 fidelity(ρ)와 detection(AUROC)을 함께 산출 — 각각 해당 섹션에 배치.", False, SUB),
], size=12, gap=6)
footnote(s, "공통 규약: SGD momentum=0 · 시드별 seed_everything · 모든 셀 config+meta(git/env)+phi+metrics 영속화. 실패 0. 총 셀 181(CNN 150 + LLM 31).")

# ── S3 비교 방법·오라클 (공용 용어집) ───────────────────────────────
s = slide()
header(s, "용어 · 공용", "비교 방법(valuation)과 오라클 정의")
rows = [["방법","유형","핵심","라운드당 비용"],
 ["Flirds (제안)","valuation","검증손실 1차+2차 테일러","HVP 1회"],
 ["Flirds-1st","valuation","1차항만 (곡률 제외)","grad 1회"],
 ["loss-heur","heuristic","클라별 손실차","fwd 1회"],
 ["GTG / FedSV","Shapley(MC)","truncated/permutation MC","수십–수백 평가"],
 ["Banzhaf","exact","2ᴺ Banzhaf","2ᴺ 평가"],
 ["ShapleyFL / ComFedSV","Shapley(변형)","uniform-util / 부분합","수십–수백"],
 ["FedIF","influence","검증그래디언트 정렬","grad 1회"],
 ["Ripple","2nd-order","고유분해 기반 (이미지만)","eigsh (수십분)"],
 ["탐지기 4종","detector","FLDetector·STD-DAGMM·FLTrust·FedDQC","모델-free/grad"]]
table(s, 0.55, 1.55, 7.05, rows, [2.0,1.2,2.5,1.5], ["l","l","l","l"], size=10.5, row_h=0.43, accent_first=True)
notes(s, 7.85, 1.62, 4.9, [
 ("오라클 (정답 기여도)", True),
 ("in-run oracle", True, ACCENT),
 "  학습 궤적을 라운드별 exact Shapley로 분해 (2ᴺ).",
 "  추정량이 재현해야 할 게임. N=5에서 2⁵.",
 "",
 ("retrain oracle", True, ACCENT),
 "  부분집합 S를 FedAvg 재학습 → 배포모델 검증손실.",
 "  '배포 가치' 게임. 이미지 N=10에서 2¹⁰=1024 재학습.",
 "",
 ("핵심: Flirds는 in-run oracle 게임의 추정량.", True),
 "retrain oracle와의 일치는 별개 질문(1차 Fidelity 섹션).",
 "intervention arm(가중·선택 전략) 정의는 2차 섹션 첫 장.",
], size=11.5, gap=4)
footnote(s, "출처: codes/core·oracle·baselines. 모든 valuation은 동일 logs=[(w_r, {client:(Δ,n_c)})] 위에서 계산 (모델·검증셋 비접근). · 내부코드 in-run oracle=(b)·retrain oracle=(a)")

# ════════ [1차 Fidelity] ═══════════════════════════════════════════

# ── S4 CNN fidelity vs in-run ───────────────────────────────────────
s = slide()
header(s, "1차 Fidelity · Image (CNN)", "in-run oracle 대비 기여도 순위 — 비교 방법 전체(10종)")
rows = [["시나리오 (CIFAR / MNIST)"] + C1_METHS]; emph = set()
for ds, label in (("cifar10","CIFAR"), ("mnist","MNIST")):
    for sc in SCEN_ORDER:
        k = _scen_key(ds, sc)
        rows.append([f"{label}·{sc}"] + [fmt(C1_SP[k].get(m), nd=2) for m in C1_METHS])
    rows.append([f"{label} 평균"] + [fmt(_avg([C1_SP[_scen_key(ds,x)].get(m) for x in SCEN_ORDER]), nd=2) for m in C1_METHS])
    emph.add(len(rows) - 1)
table(s, 0.55, 1.5, 12.23, rows, [1.75]+[1.0]*10, ["l"]+["c"]*10, size=8.3, row_h=0.315, zebra=True, emph_rows=emph)
notes(s, 0.55, 5.95, 7.4, [
 ("명세", True, ACCENT),
 "MNIST+LeNet5 / CIFAR-10+CNN · N=10 full · R=10 · SGD mom=0 lr0.01 · val=2000 · 3 seed · in-run oracle 2¹⁰",
 ("분할: iid=IID · label_skew·quantity_skew=Non-IID · label_flip·feature_noise=IID 분할+손상", True, ACCENT),
], size=10, gap=3)
notes(s, 8.2, 5.95, 4.6, [
 ("지표", True, ACCENT),
 "Spearman ρ ↑ (오라클 φ 대비 순위). 평균=데이터셋 단순평균.",
 "Flirds·Banzhaf 순위 최상; MC기반 GTG/FedSV/ShapleyFL·ComFedSV·Ripple은 noisy.",
 "2차항 고유 기여는 순위 아닌 값-거리(다음장 하단).",
], size=10, gap=3)
footnote(s, "iid·feature_noise는 in-run oracle 신호 약함(클라 동질). 비교 방법 10종 전체(로스터 11종 = 10종 + in-run oracle 기준=1.0). 출처: runs/track_c/c1/* 재집계.")

# ── S5 CNN fidelity vs retrain ──────────────────────────────────────
s = slide()
header(s, "1차 Fidelity · Image (CNN)", "retrain oracle (2¹⁰) 대비 순위, 그리고 in-run과의 괴리")
rows = [["시나리오"] + C1_METHS]; emph = set()
for ds, label in (("cifar10","CIFAR"), ("mnist","MNIST")):
    for sc in SCEN_ORDER:
        a = C1_A[_scen_key(ds, sc)]
        rows.append([f"{label}·{sc}"] + [fmt(a.get(m), nd=2) for m in C1_METHS])
    rows.append([f"{label} 평균"] + [fmt(_avg([C1_A[_scen_key(ds,x)].get(m) for x in SCEN_ORDER]), nd=2) for m in C1_METHS])
    emph.add(len(rows) - 1)
table(s, 0.55, 1.5, 12.23, rows, [1.75]+[1.0]*10, ["l"]+["c"]*10, size=8.3, row_h=0.315, zebra=True, emph_rows=emph)
notes(s, 0.55, 5.95, 7.4, [
 ("핵심 — 게임 정의 차이 (label_skew)", True, ACCENT),
 "Flirds는 in-run oracle은 거의 완벽 재현하나(앞장), retrain oracle와는 label_skew서 갈림:",
 ("  CIFAR in-run 0.98 → retrain −0.18 · MNIST 0.71 → retrain −0.29", True, ACCENT),
 "추정오차 아님 — 모든 in-run 방법 공유(희귀라벨 클라는 라운드 한계기여 낮아도 최종 재학습 모델엔 필수).",
], size=10, gap=3)
notes(s, 8.2, 5.95, 4.6, [
 ("지표 / 2차항 효과", True, ACCENT),
 "Spearman ρ ↑ (retrain oracle φ 대비). 평균=데이터셋 단순.",
 "2차항: euclid_d(값-거리) label_flip CIFAR 0.031 vs 1차 0.046",
 "(MNIST 0.207 vs 0.370) → 값 오차 1차 대비 ~25–50%↓.",
], size=10, gap=3)
footnote(s, "표=Spearman ρ ↑ (rho_a_mean, 3-seed). LLM standard reference(N=5)는 retrain oracle와 1.000 일치 → 괴리는 이질성 구동(이어지는 장). 출처: RESULTS.txt.")

# ── S6 LLM standard fidelity (1B + 3B) ──────────────────────────────
s = slide()
header(s, "1차 Fidelity · LLM standard (clean·IID)", "기여도 fidelity — 듀얼 오라클 · 1B/3B")
hdr = ["방법","Spearman ↑","Kendall ↑","cosine_d ↓","euclid_d ↓","max_diff ↓"]
def drow(F, me):
    return [me, fmt(F["spearman"].get(me),nd=3), fmt(F["kendall"].get(me),nd=3),
            f"{F['cosine_d'][me]:.1e}", f"{F['euclid_d'][me]:.1e}", f"{F['max_diff'][me]:.1e}"]
caption(s, 0.55, 1.5, 6.2, f"standard 1B (N=20) — vs in-run oracle  ({D_STD_N} seed)")
order = ["Flirds","Flirds1st","loss-heur","GTG","FedSV","ComFedSV","ShapleyFL","FedIF"]
rows = [hdr] + [drow(D_STD_F, m) for m in order if m in D_STD_F["spearman"]]
table(s, 0.55, 1.82, 6.2, rows, [1.4,1.1,1.0,1.05,1.05,1.05], ["l"]+["c"]*5, size=9.0, row_h=0.32, zebra=True)
caption(s, 7.05, 1.5, 6.25, f"reference 1B (N=5) — vs retrain oracle  ({D_ANC_N} seed)")
order2 = ["Flirds","Flirds1st","Banzhaf","loss-heur","GTG","FedSV","ShapleyFL","ComFedSV","FedIF"]
rows2 = [hdr] + [drow(D_ANC_F, m) for m in order2 if m in D_ANC_F["spearman"]]
table(s, 7.05, 1.82, 6.25, rows2, [1.4,1.1,1.0,1.05,1.05,1.05], ["l"]+["c"]*5, size=9.0, row_h=0.30, zebra=True)
caption(s, 0.55, 4.95, 8.0, f"standard 3B (N=20) — vs in-run oracle  ({D3_STD_N} seed)")
order3 = ["Flirds","Flirds1st","loss-heur","GTG","FedSV","ComFedSV","ShapleyFL","FedIF"]
rows3 = [["방법"] + [m for m in order3 if m in D3_STD_F["spearman"]]]
rows3.append(["Spearman ↑"] + [fmt(D3_STD_F["spearman"].get(m),nd=3) for m in order3 if m in D3_STD_F["spearman"]])
rows3.append(["cosine_d ↓"] + [f"{D3_STD_F['cosine_d'][m]:.1e}" for m in order3 if m in D3_STD_F["spearman"]])
table(s, 0.55, 5.27, 8.0, rows3, [1.4]+[1.0]*8, ["l"]+["c"]*8, size=8.6, row_h=0.3, accent_first=True)
notes(s, 8.75, 4.9, 4.05, [
 ("읽기", True, ACCENT),
 ("Flirds: in-run ρ=1.000 · retrain ρ=1.000 —", True, ACCENT),
 ("  IID·clean서 두 오라클 일치 → 이미지(S5)", True, ACCENT),
 ("  괴리는 이질성 구동임을 확인.", True, ACCENT),
 "1B→3B Flirds ρ=1.000 유지(GTG 0.99·FedSV 0.94).",
 "2차항: cosine_d 3.5e-7 vs 1차 4.4e-5(~125×).",
 "ShapleyFL·FedIF·ComFedSV는 스케일서 약함.",
], size=9.2, gap=2.5)
footnote(s, "alpaca-gpt4 20k IID. standard=in-run oracle, reference=retrain oracle(2⁵). 출처: runs/track_d/rundirs/* 재집계. 내부코드 standard=std20·reference=anchor5.")

# ── S7 LLM scale fidelity (ρ 통합) ──────────────────────────────────
s = slide()
header(s, "1차 Fidelity · LLM scale", "Shapley 순위 fidelity — N=5 silo · N=100 device · 3B (vs in-run oracle)")
anc_sp = DEV_ANCHOR["noisy"][1]
rows = [["방법","silo N=5 ρ noisy ↑","silo N=5 ρ poison ↑","device N=100 ρ (α=0.5) ↑","3B N=5 ρ noisy ↑"]]
for m in SCALE_METHS:
    sp = SILO_SP.get(m); m3 = M3B_SP.get(m)
    rows.append([m,
                 fmt(sp[0],nd=3) if sp else "–",
                 fmt(sp[3],nd=3) if sp else "–",
                 fmt(anc_sp.get(m),nd=3) if m in anc_sp else "–",
                 fmt(m3[0],nd=3) if m3 else "–"])
table(s, 0.55, 1.6, 8.0, rows, [1.5,1.35,1.35,1.5,1.3], ["l","c","c","c","c"], size=9.6, row_h=0.4, zebra=True)
sb(s, 8.85, 1.6, 3.95, spec=[
    "· 동일 robustness run의 ρ만 모음",
    "  (같은 run의 탐지 AUROC → Detection 섹션)",
    "· silo N=5: 1 도메인/클라(Non-IID), 3 seed",
    "· device N=100: Dir(α=0.5), per-round oracle, 1 seed",
    "· 3B: silo N=5, 1 seed",
    "· 비싼 method는 sweep 전체엔 미실행 →",
    "  device 열은 α=0.5 기준점만(비용 게이팅)",
 ], metrics=[
    "Spearman ρ ↑ (in-run oracle φ 대비 순위)",
 ], read=[
    ("N=5서 near-additive → 전 방법 ρ≈1.000;", True),
    ("  poison·스케일에서 분화:", True),
    "  poison: FedSV 0.367·1차 0.000 (Flirds 0.967 유지)",
    "  N=100: GTG 0.78·FedSV 0.75·ShapleyFL 0.58·ComFedSV −0.02",
    ("→ Flirds·loss-heur는 전 스케일 ρ=1.000 유지;", True, ACCENT),
    ("  MC/변형 Shapley는 N↑·poison서 열화.", True, ACCENT),
 ])
footnote(s, "silo noisy/fr는 전 방법 1.000(near-additive)이라 poison 열로 분화를 표시. device 값=α=0.5 anchor(per-round exact). 빈칸=해당 스케일 미실행. 출처: RESULTS.md + rundir.")

# ── S8 종합 1차 ─────────────────────────────────────────────────────
s = slide()
header(s, "1차 Fidelity · 종합", "기여도 측정 정확성 — 횡단 요약")
rows = [["측면","결과","근거"],
 ["순위 vs in-run oracle","Flirds ρ≈1.000 (전 트랙·1B→3B). Banzhaf·loss-heur도 우수; MC-SV는 noisy","cross-silo·standard 1.000 / 이미지 0.71–1.00"],
 ["값-거리 vs in-run oracle","2차항이 1차 대비 값오차 ~2× 축소; loss-heur은 순위만 비기고 크기는 빗나감","euclid_d: 이미지 ~0.5× · standard cosine 125×"],
 ["순위 vs retrain oracle","IID·clean은 1.000 일치; label_skew(이질성)에서 in-run oracle와 괴리(−)","reference(N=5) 1.000 / 이미지 label_skew −0.18~−0.29"],
 ["스케일 (N=100)","α=0.5 기준점서 Flirds·1차·loss-heur ρ=1.000; GTG/FedSV/ShapleyFL은 열화","GTG 0.78·FedSV 0.75·ShapleyFL 0.58·ComFedSV −0.02"],
 ["비용","Flirds 1 HVP/round; in-run oracle exact의 5–15× 저렴, 동일 순위","cross-silo 106s vs 530s · 3B 251s vs 1240s"]]
table(s, 0.55, 1.55, 12.23, rows, [1.7,4.7,3.5], ["l","l","l"], size=10.8, row_h=0.62, accent_first=True)
notes(s, 0.55, 5.55, 12.2, [
 ("한 줄 요약", True),
 "Flirds는 in-run oracle 게임을 순위·값 모두에서 충실히, exact 대비 저비용으로, 1B→3B·N=5→N=100서 일관되게 재현한다.",
 "2차항의 고유 기여는 순위가 아니라 값-크기 정밀도. retrain oracle 게임과의 일치는 데이터 이질성에 의존(동질↔label_skew).",
], size=11.5, gap=5)
footnote(s, "retrain oracle 괴리는 Flirds 특정 오차가 아니라 모든 in-run 방법 공유. 상세: 앞의 fidelity 페이지들.")

# ════════ [2차-① Performance / ② Convergence] ══════════════════════

# ── S9 intervention arm 정의 ────────────────────────────────────────
s = slide()
header(s, "2차 · intervention", "intervention arm 정의 — 측정한 기여도를 집계에 어떻게 쓰나")
rows = [["arm","집계 가중 규칙","설명"],
 ["vanilla","w ∝ n_c","plain FedAvg (샘플 수만 가중 — baseline)"],
 ["flirds_mult","w ∝ n_c · s_c","크기가중 × Flirds 점수 (soft, 기본형)"],
 ["flirds_repl","w ∝ s_c","크기 무시, 점수로 대체"],
 ["flirds_add","w ∝ ½·(s_c/Σs) + ½·(n_c/Σn)","점수-비중과 크기-비중 5:5 혼합 (λ=0.5)"],
 ["flirds_select","점수 softmax로 코호트 선발 → 저기여 제외","hard selection (S-FedAvg式)"],
 ["shapleyfl / fedif","동일 기계, 점수만 교체","ShapleyFL 라운드 exact Shapley / FedIF 영향함수"],
 ["sfedavg","자체 점수 softmax 선택","S-FedAvg baseline"]]
table(s, 0.55, 1.55, 8.1, rows, [1.6,3.0,3.4], ["l","l","l"], size=10.2, row_h=0.52, accent_first=True)
notes(s, 8.85, 1.62, 3.95, [
 ("'arm' 이란", True, ACCENT),
 "비교하는 하나의 조건/변형 (임상시험 treatment arm).",
 "여기선 'φ를 어떻게 활용할지'의 각 전략이 한 arm.",
 "",
 ("점수 s_c", True, ACCENT),
 "Flirds 추정점수의 EMA (β=0.5, 높을수록 기여 큼).",
 "clean에선 점수 평탄 → vanilla로 복귀 = do-no-harm.",
 "",
 ("LLM standard 표기", True, ACCENT),
 "flirds_w=mult(β.5) · flirds_sel=select ·",
 "shapleyfl_w · fedif_w (β=0.7).",
], size=10, gap=3)
footnote(s, "출처: codes/flirds/fl/intervene.py rule_weights. repl/add 는 dir1(크기 skew)서만 mult 와 갈림 → iid/shard 는 6 arm 공통, dir1 만 8 arm.")

# ── S10 CNN intervention — IID vs Non-IID partition (merged) ────────
s = slide()
header(s, "2차-① Performance · Image (CNN)", "intervention — IID vs Non-IID partition (clean, do-no-harm)")
rows = [["dataset · partition"] + C2_COMMON]; emph = set()
for ds in ("cifar10","fmnist"):
    for part in ("iid","dir1","shard"):
        d = C2F.get((ds,part,"clean","strmain"), {})
        rows.append([f"{DS_LABEL[ds]}·{PART_LABEL[part]}"] + [fmt(d.get(a,(None,))[0]) for a in C2_COMMON])
table(s, 0.55, 1.6, 8.5, rows, [2.3,1.0,1.05,1.0,0.95,0.95,0.95], ["l"]+["c"]*6, size=9.6, row_h=0.46, zebra=True)
sb(s, 9.25, 1.6, 3.55, spec=[
    "· CIFAR-10 / FMNIST + CNN · N=100",
    "· C=0.1 (라운드당 10) · R=120 · lr0.01",
    "· partition iid(IID)·dir1(Dir α=1)·shard(2클래스/클라)",
    "· clean = 위협 없음 → do-no-harm 검증 · 3 seed",
    "· arm 6종 = iid/shard 공통(dir1만 repl/add 추가)",
 ], metrics=[
    "test acc ↑ (오염 없음 → 탐지 AUROC 없음)",
    "비교 대상 = intervention arm(가중·선택 전략)",
 ], read=[
    ("clean: 전 weighting arm ≈ vanilla → 분할 무관 무해.", True),
    "  (CIFAR shard만 flirds_select<vanilla = hard-drop 부작용)",
    "분할 난이도가 절대 정확도를 지배(shard=극단 Non-IID):",
    "  CIFAR 0.65→0.48 · FMNIST 0.86→0.69.",
    ("→ 가중은 어느 분할에서도 clean 성능을 해치지 않음.", True, ACCENT),
 ])
footnote(s, "한 표에 두 데이터셋×세 분할 통합(앞 fidelity 표와 동일 형식). 위협 하 결과는 다음 장. arm 정의=앞 장. 출처: runs/track_c/RESULTS.txt. 내부코드 intervention=C2.")

# ── S11 CNN intervention — under threat ─────────────────────────────
s = slide()
header(s, "2차-① Performance · ③ Detection · Image (CNN)", "intervention — under threat (CIFAR-10, dir1=Non-IID, 8 arm)")
rows = [["arm"] + C2_THREATS + ["위협평균 acc"]]
for a in C2_DIR1ARM:
    row = [a]; accs = []
    for ti, t in enumerate(C2_THREATS):
        acc, au = C2F.get(("cifar10","dir1",t,"strmain"), {}).get(a, (None, None))
        row.append((f"{acc:.3f} " + (f"({au:.2f})" if au is not None else "(–)")) if acc is not None else "–")
        if ti > 0 and acc is not None: accs.append(acc)
    row.append(fmt(_avg(accs)))
    rows.append(row)
table(s, 0.55, 1.6, 8.25, rows, [1.4,1.45,1.5,1.5,1.5,1.3], ["l","c","c","c","c","c"], size=9.6, row_h=0.46, accent_first=True)
sb(s, 9.05, 1.6, 3.8, spec=[
    "· CIFAR-10 + CNN · N=100, C=0.1 · R=120",
    "· partition dir1 (Dirichlet α=1, Non-IID)",
    "· arm 8종 (앞 장 정의)",
    "· strength sweep: grad0.05·flip0.6/0.8 별도",
 ], metrics=[
    "값 = test acc ↑ · ( ) = 탐지 AUROC ↑",
    "clean = 위협 無 → (–) · 위협평균 = 오염3종 acc",
 ], read=[
    ("under threat: 전 weighting arm > vanilla.", True),
    "grad_noise: shapleyfl0.550·fedif0.527 최고",
    "  > flirds_mult0.433 (vanilla 0.245).",
    "label_flip: flirds_mult0.585 최고.",
    "free_rider: fedif AUROC0.996, flirds0.385(약).",
    ("→ soft(mult) > hard(select). robust+무해,", True, ACCENT),
    ("  정확도 최고는 위협별로 갈림.", True, ACCENT),
 ])
footnote(s, "free-rider(zero-delta) AUROC가 N=100 CNN서 낮음(0.33–0.48) — LLM(1.0)과 상충, 미해결(한계 참조). FMNIST·strength sweep 동일 경향. 출처: RESULTS.txt.")

# ── S12 LLM standard performance + convergence ──────────────────────
s = slide()
header(s, "2차-① Performance · ② Convergence · LLM standard", "intervention arm — MMLU · ROUGE-L · 수렴 (standard 1B N=20)")
arm_order = ["base","vanilla","flirds_w","flirds_sel","shapleyfl_w","fedif_w"]
rows = [["arm","MMLU (0-shot) ↑","ROUGE-L ↑","final val-loss ↓","rounds→target ↓"]]
for a in arm_order:
    if a not in D_STD_A: continue
    A = D_STD_A[a]
    rows.append([a, fmt(A.get("mmlu"),nd=4), fmt(A.get("rouge_l"),nd=4),
                 fmt(A.get("final_val_loss"),nd=4) if A.get("final_val_loss") else "–",
                 f"{A['rounds_to_target']:.0f}" if A.get("rounds_to_target") else "–"])
table(s, 0.55, 1.65, 8.2, rows, [1.5,1.7,1.2,1.6,1.6], ["l","c","c","c","c"], size=10.5, row_h=0.46, accent_first=True)
sb(s, 9.0, 1.62, 3.8, spec=[
    "· Llama-3.2-1B / alpaca-gpt4 20k (IID)",
    "· standard N=20, 2/round · R=200",
    "· 10 steps×batch16 · lr1e-3 · seq512",
    "· arm 6종 · 3 seed · clean·IID",
    "· 3B는 N=20 intervention 미실행(N=5 reference만)",
 ], metrics=[
    "MMLU(0-shot) ↑ · ROUGE-L ↑   ① Performance",
    "final val-loss ↓ · rounds→target ↓   ② Convergence",
 ], read=[
    ("base = 학습 전. FL-SFT 후 ROUGE-L +0.067", True),
    "  (0.217→0.284), MMLU −0.008.",
    "전 weighting arm ≈ vanilla (±0.001, ±1라운드).",
    ("→ clean·IID서 do-no-harm parity 확인", True, ACCENT),
    ("  (고칠 오염 없을 때 성능·수렴 무해).", True, ACCENT),
 ])
footnote(s, "rounds-to-target = 검증손실 목표 도달 라운드(전 arm ~199/200). reference(N=5, 1B·3B)도 동일 parity. 출처: 재집계.")

# ════════ [2차-③ Detection] ════════════════════════════════════════

# ── S13 LLM cross-silo detection ────────────────────────────────────
s = slide()
header(s, "2차-③ Detection · LLM cross-silo (N=5, 1B)", "오염 클라 탐지 AUROC (Non-IID 5도메인)")
rows = [["방법","noisy ↑","fr-rand ↑","fr-zero ↑","poison ↑","평균 ↑","런타임 ↓"]]
for me in SILO_M:
    au = SILO_AUROC[me]
    rows.append([me, fmt(au[0]),fmt(au[1]),fmt(au[2]),fmt(au[3]), fmt(_avg(au)), f"{SILO_RT[me]}s"])
table(s, 0.55, 1.55, 8.3, rows, [1.7,1.0,1.0,1.0,1.0,0.95,1.0], ["l"]+["c"]*6, size=9.2, row_h=0.31, zebra=True)
sb(s, 9.05, 1.55, 3.8, spec=[
    "· Llama-3.2-1B(+3B) LoRA · N=5",
    "· 1 도메인/클라 → Non-IID",
    "· R=10 · val=100 · lr1e-3(poison 2e-3)",
    "· fp32 · eager · 3 seed",
    "· 위협: noisy(품질)/fr(zero,random)/",
    "        poison(백도어 스케일교체)",
 ], metrics=[
    "AUROC ↑ (오염=높은 의심점수 분리, 1.0=완전)",
    "평균 = 4위협 단순평균 · 런타임 ↓",
    "fidelity(ρ)는 1차 Fidelity 섹션(S7)에",
 ], read=[
    ("noisy·free-rider: 전 valuation AUROC 1.0", True),
    "  (N=5 near-additive); 우위 = 비용.",
    ("poison ASR≈1.0: Flirds-1st 0.000 완전회피 →", True, ACCENT),
    ("  2차 0.917 회복(불안정). loss-heur·전 탐지기 1.0.", True, ACCENT),
 ])
footnote(s, "poison Flirds(2차)는 run간 비결정(동일 config·seed서 0.42↔0.92). 평균=4위협 단순평균(참고용). 출처: RESULTS.md. 내부코드 cross-silo=silo5.")

# ── S14 LLM cross-device detection ──────────────────────────────────
s = slide()
header(s, "2차-③ Detection · LLM cross-device (N=100, 1B)", "α-sweep 탐지 AUROC (α = 분할 이질성)")
caption(s, 0.55, 1.52, 6.5, "noisy 탐지 AUROC ↑  (α-sweep)")
rows = [["방법"] + [f"α={a}" for a in ALPHAS] + ["평균"]]
for me in ["Flirds","FedIF","STD-DAGMM","FLTrust","FedDQC","FLDetector","ComFedSV"]:
    rows.append([me] + [fmt(x) for x in DEV_NOISY[me]] + [fmt(_avg(DEV_NOISY[me]))])
table(s, 0.55, 1.84, 6.55, rows, [1.5,0.85,0.85,0.85,0.85,0.85,0.9], ["l"]+["c"]*6, size=9.6, row_h=0.355, zebra=True)
caption(s, 7.3, 1.52, 5.5, "poison 탐지 AUROC ↑")
rows2 = [["방법","α=0.0","α=0.5","평균"]]
for me in ["Flirds","Flirds1st","loss-heur","FLDetector","STD-DAGMM","FedDQC","FLTrust","FedIF"]:
    rows2.append([me, fmt(DEV_POISON[me][0]), fmt(DEV_POISON[me][1]), fmt(_avg(DEV_POISON[me]))])
table(s, 7.3, 1.84, 5.5, rows2, [1.8,0.95,0.95,0.95], ["l","c","c","c"], size=9.6, row_h=0.335, zebra=True)
notes(s, 0.55, 4.95, 6.9, [
 ("명세", True, ACCENT),
 "· Llama-3.2-1B LoRA · N=100, K=10/round · R=30 · 1–3 seed",
 "· 분할 = per-client Dirichlet(α), α∈{0,0.01,0.1,0.5,5.0}",
 ("· α = 분할 이질성: α→0 도메인 분리(극단 Non-IID),", True, ACCENT),
 ("   α→∞ 균일혼합(IID). α=5.0 ≈ near-IID.", True, ACCENT),
], size=9.6, gap=2.5)
notes(s, 7.75, 4.95, 5.05, [
 ("읽기", True, ACCENT),
 "free-rider(random·zero, 전 α): Flirds·1차·",
 "  loss-heur·FLTrust = 1.000.",
 "noisy: FedDQC 영역(스케일 1.0), Flirds ~0.60.",
 ("poison(install): Flirds 1.0(both α), 1차도", True, ACCENT),
 ("  α=0서 1.0 — 1차가 회피되지 않음.", True, ACCENT),
], size=9.6, gap=2.5)
footnote(s, "fidelity(ρ α=0.5)는 1차 Fidelity 섹션(S7). cross-device φ는 선택된 클라만 기록(~96–98/100). 출처: RESULTS.md. 내부코드 device100.")

# ── S15 LLM 3B detection (+ CNN 탐지 메모) ──────────────────────────
s = slide()
header(s, "2차-③ Detection · LLM 3B (N=5)", "Llama-3.2-3B 탐지 AUROC (Non-IID 5도메인)")
rows = [["방법","noisy ↑","fr-rand ↑","fr-zero ↑","poison ↑","평균 ↑","런타임 ↓"]]
for me in M3B:
    au = M3B_AUROC[me]
    rows.append([me, fmt(au[0]),fmt(au[1]),fmt(au[2]),fmt(au[3]), fmt(_avg(au)), f"{M3B_RT[me]}s"])
table(s, 0.55, 1.6, 8.3, rows, [1.7,1.0,1.0,1.0,1.0,0.95,1.0], ["l"]+["c"]*6, size=9.6, row_h=0.4, zebra=True)
sb(s, 9.05, 1.6, 3.8, spec=[
    "· Llama-3.2-3B LoRA · cross-silo N=5",
    "· Non-IID(도메인 분리) · R=10 · seed 1개(잠정)",
    "· 비싼 method(Banzhaf/GTG/FedSV/ShapleyFL)는",
    "  3B서 비용상 미실행",
 ], metrics=[
    "AUROC ↑ · 평균=4위협 단순 · 런타임 ↓",
    "fidelity(ρ)는 1차 Fidelity 섹션(S7)",
 ], read=[
    ("noisy·fr: Flirds AUROC 1.0(1B→3B 유지),", True),
    "  251s vs in-run oracle 1240s (~5× ↓).",
    ("poison: Flirds·1차 모두 0.000 — 1B의 2차", True, ACCENT),
    ("  회복이 3B(1seed)선 안 나타남(미확정).", True, ACCENT),
    ("CNN 탐지(N=100): intervention arm의 누적 점수로", False, SUB),
    ("  측정 — S11 괄호값. free-rider만 0.4(상충).", False, SUB),
 ])
footnote(s, "3B는 seed 1개 → 잠정. 7B·N=10 retrain oracle은 설계상 deferred(ⓒ). 출처: RESULTS.md.")

# ── S16 종합 2차 ────────────────────────────────────────────────────
s = slide()
header(s, "2차 실효성 · 종합", "① Performance → ② Convergence → ③ Detection")
rows = [["질문","요약"],
 ["① Performance","under threat 전 weighting arm > vanilla(이미지, 전 분할). clean/IID는 무해 parity(이미지·LLM). 정확도 최고는 위협별로 갈림 — grad_noise서 shapleyfl/fedif, label_flip서 flirds. soft>hard. 분할 난이도(shard≪iid)가 절대성능을 지배."],
 ["② Convergence","IID·clean(LLM standard)서 전 arm 동일 수렴(rounds→target ~199/200, parity). 고칠 오염 없을 때 가중이 수렴을 해치지 않음."],
 ["③ Detection","free-rider: LLM AUROC 1.0(전 α)이나 N=100 CNN서 0.4(상충, 한계). noisy: FedDQC 영역(1.0), Flirds ~0.6. poison: 1차 회피→2차 부분회복(1B)·실패(3B); loss-heur·전 탐지기 1.0."]]
table(s, 0.55, 1.55, 12.23, rows, [1.5,9.0], ["l","l"], size=11.3, row_h=1.0, accent_first=True)
notes(s, 0.55, 5.95, 12.2, [
 ("Detection은 위계상 마지막 — 기여도≠탐지. clean-val-loss를 낮추는 공격자를 φ가 '기여 높음'으로 보는 것은 valuation의 정직한 답이며,", True),
 "  바로 그 지점(clean-preserving poison)이 valuation 접근의 경계다 (한계 페이지).",
], size=11.5, gap=5)

# ── S17 한계·미해결 ─────────────────────────────────────────────────
s = slide()
header(s, "한계 · 미해결", "솔직한 항목 (같은 톤으로)")
rows = [["항목","상태","내용"],
 ["오라클 게임 괴리 (retrain↔in-run)","ⓑ 실측","label_skew(이질성)서 in-run φ가 재학습 가치와 역상관. 모든 in-run 방법 공유. 프레이밍 필요."],
 ["free-rider 탐지 상충","ⓑ 미해결","동일 zero-delta가 LLM AUROC 1.0 vs N=100 CNN 0.4(IID칸 포함). 버그 배제용 드릴다운 필요."],
 ["poison 2차 불안정","ⓑ 잠정","1B서 1차 회피(0.000)→2차 부분회복하나 run간 0.42↔0.92(동일 config·seed). 3B(1seed)선 회복 실패."],
 ["미실행 (deferred)","ⓒ 미실행","LLM standard 7B · LLM N=10 retrain oracle · 7B robustness 전반 — 설계상 보류 (3B standard는 완료)."]]
table(s, 0.55, 1.55, 12.23, rows, [1.9,1.1,7.0], ["l","c","l"], size=11, row_h=0.8, accent_first=True)
notes(s, 0.55, 5.7, 12.2, [
 ("다음 단계 (제안)", True),
 "  1차: 오라클 게임 괴리(retrain↔in-run) 규명(φ 직접 분석, 희귀라벨 가설) — 핵심질문 위계상 우선.",
 "  병행: free-rider CNN-vs-LLM 상충 드릴다운(버그/실재 판별) · poison 2차 비결정성 재현 확인.",
], size=11.5, gap=6)
footnote(s, "ⓐ구현 ⓑ실측 ⓒ미실행. 잠정·미해결치도 결과와 같은 비중으로 기재.")

prs.save(str(OUT))
print(f"saved: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
