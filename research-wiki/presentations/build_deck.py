# -*- coding: utf-8 -*-
"""Flirds 연구 진행 보고 — 12-slide deck (python-pptx -> LibreOffice PDF)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
INK   = (15, 23, 42)      # main text (slate-900)
NAVY  = (17, 38, 64)      # dark headers / title bg
PRIMARY = (30, 78, 121)   # blue accent
TEAL  = (13, 118, 110)    # done / positive
AMBER = (176, 112, 16)    # in-progress / highlight
GRAY  = (90, 100, 116)    # secondary text
LINE  = (212, 218, 226)   # hairline
CARD  = (241, 245, 249)   # card bg
EDGE  = (205, 214, 226)   # card edge
WHITE = (255, 255, 255)
LIGHTTXT = (206, 216, 230) # body text on navy

FONT = "Noto Sans CJK KR"
ML = 0.62                  # left margin
CW = 13.333 - 2 * ML       # content width

def C(rgb): return RGBColor(*rgb)

def set_run_font(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def style_run(run, size, bold=False, color=INK, italic=False, name=FONT):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = C(color)
    set_run_font(run, name)

def _para(tf, first):
    if first:
        p = tf.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
    else:
        p = tf.add_paragraph()
    return p

def add_text(tf, text, size=13, bold=False, color=INK, align=PP_ALIGN.LEFT,
             first=False, sb=0, sa=0, line=1.0, name=FONT):
    p = _para(tf, first)
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after = Pt(sa)
    p.line_spacing = line
    r = p.add_run(); r.text = text
    style_run(r, size, bold, color, name=name)
    return p

def bullet(tf, text, lead=None, size=12.5, color=INK, mcolor=PRIMARY,
           first=False, sb=7, level=0, marker="•", line=1.05):
    p = _para(tf, first)
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(sb); p.line_spacing = line
    marL = 0.27 + level * 0.30
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(marL * 914400)))
    pPr.set("indent", str(int(-0.27 * 914400)))
    rm = p.add_run(); rm.text = marker + "  "; style_run(rm, size, True, mcolor)
    if lead:
        rl = p.add_run(); rl.text = lead; style_run(rl, size, True, color)
    rt = p.add_run(); rt.text = text; style_run(rt, size, False, color)
    return p

def box(slide, left, top, w, h, anchor=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if anchor: tf.vertical_anchor = anchor
    return tf

def rect(slide, left, top, w, h, fill, edge=None, ew=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    if edge is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = C(edge); sh.line.width = Pt(ew)
    sh.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sh.adjustments[0] = radius
        except Exception: pass
    return sh

def card(slide, left, top, w, h, fill=CARD, edge=EDGE, radius=0.045):
    return rect(slide, left, top, w, h, fill, edge=edge, ew=0.75,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)

def pill(slide, left, top, text, fill, tcolor=WHITE, w=0.92, h=0.30, size=10.5):
    sh = rect(slide, left, top, w, h, fill, edge=None,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; style_run(r, size, True, tcolor)
    return sh

def header(slide, title, eyebrow=None):
    if eyebrow:
        tf = box(slide, ML, 0.34, 11.5, 0.3)
        add_text(tf, eyebrow, size=11, bold=True, color=PRIMARY, first=True)
        ty = 0.60
    else:
        ty = 0.50
    tf = box(slide, ML, ty, 12.0, 0.7)
    add_text(tf, title, size=23, bold=True, color=NAVY, first=True)
    rect(slide, ML, ty + 0.62, 0.85, 0.052, TEAL)

def footer(slide, page):
    rect(slide, ML, 7.02, CW, 0.011, LINE)
    tf = box(slide, ML, 7.06, 9, 0.3)
    add_text(tf, "Flirds · 연구 진행 보고 · 2026-06-04", size=9, color=GRAY, first=True)
    tf2 = box(slide, 11.4, 7.06, 1.5, 0.3)
    add_text(tf2, "%02d / 12" % page, size=9, color=GRAY, align=PP_ALIGN.RIGHT, first=True)

def make_table(slide, matrix, left, top, total_w, fracs, row_h, font=11,
               header_row=True, aligns=None, hfill=NAVY, body_alt=CARD):
    rows = len(matrix); cols = len(matrix[0])
    gt = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                Inches(total_w), Inches(row_h * rows)).table
    gt.first_row = header_row; gt.horz_banding = False; gt.first_col = False
    for j, fr in enumerate(fracs):
        gt.columns[j].width = Emu(int(total_w * 914400 * fr))
    for i in range(rows):
        gt.rows[i].height = Inches(row_h)
        for j in range(cols):
            cell = gt.cell(i, j)
            cell.margin_left = Inches(0.09); cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            is_head = header_row and i == 0
            if is_head:
                cell.fill.fore_color.rgb = C(hfill)
            else:
                cell.fill.fore_color.rgb = C(WHITE if (i % 2 == 1) else body_alt)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.alignment = aligns[j] if aligns else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(matrix[i][j])
            style_run(r, font, bold=is_head, color=(WHITE if is_head else INK))
    return gt

def frun(p, text, base=0, size=17, bold=False, color=INK):
    r = p.add_run(); r.text = text; style_run(r, size, bold, color)
    if base:
        r._r.get_or_add_rPr().set("baseline", str(base))
    return r

SUP, SUB = 30000, -25000

# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def new_slide():
    return prs.slides.add_slide(BLANK)

# ---------- Slide 1 : Title ----------
s = new_slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 0.16, 7.5, TEAL)
tf = box(s, 1.1, 2.18, 11, 0.4)
add_text(tf, "연구 진행 보고   ·   2026-06-04", size=13, bold=True, color=AMBER, first=True)
tf = box(s, 1.05, 2.55, 11, 1.2)
add_text(tf, "Flirds", size=58, bold=True, color=WHITE, first=True)
tf = box(s, 1.1, 3.78, 11.2, 0.6)
add_text(tf, "Federated Learning  +  In-Run Data Shapley", size=23, bold=False, color=LIGHTTXT, first=True)
rect(s, 1.12, 4.46, 2.0, 0.045, TEAL)
tf = box(s, 1.1, 4.66, 10.9, 1.2)
add_text(tf, "Vanilla FedAvg가 이미 받는 클라이언트 업데이트 Δw_k 만으로, 라운드별 검증손실 변화의 "
              "1차+2차 Taylor 전개로 client-level Shapley 를 닫힌형 계산 — 통신 추가비용 0, 역헤시안 없음.",
         size=14, color=LIGHTTXT, first=True, line=1.25)
tf = box(s, 1.1, 6.45, 11, 0.4)
add_text(tf, "Yonghee Choi  ·  Edge AI Lab  ·  컴퓨팅 DGX B200 × 4", size=12, color=(150,165,188), first=True)

# ---------- Slide 2 : 한눈에 보기 ----------
s = new_slide()
header(s, "한눈에 보기", "EXECUTIVE SUMMARY")
# left column — what / how / strength / novelty
tf = box(s, ML, 1.55, 6.05, 5.0)
bullet(tf, "FL에서 client별 데이터 기여도 φ_k 를 측정. IRDS(In-Run Data Shapley)를 FL + LoRA(PEFT)로 확장.",
       lead="무엇:  ", first=True, sb=0)
bullet(tf, "라운드별 검증손실 변화의 1차+2차 Taylor 닫힌형. Δw_k 만 사용 → 통신 추가비용 0.",
       lead="어떻게:  ")
bullet(tf, "역헤시안 없이 forward HVP(H·Δw). LLM에서 IF가 무너지는 원인(iHVP collapse)을 구조적으로 회피.",
       lead="강점:  ")
bullet(tf, "client-level + in-run + 1·2차 Taylor + HVP 상호작용 + zero-comm + LoRA/LLM 교집합이 prior art에 비어 있음.",
       lead="novelty:  ")
# right column — status ladder
rx = 7.05
tf = box(s, rx, 1.5, 5.6, 0.3)
add_text(tf, "현재 위치", size=12, bold=True, color=GRAY, first=True)
rows = [("완료", TEAL, "Phase 0 — CNN baseline 4종 재현"),
        ("완료", TEAL, "Phase 0.5 — estimator + dual oracle · 전 게이트 green"),
        ("진행", AMBER, "Phase 1 — LLM 1B · 5단계 빌드+검증 · SMOKE green"),
        ("예정", GRAY, "Phase 2 / 3 — full baseline · 실험 매트릭스 실행")]
y = 1.92
for label, col, txt in rows:
    pill(s, rx, y + 0.02, label, col, w=0.84, h=0.30)
    tf = box(s, rx + 1.02, y, 4.6, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    add_text(tf, txt, size=12, color=INK, first=True, line=1.0)
    y += 0.62
# bottom callout
cc = card(s, ML, 6.05, CW, 0.78, fill=NAVY, edge=NAVY)
tf = cc.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_text(tf, "다음 액션", size=11, bold=True, color=AMBER, first=True)
add_text(tf, "1B FULL clean run 실행 (MINI de-risk 먼저, ~5–7h) → 이후 SV baselines LLM 이식 → Phase 2/3",
         size=13.5, bold=True, color=WHITE, sb=1)
footer(s, 2)

# ---------- Slide 3 : 문제 정의 & 접근 ----------
s = new_slide()
header(s, "문제 정의 — FL에서 누구의 데이터가 기여했는가", "PROBLEM")
tf = box(s, ML, 1.55, CW, 0.9)
add_text(tf, "학습된 FL 모델에 대해 각 client(데이터 보유자)에게 기여도 φ_k 를 부여한다. "
             "용도: 데이터 가치평가 · 저품질/free-rider 탐지 · 공정 보상 · 데이터 선택.",
         size=13.5, color=INK, first=True, line=1.2)
tf = box(s, ML, 2.5, CW, 0.35)
add_text(tf, "기존 방법의 벽", size=13, bold=True, color=PRIMARY, first=True)
# three problem cards
cards = [
    ("Retraining 기반 Shapley", "U(S) 평가에 2^N 재학습 필요 → LLM 스케일에서 불가능."),
    ("통신 / 계산 오버헤드", "FedSV는 O(Tm²) 서버 평가, ComFedSV는 추가 all-client 라운드."),
    ("Gradient / IF 방법", "heterogeneous FL서 취약(DataInf가 Fed-WildChat서 실패) · LLM서 iHVP collapse."),
]
cw3 = (CW - 0.6) / 3
for i, (t, b) in enumerate(cards):
    x = ML + i * (cw3 + 0.3)
    card(s, x, 2.85, cw3, 1.5)
    tf = box(s, x + 0.18, 3.0, cw3 - 0.36, 1.25)
    add_text(tf, t, size=12.5, bold=True, color=NAVY, first=True)
    add_text(tf, b, size=11.5, color=INK, sb=5, line=1.12)
# approach band
cc = card(s, ML, 4.7, CW, 1.55, fill=(235, 246, 244), edge=(180, 214, 208))
tf = box(s, ML + 0.22, 4.84, CW - 0.5, 1.3)
add_text(tf, "Flirds의 접근", size=13, bold=True, color=TEAL, first=True)
bullet(tf, "vanilla FedAvg 위에 post-hoc로 동작 — 학습 알고리즘 변경 없음.", color=INK, mcolor=TEAL, sb=6)
bullet(tf, "서버가 이미 받는 Δw_k(LoRA) 만 사용 → 통신·계산 추가비용 0.", color=INK, mcolor=TEAL, sb=5)
bullet(tf, "forward HVP(H·Δw)만 — H⁻¹ 없음 → LLM-IF의 실패 원인 회피.", color=INK, mcolor=TEAL, sb=5)
footer(s, 3)

# ---------- Slide 4 : 방법 ----------
s = new_slide()
header(s, "방법 — 라운드별 손실변화의 1·2차 Taylor", "METHOD")
# hero formula card
card(s, ML, 1.6, CW, 1.55, fill=NAVY, edge=NAVY)
tf = box(s, ML + 0.3, 1.86, CW - 0.6, 0.7, anchor=MSO_ANCHOR.MIDDLE)
p = _para(tf, True); p.alignment = PP_ALIGN.CENTER
frun(p, "φ", color=WHITE); frun(p, "k", SUB, 11, color=WHITE); frun(p, "(r)", SUP, 11, color=WHITE)
frun(p, "  ≈  −∇ℓ(w", color=WHITE); frun(p, "r", SUP, 11, color=WHITE)
frun(p, ", z", color=WHITE); frun(p, "val", SUB, 11, color=WHITE)
frun(p, ") · Δw", color=WHITE); frun(p, "k", SUB, 11, color=WHITE)
frun(p, "   +   ½ Δw", color=AMBER); frun(p, "k", SUB, 11, color=AMBER); frun(p, "T", SUP, 11, color=AMBER)
frun(p, " H", color=AMBER); frun(p, "(val)", SUP, 11, color=AMBER)
frun(p, "(w", color=AMBER); frun(p, "r", SUP, 11, color=AMBER); frun(p, ") ΔW", color=AMBER); frun(p, "(r)", SUP, 11, color=AMBER)
tf = box(s, ML + 0.3, 2.62, CW - 0.6, 0.4, anchor=MSO_ANCHOR.MIDDLE)
p = _para(tf, True); p.alignment = PP_ALIGN.CENTER
frun(p, "ΔW", 0, 12.5, color=LIGHTTXT); frun(p, "(r)", SUP, 8.5, color=LIGHTTXT)
frun(p, " = Σ", 0, 12.5, color=LIGHTTXT); frun(p, "j", SUB, 8.5, color=LIGHTTXT)
frun(p, " Δw", 0, 12.5, color=LIGHTTXT); frun(p, "j", SUB, 8.5, color=LIGHTTXT)
frun(p, "  (참여 client 합)        φ", 0, 12.5, color=LIGHTTXT); frun(p, "k", SUB, 8.5, color=LIGHTTXT)
frun(p, " = Σ", 0, 12.5, color=LIGHTTXT); frun(p, "r", SUB, 8.5, color=LIGHTTXT)
frun(p, " φ", 0, 12.5, color=LIGHTTXT); frun(p, "k", SUB, 8.5, color=LIGHTTXT); frun(p, "(r)", SUP, 8.5, color=LIGHTTXT)
frun(p, "  (라운드 합)", 0, 12.5, color=LIGHTTXT)
# bullets
tf = box(s, ML, 3.45, CW, 2.2)
bullet(tf, "라운드당 비용: HVP 1회로 u := H·ΔW^(r) 계산 후 N개 내적 Δw_k·u — LoRA 차원이라 저렴.",
       lead="비용:  ", first=True, sb=0)
bullet(tf, "1차항 = 정렬(alignment), 2차항 = client 간 상호작용(곡률). 비-IID FL에서 2차항이 본질적.",
       lead="2차항:  ")
bullet(tf, "forward HVP(H·Δw)만 — H⁻¹ 절대 계산 안 함 → IF의 iHVP collapse 원인 회피.",
       lead="안정성:  ")
# note band
cc = card(s, ML, 5.75, CW, 0.72, fill=CARD, edge=EDGE)
tf = box(s, ML + 0.2, 5.86, CW - 0.4, 0.55, anchor=MSO_ANCHOR.MIDDLE)
add_text(tf, "확정: 곡률 = 참 Hessian (GGN/Fisher 변형 검증 후 기각) · 최적화 = plain SGD(momentum 0) "
             "— 이때 비로소 2차항이 1차항을 이김.", size=11.5, color=GRAY, first=True, line=1.1)
footer(s, 4)

# ---------- Slide 5 : 차별점 & novelty ----------
s = new_slide()
header(s, "무엇이 새로운가 — 비어 있는 교집합", "POSITIONING")
matrix = [
    ["방법", "단위", "통신 오버헤드", "in-run", "비고"],
    ["FedSV (2020)", "client", "0 (단 O(Tm²) 평가)", "per-round 대리", "federated Shapley 원조"],
    ["Ripple (2026)", "sample", "0", "yes (Jacobian)", "최근접 경쟁자"],
    ["FedIF (2025)", "client", "0", "yes (1차 TracIn)", "CNN 한정 · aggregation-side"],
    ["FedTSV (2026)", "client", "+ val pass", "yes", "aggregation (가치평가 아님)"],
    ["Flirds (본 연구)", "client", "0", "yes (1·2차 Taylor)", "LoRA + 2차항 + zero-comm"],
]
make_table(s, matrix, ML, 1.6, CW, [0.17, 0.10, 0.21, 0.17, 0.35], 0.50, font=11.5,
           aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT])
# highlight last row by overlaying tint already? keep table; add accent note
tf = box(s, ML, 4.65, CW, 0.35)
add_text(tf, "novelty", size=13, bold=True, color=PRIMARY, first=True)
tf = box(s, ML, 5.0, CW, 1.9)
bullet(tf, "교집합(client-level + in-run + 1·2차 Taylor + HVP 상호작용 + zero-comm + LoRA/LLM)을 "
           "점유한 단일 논문 없음 — 2026-06-03 재조사(내부 4 + 외부 survey 교차검증)로 확인.",
       first=True, sb=0)
bullet(tf, "cross-domain valuation-fairness: 공유 val-loss Shapley의 도메인 간 공정성은 prior art서 "
           "under-addressed → 추가 novelty hook (모든 도메인 free-form 통일로 대응).")
bullet(tf, "FedIF/FedTSV는 aggregation-side(가중치 변경), Flirds는 valuation-side(post-hoc 가치) — 같은 입력, 다른 출력.")
footer(s, 5)

# ---------- Slide 6 : 이론 ----------
s = new_slide()
header(s, "이론 — 중앙집중 등가 + drift residual", "THEORY")
tf = box(s, ML, 1.55, CW, 1.0)
bullet(tf, "1-step SGD에서 한 client에 속한 data-level Shapley의 합 = client-level Shapley "
           "(1·2차항 모두). → 중앙집중 1-step에서 granularity 선택은 grouping-invariant.",
       lead="① 등가:  ", first=True, sb=0, line=1.15)
bullet(tf, "FL 멀티스텝: Δw_k 는 E-step 궤적의 끝점 → Taylor 전개 시 '중앙집중 등가항' + "
           "'client drift residual' O(ηE·|H|·궤적길이) 로 분해.",
       lead="② FL 편차:  ", sb=8, line=1.15)
# two proposition cards
cw2 = (CW - 0.3) / 2
for i, (t, b) in enumerate([
    ("Proposition 1", "Flirds client-level SV = (중앙집중 data-level SV를 client로 합산) + drift residual. "
                      "→ FL의 중앙집중 IRDS로부터의 이탈을 정량화."),
    ("Proposition 2", "drift residual은 라운드별 local 궤적 반경의 3차식으로 bound. "
                      "E=1서 소멸, non-IID·큰 local epoch서 증가."),
]):
    x = ML + i * (cw2 + 0.3)
    card(s, x, 3.55, cw2, 1.55, fill=(237, 242, 248), edge=(198, 212, 230))
    tf = box(s, x + 0.2, 3.7, cw2 - 0.4, 1.3)
    add_text(tf, t, size=13, bold=True, color=PRIMARY, first=True)
    add_text(tf, b, size=12, color=INK, sb=6, line=1.18)
# implication
cc = card(s, ML, 5.4, CW, 1.05, fill=(235, 246, 244), edge=(180, 214, 208))
tf = box(s, ML + 0.22, 5.54, CW - 0.5, 0.85, anchor=MSO_ANCHOR.MIDDLE)
add_text(tf, "핵심 함의", size=12, bold=True, color=TEAL, first=True)
add_text(tf, "IRDS는 \"2차항이 marginal\"이라 보고 — 그러나 이는 centralized per-step(미소 η)의 산물. "
             "FL의 per-round 멀티스텝이 바로 2차항이 본질적으로 작동하는 무대 (실험으로 검증).",
         size=12, color=INK, sb=3, line=1.18)
footer(s, 6)

# ---------- Slide 7 : 설계 결정 ----------
s = new_slide()
header(s, "확정된 핵심 설계 결정", "LOCKED DECISIONS")
matrix = [
    ["항목", "결정"],
    ["단위 (granularity)", "client-level"],
    ["서버 입력", "Δw_k (LoRA) 만 — 추가 통계/통신 0"],
    ["근사 차수", "1차 + 2차 Taylor, 항상 둘 다"],
    ["곡률", "참 Hessian (GGN 기각)"],
    ["최적화", "plain SGD (momentum 0)"],
    ["검증셋", "server-side held-out · 도메인당 200 × 5 = 1000 stratified"],
    ["PEFT / 정규화", "LoRA / participation 정규화 없음(기본)"],
    ["모델", "1B Llama-3.2 · 3B Llama-3.2 · 7B Llama-2"],
]
make_table(s, matrix, ML, 1.55, 7.4, [0.34, 0.66], 0.435, font=11.5,
           aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT])
# dual oracle on the right
rx = 8.25
tf = box(s, rx, 1.5, 4.5, 0.3)
add_text(tf, "Dual Oracle (ground truth)", size=12.5, bold=True, color=PRIMARY, first=True)
card(s, rx, 1.9, CW - (rx - ML), 1.55, fill=CARD, edge=EDGE)
tf = box(s, rx + 0.2, 2.04, CW - (rx - ML) - 0.4, 1.3)
add_text(tf, "(a) Exact retrain SV", size=12, bold=True, color=NAVY, first=True)
add_text(tf, "U(S)=S만으로 FL 학습 (데이터가치 표준). 1B N∈{5,10} · 3B N=5 · 7B ✗(compute).",
         size=11, color=INK, sb=3, line=1.15)
card(s, rx, 3.6, CW - (rx - ML), 1.7, fill=CARD, edge=EDGE)
tf = box(s, rx + 0.2, 3.74, CW - (rx - ML) - 0.4, 1.45)
add_text(tf, "(b) IRDS-定 in-run SV", size=12, bold=True, color=NAVY, first=True)
add_text(tf, "Flirds-correct oracle. cross-silo N=10 exact enum(1024 subset) · cross-device N=100 MC. "
             "비용 2^N·R·|val|·seq → N5↔N10 = 32× · fp32 필수.",
         size=11, color=INK, sb=3, line=1.15)
cc = card(s, rx, 5.45, CW - (rx - ML), 1.0, fill=NAVY, edge=NAVY)
tf = box(s, rx + 0.2, 5.56, CW - (rx - ML) - 0.4, 0.8, anchor=MSO_ANCHOR.MIDDLE)
add_text(tf, "원본 lock (conv4) + N1–N4 / Q1–Q3 / ④⑤ 까지 설계 전부 동결 — "
             "구현 중 새 정보 없는 한 재론 안 함.", size=10.5, color=WHITE, first=True, line=1.15)
footer(s, 7)

# ---------- Slide 8 : 진행 ① CNN ----------
s = new_slide()
header(s, "진행 현황 ① — Phase 0 / 0.5 완료 (CNN)", "PROGRESS · 1")
# Phase 0
pill(s, ML, 1.58, "완료", TEAL, w=0.84, h=0.30)
tf = box(s, ML + 1.0, 1.55, 10, 0.35)
add_text(tf, "Phase 0 — FL-Shapley baseline 4종 self-build 재현", size=13.5, bold=True, color=NAVY, first=True)
matrix = [
    ["GTG-Shapley", "FedSV", "ComFedSV", "Ripple"],
    ["recon cosine 0.99", "perm-MC 0.998", "Spearman {1.0,.96,.85,.84}", "noisy-AUROC 1.0"],
]
make_table(s, matrix, ML, 2.0, CW, [0.25, 0.25, 0.25, 0.25], 0.42, font=11,
           aligns=[PP_ALIGN.CENTER]*4, hfill=PRIMARY)
tf = box(s, ML, 2.92, CW, 0.3)
add_text(tf, "Ripple 속도이득 62× vs AFedSV+ / 49× vs FedSV (ground-truth-SV 지표 없음 → task-driven 검증).",
         size=10.5, color=GRAY, first=True)
# Phase 0.5
pill(s, ML, 3.42, "완료", TEAL, w=0.84, h=0.30)
tf = box(s, ML + 1.0, 3.39, 10, 0.35)
add_text(tf, "Phase 0.5 — Flirds estimator + dual oracle · 전 sanity gate green", size=13.5, bold=True, color=NAVY, first=True)
tf = box(s, ML, 3.85, CW, 1.7)
bullet(tf, "estimator ≈ (b) in-run oracle: Spearman 1.0 (3-seed 1+2차 0.96 > 1차 0.92).", first=True, sb=0)
bullet(tf, "noisy-client AUROC 1.0 · (b) Shapley efficiency = 0, symmetry = 0 (exact).")
bullet(tf, "HVP jvp-vs-double-backward 9.8e-6 · N=2 singleton 3e-3 · 재현성 bitwise-0.")
# decision band
cc = card(s, ML, 5.55, CW, 0.92, fill=(235, 246, 244), edge=(180, 214, 208))
tf = box(s, ML + 0.22, 5.66, CW - 0.5, 0.72, anchor=MSO_ANCHOR.MIDDLE)
add_text(tf, "결정 (데이터로 확인)", size=11, bold=True, color=TEAL, first=True)
add_text(tf, "2차항 = 참 Hessian · momentum 제거(plain SGD) 시 2차항이 1차항을 이김 → "
             "\"FL per-round가 2차항의 무대\"라는 가설을 CNN에서 실증.",
         size=11.5, color=INK, sb=2, line=1.15)
footer(s, 8)

# ---------- Slide 9 : 진행 ② LLM ----------
s = new_slide()
header(s, "진행 현황 ② — Phase 1 (LLM 1B) 거의 완료", "PROGRESS · 2")
tf = box(s, ML, 1.5, CW, 0.3)
add_text(tf, "estimator/oracle 무변경. 5단계 모두 빌드 + 검증 (real Llama-3.2-1B 궤적).",
         size=12, color=GRAY, first=True)
stages = [
    ("backend-agnostic estimator/oracle + partial-participation + per-layer φ", "CNN bit-identical 회귀"),
    ("LLM backend + FL loop (TRL SFTTrainer + forced SGD)", "1B est≈oracle 1.70e-6"),
    ("5-domain free-form data layer + val micro-batching + per-domain norm", "chunked==single 3.8e-8"),
    ("② corruptor: answer_swap(noisy) + free_rider(zero/random)", "free-rider φ = 정확히 0"),
    ("#7 first-clean-run 인프라 (eval · run_logger · orchestrator)", "SMOKE 1.6e-7 · AUROC 1.0"),
]
y = 1.92
for i, (t, num) in enumerate(stages):
    pill(s, ML, y + 0.04, str(i + 1), TEAL, w=0.34, h=0.30, size=11)
    tf = box(s, ML + 0.5, y, 8.6, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    add_text(tf, t, size=12, color=INK, first=True, line=1.0)
    tf = box(s, 9.3, y, 3.4, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    add_text(tf, num, size=11, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT, first=True)
    y += 0.52
# musts + remaining
cc = card(s, ML, 4.62, CW, 0.86, fill=CARD, edge=EDGE)
tf = box(s, ML + 0.2, 4.72, CW - 0.4, 0.68, anchor=MSO_ANCHOR.MIDDLE)
add_text(tf, "3 LLM musts (CNN엔 없던 것)", size=11, bold=True, color=NAVY, first=True)
add_text(tf, "eager attention(forward-AD) · FL state = named_parameters key · embedding require-grad hook clear (functorch ↔ HF).",
         size=11, color=INK, sb=2, line=1.12)
cc = card(s, ML, 5.6, CW, 0.88, fill=NAVY, edge=NAVY)
tf = box(s, ML + 0.2, 5.71, CW - 0.4, 0.7, anchor=MSO_ANCHOR.MIDDLE)
add_text(tf, "남은 1액션", size=11, bold=True, color=AMBER, first=True)
add_text(tf, "FULL scale run 실행 — N=5, R≈30, 3 seed, ORACLE_B(N=5), ~5–7h (MINI de-risk 먼저).",
         size=12.5, bold=True, color=WHITE, sb=1)
footer(s, 9)

# ---------- Slide 10 : 데이터 & 프로토콜 ----------
s = new_slide()
header(s, "데이터셋 & 평가 프로토콜", "DATA · PROTOCOL")
tf = box(s, ML, 1.5, CW, 0.3)
add_text(tf, "Cross-silo 5-domain — 전부 free-form instruction→response (공유 val-loss Shapley 공정성)",
         size=12.5, bold=True, color=NAVY, first=True)
matrix = [
    ["domain", "medical", "legal", "finance", "math", "general"],
    ["dataset", "med flashcards", "ibunescu QA", "FiQA", "AQUA-RAT", "Dolly"],
]
make_table(s, matrix, ML, 1.85, CW, [0.12, 0.19, 0.18, 0.14, 0.18, 0.19], 0.40, font=11,
           aligns=[PP_ALIGN.LEFT]+[PP_ALIGN.CENTER]*5, hfill=PRIMARY)
tf = box(s, ML, 2.74, CW, 0.6)
bullet(tf, "크기: 도메인당 train 12k / val 200 / test 2k (상호 disjoint) · cross-device: Fed-WildChat + FedHDS (N=100, K=10).",
       first=True, sb=0, size=11.5)
# eval
tf = box(s, ML, 3.4, CW, 0.32)
add_text(tf, "평가 — utility ≠ downstream 분리", size=12.5, bold=True, color=PRIMARY, first=True)
tf = box(s, ML, 3.72, CW, 1.2)
bullet(tf, "utility = val-loss (estimator/oracle가 소비) ≠ downstream = per-domain ROUGE-L + math(AQUA) EM.",
       first=True, sb=0, size=11.5)
bullet(tf, "selection-convergence: φ→top-K→retrain vs full/random (MATES 템플릿). "
           "caveat: heterogeneous FL서 selection이 random에 질 수 있음 — 넘어야 할 bar.", size=11.5)
# rigor band
cc = card(s, ML, 5.05, CW, 1.4, fill=CARD, edge=EDGE)
tf = box(s, ML + 0.22, 5.18, CW - 0.5, 1.18)
add_text(tf, "보고 엄밀성 (Protocol)", size=12, bold=True, color=NAVY, first=True)
bullet(tf, "수치: bf16 train / fp32 eval + fp32 내적 · ≥ 3 seed mean±std · 95% bootstrap CI.", sb=5, size=11.5)
bullet(tf, "sanity gate: E=1 ⇒ residual≈0 · N=2 ⇒ singleton SV = φ_k · 재현성 bitwise.", sb=4, size=11.5)
bullet(tf, "추적: local run-dir (config + git SHA + env hash + per-round φ parquet), W&B 미사용.", sb=4, size=11.5)
footer(s, 10)

# ---------- Slide 11 : 실험 계획 ----------
s = new_slide()
header(s, "실험 계획 — Section 3 (18 항목)", "EXPERIMENT PLAN")
groups = [
    ("★★★ Spine", PRIMARY,
     "baseline 10종 + detection 2(FLDetector·STD-DAGMM) · dual oracle (a)/(b) · "
     "Ripple head-to-head + 이론적 reduction 시도."),
    ("★★ 특성화 / ablation", PRIMARY,
     "α-sweep × E-sweep drift 매트릭스(16셀) · Q2 variants(3×3) · ②③ 특성화 · "
     "적대적 stress · non-IID valuation bias."),
    ("★ Scale", PRIMARY,
     "7B instruction-tuning bench (LESS · FedDQC 직접 비교)."),
]
y = 1.6
for t, col, b in groups:
    tf = box(s, ML, y, CW, 0.3)
    add_text(tf, t, size=12.5, bold=True, color=col, first=True)
    tf = box(s, ML, y + 0.30, CW, 0.7)
    add_text(tf, b, size=12, color=INK, first=True, line=1.15)
    y += 1.02
# flagship card
cc = card(s, ML, 4.72, CW, 1.0, fill=(252, 246, 233), edge=(226, 198, 140))
tf = box(s, ML + 0.22, 4.84, CW - 0.5, 0.8, anchor=MSO_ANCHOR.MIDDLE)
add_text(tf, "◆ 차별화 실험 (flagship 후보)", size=12, bold=True, color=AMBER, first=True)
add_text(tf, "PGD / direction-aligned poison: 1차항으로 못 잡는 공격을 2차(곡률)항이 분리하는가? "
             "— FedIF의 blind spot, Flirds-1st-only와 대조하여 \"2차 > 1차\"를 직접 입증.",
         size=11.5, color=INK, sb=2, line=1.15)
# matrix note
cc = card(s, ML, 5.85, CW, 0.62, fill=NAVY, edge=NAVY)
tf = box(s, ML + 0.2, 5.9, CW - 0.4, 0.5, anchor=MSO_ANCHOR.MIDDLE)
add_text(tf, "매트릭스: 1B / 3B / 7B 전 셀 (3 seed mean±std + 95% CI). "
             "est-vs-oracle fidelity는 비용 때문에 N으로 capped (N5↔N10 = 32×, fp32).",
         size=11, color=WHITE, first=True, line=1.1)
footer(s, 11)

# ---------- Slide 12 : 다음 단계 & 한계 ----------
s = new_slide()
header(s, "다음 단계 & 한계", "NEXT · LIMITATIONS")
# next steps (left)
tf = box(s, ML, 1.55, 6.1, 0.3)
add_text(tf, "다음 단계", size=13, bold=True, color=PRIMARY, first=True)
steps = [
    ("즉시", AMBER, "1B FULL clean run (N=5, R≈30, 3 seed, ~5–7h) — MINI de-risk 먼저"),
    ("이후", PRIMARY, "SV baselines LLM 이식 (GTG/FedSV/ComFedSV/Ripple)"),
    ("Phase 2", GRAY, "full baseline + Data Banzhaf + ShapleyFL + detection + 3B/7B"),
    ("Phase 3", GRAY, "실험 매트릭스 실행 + Ripple 이론 reduction → 논문"),
]
y = 1.95
for label, col, txt in steps:
    pill(s, ML, y + 0.02, label, col, w=1.06, h=0.30, size=10)
    tf = box(s, ML + 1.2, y, 4.85, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    add_text(tf, txt, size=11.5, color=INK, first=True, line=1.05)
    y += 0.62
# limitations (right)
rx = 7.05
tf = box(s, rx, 1.55, 5.6, 0.3)
add_text(tf, "한계 (정직하게 보고)", size=13, bold=True, color=PRIMARY, first=True)
lims = [
    ("Privacy", "서버가 개별 Δw_k 를 봐야 함 → secure aggregation과 비호환. client-level 가치평가의 본질적 한계 (고칠 결함 아님)."),
    ("noise vs OOD-good", "signed value 안에서 분리하는 FL 방법은 없음 (detector들 non-IID서 붕괴) → characterized limitation으로 보류."),
    ("non-IID bias", "FL-Shapley가 maverick/희귀도메인 client를 과소평가 → α-sweep 필수 측정 의무."),
]
y = 1.95
for t, b in lims:
    card(s, rx, y, 5.65, 1.32, fill=CARD, edge=EDGE)
    tf = box(s, rx + 0.18, y + 0.12, 5.3, 1.1)
    add_text(tf, t, size=11.5, bold=True, color=NAVY, first=True)
    add_text(tf, b, size=11, color=INK, sb=3, line=1.12)
    y += 1.46
footer(s, 12)

# ============================================================
OUTDIR = os.path.dirname(os.path.abspath(__file__))
os.makedirs(OUTDIR, exist_ok=True)
pptx_path = os.path.join(OUTDIR, "flirds-progress-2026-06-04.pptx")
prs.save(pptx_path)
print("SAVED", pptx_path, len(prs.slides.__iter__.__self__._sldIdLst))
