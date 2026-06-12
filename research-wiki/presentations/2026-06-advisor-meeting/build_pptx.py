#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML 데크 → 편집 가능한 PPTX 변환기 (상세판 + 요약판).

사용:
    python build_pptx.py        # 두 데크 모두 변환

산출 (이 디렉토리):
    flirds-advisor-2026-06.pptx        ← flirds-advisor-2026-06.html
    flirds-advisor-2026-06-brief.pptx  ← flirds-advisor-2026-06-brief.html

소스는 빌드된 HTML(단일 소스) — 내용 갱신 흐름:
    build.py / build_brief.py 수정 → python build.py(/_brief.py) → python build_pptx.py
텍스트·표·서식(볼드/작은 회색/위첨자)을 그대로 PPTX 네이티브 요소로 옮기므로
PowerPoint에서 자유롭게 편집 가능. 16:9 (13.33in × 7.5in = 1280×720px@96dpi).
"""

import math
import re
import sys
from html.parser import HTMLParser
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent

INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x17, 0x50, 0x8C)
PH = RGBColor(0x99, 0x99, 0x99)
LINE = RGBColor(0xC8, 0xC8, 0xC8)
FAINT = RGBColor(0x88, 0x88, 0x88)
HDRBG = RGBColor(0xF4, 0xF4, 0xF4)
BOXBG = RGBColor(0xFA, 0xFA, 0xFA)

FONT = "Malgun Gothic"
FONT_MATH = "Cambria"

CONFIGS = {
    "flirds-advisor-2026-06.html": dict(
        out="flirds-advisor-2026-06.pptx", margin=56, top=44,
        body=16.5, dense=14.5, lh=1.52, lh_dense=1.45,
        table=14.5, table_dense=13.5, t_lh=1.4,
        bt=15, kicker=14, h2=27, formula=19, fnote=13, box=15,
        li_gap=7, li_gap_tight=4, col_gap=36),
    "flirds-advisor-2026-06-brief.html": dict(
        out="flirds-advisor-2026-06-brief.pptx", margin=60, top=48,
        body=20, dense=17.5, lh=1.6, lh_dense=1.55,
        table=17, table_dense=15.5, t_lh=1.45,
        bt=16, kicker=15, h2=30, formula=23, fnote=14, box=18,
        li_gap=10, li_gap_tight=10, col_gap=40),
}

PAGE_W, PAGE_H = 1280, 720
px = lambda v: Emu(int(round(v * 9525)))          # 1px @96dpi = 9525 EMU
fpt = lambda v: Pt(round(v * 0.75, 1))             # px → pt


# ──────────────────────────────────────────────────────────────────
# 미니 HTML DOM
# ──────────────────────────────────────────────────────────────────
class Node:
    def __init__(self, tag, attrs=()):
        self.tag = tag
        self.attrs = dict(attrs)
        self.children = []          # Node 또는 str

    @property
    def cls(self):
        return (self.attrs.get("class") or "").split()

    def find_all(self, tag=None, cls=None, rec=True):
        out = []
        for c in self.children:
            if isinstance(c, Node):
                if (tag is None or c.tag == tag) and (cls is None or cls in c.cls):
                    out.append(c)
                if rec:
                    out += c.find_all(tag, cls, rec)
        return out

    def text(self):
        parts = []
        for c in self.children:
            parts.append(c if isinstance(c, str) else c.text())
        return "".join(parts)


VOID = {"br", "hr", "meta", "img", "link", "input"}


class MiniDOM(HTMLParser):
    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.root = Node("root")
        self.stack = [self.root]

    def handle_starttag(self, tag, attrs):
        n = Node(tag, attrs)
        self.stack[-1].children.append(n)
        if tag not in VOID:
            self.stack.append(n)

    def handle_endtag(self, tag):
        for i in range(len(self.stack) - 1, 0, -1):
            if self.stack[i].tag == tag:
                del self.stack[i:]
                break

    def handle_data(self, data):
        if data:
            self.stack[-1].children.append(data)


def parse_html(path: Path) -> Node:
    p = MiniDOM()
    p.feed(path.read_text(encoding="utf-8"))
    return p.root


# ──────────────────────────────────────────────────────────────────
# 인라인 → run 목록   run = (text, dict(bold, scale, color, base, italic, math))
# ──────────────────────────────────────────────────────────────────
BREAK = ("\x00BR\x00", None)


def inline_runs(node, st=None):
    st = dict(bold=False, scale=1.0, color=None, base=0, italic=False) if st is None else dict(st)
    out = []
    for c in node.children:
        if isinstance(c, str):
            t = re.sub(r"\s+", " ", c)
            if t:
                out.append((t, dict(st)))
            continue
        s = dict(st)
        if c.tag == "br":
            out.append(BREAK)
            continue
        if c.tag == "b":
            s["bold"] = True
        elif c.tag in ("i", "var"):
            s["italic"] = True
        elif c.tag == "sub":
            s["base"] = -25000
        elif c.tag == "sup":
            s["base"] = 30000
        elif c.tag == "span":
            cl = c.cls
            if "lbl" in cl:
                s["bold"] = True
            if "sm" in cl:
                s["scale"] = 0.82
                s["color"] = s["color"] or SUB
            if "muted" in cl:
                s["color"] = s["color"] or SUB
            if "ph" in cl:
                s["color"] = PH
                s["italic"] = True
            if "note" in cl:        # formula 안의 note는 별도 처리
                continue
        out.append(("", None))      # 자리표시 없음 — 아래에서 재귀 결과 붙임
        out.pop()
        out += inline_runs(c, s)
    return out


def runs_to_para(p, runs, size_px, default_color=INK, font=FONT):
    for t, s in runs:
        if (t, s) == BREAK:
            p._p.append(p._p.makeelement(qn("a:br"), {}))
            continue
        if not t:
            continue
        r = p.add_run()
        r.text = t
        f = r.font
        f.name = font
        f.size = fpt(size_px * (s["scale"] if s else 1.0))
        f.bold = bool(s and s["bold"])
        f.italic = bool(s and s["italic"])
        f.color.rgb = (s and s["color"]) or default_color
        if s and s["base"]:
            r._r.get_or_add_rPr().set("baseline", str(s["base"]))


# ──────────────────────────────────────────────────────────────────
# 높이 추정 (블록 스택 배치용)
# ──────────────────────────────────────────────────────────────────
def vis_w(ch):
    o = ord(ch)
    if o >= 0x2E80 or 0x1100 <= o <= 0x11FF or ch in "①②③④⑤⑥⑦ⓐⓑⓒ→↔≡≥≤·×—≈±":
        return 1.0
    if ch == " ":
        return 0.30
    return 0.52


def est_lines(runs, width_px, size_px):
    segs, cur = [], 0.0
    for t, s in runs:
        if (t, s) == BREAK:
            segs.append(cur)
            cur = 0.0
            continue
        sc = s["scale"] if s else 1.0
        cur += sum(vis_w(ch) for ch in t) * size_px * sc
    segs.append(cur)
    return sum(max(1, math.ceil(w * 1.10 / max(width_px, 1))) for w in segs)


# ──────────────────────────────────────────────────────────────────
# 도형 헬퍼
# ──────────────────────────────────────────────────────────────────
def add_tb(sl, x, y, w, h):
    tb = sl.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def add_rect(sl, x, y, w, h, fill=None, line_color=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = px(line_w)
    return sp


def set_cell_borders(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in reversed(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
        ln = tcPr.makeelement(qn(tag), {"w": "9525", "cap": "flat"})
        sf = ln.makeelement(qn("a:solidFill"), {})
        sf.append(sf.makeelement(qn("a:srgbClr"), {"val": "C8C8C8"}))
        ln.append(sf)
        tcPr.insert(0, ln)


def para(tf, first):
    return tf.paragraphs[0] if first else tf.add_paragraph()


# ──────────────────────────────────────────────────────────────────
# 블록 렌더러
# ──────────────────────────────────────────────────────────────────
def style_width(node):
    m = re.search(r"width:\s*([\d.]+)%", node.attrs.get("style", ""))
    return float(m.group(1)) / 100 if m else None


def render_table(sl, node, x, y, w, cfg):
    rows = node.find_all("tr", rec=False) or node.find_all("tr")
    ncol = max(len(r.find_all(rec=False)) for r in rows)
    size = cfg["table_dense"] if "dense" in node.cls else cfg["table"]
    pad_h, pad_v = 2 * 10, 2 * 6

    widths, fixed = [None] * ncol, 0.0
    for i, c in enumerate(rows[0].find_all(rec=False)):
        widths[i] = style_width(c)
        fixed += widths[i] or 0
    free = [i for i, v in enumerate(widths) if v is None]
    for i in free:
        widths[i] = (1.0 - fixed) / len(free)
    col_px = [max(40, w * v) for v in widths]

    row_h = []
    for r in rows:
        cells = r.find_all(rec=False)
        lines = 1
        for i, c in enumerate(cells):
            lines = max(lines, est_lines(inline_runs(c), col_px[i] - pad_h, size))
        row_h.append(lines * size * cfg["t_lh"] + pad_v)

    gfx = sl.shapes.add_table(len(rows), ncol, px(x), px(y), px(w), px(sum(row_h)))
    tbl = gfx.table
    tbl.first_row = tbl.horz_banding = False
    for i, v in enumerate(col_px):
        tbl.columns[i].width = px(v)
    for ri, r in enumerate(rows):
        tbl.rows[ri].height = px(row_h[ri])
        for ci, c in enumerate(r.find_all(rec=False)):
            cell = tbl.cell(ri, ci)
            set_cell_borders(cell)
            cell.fill.solid()
            cell.fill.fore_color.rgb = HDRBG if c.tag == "th" else RGBColor(0xFF, 0xFF, 0xFF)
            cell.margin_left = cell.margin_right = px(10)
            cell.margin_top = cell.margin_bottom = px(5)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            runs = inline_runs(c)
            if c.tag == "th":
                runs = [(t, dict(s, bold=True) if s else s) for t, s in runs]
            runs_to_para(p, runs, size)
    return sum(row_h)


def render_stack(sl, blocks, x, y, w, base, cfg, first_bt=True):
    """블록 목록을 위에서 아래로 배치. 새 y를 반환."""
    for b in blocks:
        if isinstance(b, str):
            continue
        if b.tag == "div" and "bt" in b.cls:
            if not first_bt:
                y += 14
            first_bt = False
            tf = add_tb(sl, x, y, w, cfg["bt"] * 1.5)
            runs = [(t, dict(s, bold=True)) for t, s in inline_runs(b) if (t, s) != BREAK]
            runs_to_para(tf.paragraphs[0], runs, cfg["bt"], default_color=ACCENT)
            y += cfg["bt"] * 1.45 + 6

        elif b.tag == "ul":
            items = [c for c in b.children if isinstance(c, Node) and c.tag == "li"]
            h_total = 0
            tf = add_tb(sl, x, y, w, 10)
            for i, li in enumerate(items):
                p = para(tf, i == 0)
                pPr = p._p.get_or_add_pPr()
                pPr.set("marL", str(int(14 * 9525)))
                pPr.set("indent", str(int(-14 * 9525)))
                gap = cfg["li_gap_tight"] if "tight" in li.cls else cfg["li_gap"]
                p.space_after = fpt(gap)
                runs = [("–  ", dict(bold=False, scale=1.0, color=SUB, base=0, italic=False))]
                runs += inline_runs(li)
                runs_to_para(p, runs, base)
                lh = cfg["lh_dense"] if base == cfg["dense"] else cfg["lh"]
                h_total += est_lines(runs, w - 14, base) * base * lh + gap
            y += h_total + 2

        elif b.tag == "table":
            y += 4
            y += render_table(sl, b, x, y, w, cfg) + 10

        elif b.tag == "p":
            runs = inline_runs(b)
            size = base
            if "sm" in b.cls:
                size = base * 0.82
            tf = add_tb(sl, x, y, w, 10)
            tf.paragraphs[0]
            runs_to_para(tf.paragraphs[0], runs, size,
                         default_color=SUB if ("muted" in b.cls or "sm" in b.cls) else INK)
            y += est_lines(runs, w, size) * size * cfg["lh"] + 10

        elif b.tag == "div" and "formula" in b.cls:
            note = next(iter(b.find_all("span", "note")), None)
            main_runs = inline_runs(b)
            n_lines = est_lines(main_runs, w - 36, cfg["formula"])
            note_runs = inline_runs(note) if note is not None else []
            note_lines = est_lines(note_runs, w - 36, cfg["fnote"]) if note_runs else 0
            h = n_lines * cfg["formula"] * 1.5 + note_lines * cfg["fnote"] * 1.5 + 26
            add_rect(sl, x, y, w, h, fill=BOXBG, line_color=LINE)
            tf = add_tb(sl, x + 18, y + 12, w - 36, h - 24)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            runs_to_para(p, main_runs, cfg["formula"], font=FONT_MATH)
            if note_runs:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = fpt(6)
                runs_to_para(p2, note_runs, cfg["fnote"], default_color=SUB)
            y += h + 12

        elif b.tag == "div" and "box" in b.cls:
            y += 12
            runs = inline_runs(b)
            h = est_lines(runs, w - 28, cfg["box"]) * cfg["box"] * 1.5 + 22
            add_rect(sl, x, y, w, h, fill=None, line_color=INK)
            tf = add_tb(sl, x + 14, y + 10, w - 28, h - 20)
            runs_to_para(tf.paragraphs[0], runs, cfg["box"])
            y += h

        elif b.tag == "div" and "cols" in b.cls:
            cols = [c for c in b.children if isinstance(c, Node) and c.tag == "div"]
            cw = (w - cfg["col_gap"] * (len(cols) - 1)) / len(cols)
            y_end = y
            for i, col in enumerate(cols):
                ye = render_stack(sl, col.children, x + i * (cw + cfg["col_gap"]),
                                  y, cw, base, cfg, first_bt=True)
                y_end = max(y_end, ye)
            y = y_end

        elif b.tag == "div":      # 알 수 없는 div — 내용만 재귀
            y = render_stack(sl, b.children, x, y, w, base, cfg, first_bt)
    return y


# ──────────────────────────────────────────────────────────────────
# 슬라이드
# ──────────────────────────────────────────────────────────────────
def add_foot(sl, sec, cfg, label, n, total):
    m = cfg["margin"]
    tf = add_tb(sl, m, PAGE_H - 30, 700, 18)
    runs_to_para(tf.paragraphs[0], [(label, None)], 11.5, default_color=FAINT)
    tf = add_tb(sl, PAGE_W - m - 120, PAGE_H - 30, 120, 18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    runs_to_para(p, [(f"{n} / {total}", None)], 11.5, default_color=FAINT)


def render_cover(sl, sec, cfg):
    m = cfg["margin"]
    h1 = next(iter(sec.find_all("h1")), None)
    tf = add_tb(sl, m, 215, PAGE_W - 2 * m, 70)
    runs_to_para(tf.paragraphs[0], [(h1.text().strip(), dict(bold=True, scale=1.0, color=None, base=0, italic=False))], 44)
    sub = next(iter(sec.find_all("div", "subtitle")), None)
    tf = add_tb(sl, m, 290, PAGE_W - 2 * m, 34)
    runs_to_para(tf.paragraphs[0], [(sub.text().strip(), None)], 21, default_color=SUB)
    meta = next(iter(sec.find_all("div", "meta")), None)
    tf = add_tb(sl, m, 326, PAGE_W - 2 * m, 26)
    runs_to_para(tf.paragraphs[0], [(meta.text().strip(), None)], 16, default_color=SUB)

    conv = next(iter(sec.find_all("div", "conv")), None)
    if conv is not None:
        t = next(iter(conv.find_all("div", "t")), None)
        # div.t를 제외한 나머지를 래퍼에 담아 inline_runs로 — br(BREAK)·span 스타일 보존
        wrapper = Node("div")
        wrapper.children = [c for c in conv.children
                            if not (isinstance(c, Node) and c.tag == "div" and "t" in c.cls)]
        body_runs = [r for r in inline_runs(wrapper)]
        while body_runs and body_runs[0] != BREAK and not body_runs[0][0].strip():
            body_runs.pop(0)
        w = 880
        lines = est_lines(body_runs, w - 36, 14.5)
        h = 24 + 22 + lines * 14.5 * 1.6 + 10
        add_rect(sl, m, 396, w, h, fill=None, line_color=LINE)
        tf = add_tb(sl, m + 18, 396 + 12, w - 36, h - 24)
        p = tf.paragraphs[0]
        runs_to_para(p, [(t.text().strip(), dict(bold=True, scale=1.0, color=None, base=0, italic=False))], 15)
        p2 = tf.add_paragraph()
        p2.space_before = fpt(4)
        runs_to_para(p2, body_runs, 14.5)


def render_slide(prs, sec, cfg, n, total, foot_label):
    sl = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    if "cover" in sec.cls:
        render_cover(sl, sec, cfg)
        add_foot(sl, sec, cfg, foot_label, n, total)
        return
    m, y = cfg["margin"], cfg["top"]
    kicker = next(iter(sec.find_all("div", "kicker")), None)
    h2 = next(iter(sec.find_all("h2")), None)
    body = next(iter(sec.find_all("div", "body")), None)
    w = PAGE_W - 2 * m

    tf = add_tb(sl, m, y, w, cfg["kicker"] * 1.4)
    runs_to_para(tf.paragraphs[0], [(kicker.text().strip(), None)], cfg["kicker"], default_color=ACCENT)
    y += cfg["kicker"] * 1.35 + 4
    tf = add_tb(sl, m, y, w, cfg["h2"] * 1.4)
    runs_to_para(tf.paragraphs[0],
                 [(h2.text().strip(), dict(bold=True, scale=1.0, color=None, base=0, italic=False))], cfg["h2"])
    y += cfg["h2"] * 1.35 + 10
    add_rect(sl, m, y, w, 1.2, fill=LINE)
    y += 1 + (20 if cfg is CONFIGS["flirds-advisor-2026-06-brief.html"] else 16)

    base = cfg["dense"] if "dense" in body.cls else cfg["body"]
    render_stack(sl, body.children, m, y, w, base, cfg)
    add_foot(sl, sec, cfg, foot_label, n, total)


def build(html_name: str):
    cfg = CONFIGS[html_name]
    src = HERE / html_name
    root = parse_html(src)
    title = next(iter(root.find_all("title")), None)
    foot_label = title.text().split("·")[0].strip() + " · " + title.text().split("·")[-1].strip() \
        if title is not None else ""
    slides = root.find_all("section", "slide")
    prs = Presentation()
    prs.slide_width = px(PAGE_W)
    prs.slide_height = px(PAGE_H)
    for i, sec in enumerate(slides):
        render_slide(prs, sec, cfg, i + 1, len(slides), foot_label)
    out = HERE / cfg["out"]
    prs.save(out)
    print(f"[ok] PPTX → {out.name} ({len(slides)}장)")


if __name__ == "__main__":
    targets = sys.argv[1:] or list(CONFIGS)
    for name in targets:
        build(name)
